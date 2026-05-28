"""APOLLO PPTX生成スクリプト — 炭酸カルシウム技術動向分析 2026

session_dir で実行: python3 reports/generate_pptx.py
出力: reports/apollo_report_20260527.pptx
"""
import os, re, sys
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_AUTO_SIZE
from pptx.enum.shapes import MSO_SHAPE
from PIL import Image
from lxml import etree

# ============ デザインシステム ============
JA_FONT = "Meiryo UI"
LATIN_FONT = "Meiryo UI"
A_NS = 'http://schemas.openxmlformats.org/drawingml/2006/main'

NAVY = RGBColor(0x1B, 0x2A, 0x4A)
BLUE = RGBColor(0x2E, 0x50, 0x90)
ACCENT = RGBColor(0x3B, 0x7D, 0xD8)
DARK_GRAY = RGBColor(0x33, 0x33, 0x33)
MEDIUM_GRAY = RGBColor(0x66, 0x66, 0x66)
LIGHT_GRAY = RGBColor(0xF2, 0xF2, 0xF2)
BORDER_GRAY = RGBColor(0xCC, 0xCC, 0xCC)
KEY_MSG_BG = RGBColor(0xE8, 0xF0, 0xFE)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
RED_ACCENT = RGBColor(0xD6, 0x45, 0x45)
GREEN_ACCENT = RGBColor(0x2E, 0x8B, 0x57)
AMBER = RGBColor(0xD4, 0xA0, 0x17)
GHOST_NAVY = RGBColor(0x2A, 0x3A, 0x5A)
BOTTOM_BAR_Y = 6.92

SESSION = os.path.dirname(os.path.abspath(__file__)).replace("/reports", "")
TEMPLATE = os.path.join(SESSION, "capcom_schema/templates/apollo_template.pptx")
SNAP = os.path.join(SESSION, "snapshots")

# ============ コアユーティリティ ============
def _apply_font(run):
    run.font.name = LATIN_FONT
    rPr = run._r.get_or_add_rPr()
    rPr.set('lang', 'ja-JP')
    rPr.set('altLang', 'en-US')
    ea = rPr.find(f'{{{A_NS}}}ea')
    if ea is None:
        ea = etree.SubElement(rPr, f'{{{A_NS}}}ea')
    ea.set('typeface', JA_FONT)

def _apply_kinsoku(p):
    pPr = p._p.get_or_add_pPr()
    pPr.set('eaLnBrk', '1')
    pPr.set('hangingPunct', '1')

def set_text(p, text, size, color, bold=False, line_spacing=None):
    p.text = ""
    run = p.add_run()
    run.text = text
    run.font.size = size
    run.font.color.rgb = color
    run.font.bold = bold
    _apply_font(run)
    _apply_kinsoku(p)
    if line_spacing:
        p.line_spacing = line_spacing

def add_rich_runs(p, text, base_size=Pt(14), base_color=DARK_GRAY,
                  bold_color=None, force_bold=False, line_spacing=1.4):
    p.clear()
    bold_color = bold_color or base_color
    _apply_kinsoku(p)
    p.line_spacing = line_spacing
    parts = re.split(r'(\*\*.*?\*\*)', text)
    for part in parts:
        if not part: continue
        run = p.add_run()
        if part.startswith('**') and part.endswith('**'):
            run.text = part[2:-2]
            run.font.bold = True
            run.font.color.rgb = bold_color
        else:
            run.text = part
            run.font.bold = force_bold
            run.font.color.rgb = bold_color if force_bold else base_color
        run.font.size = base_size
        _apply_font(run)

def add_title_shape(slide, text, x=0.5, y=0.15, w=12.3):
    n = len(text)
    if n <= 30: font_size, box_h = Pt(24), 0.65
    elif n <= 50: font_size, box_h = Pt(22), 0.75
    else: font_size, box_h = Pt(20), 0.90
    txBox = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(box_h))
    tf = txBox.text_frame
    tf.word_wrap = True
    tf.auto_size = MSO_AUTO_SIZE.NONE
    add_rich_runs(tf.paragraphs[0], text, base_size=font_size,
                  base_color=NAVY, bold_color=NAVY, force_bold=True, line_spacing=1.3)
    line_y = y + box_h + 0.05
    line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(x), Inches(line_y), Inches(w), Emu(12700))
    line.fill.solid(); line.fill.fore_color.rgb = ACCENT; line.line.fill.background()
    return line_y + 0.08

def add_sub_message(slide, message, x=0.5, y=None, w=12.3):
    if y is None: y = 1.00
    chars_per_line = 40
    num_lines = max(1, -(-len(message) // chars_per_line))
    box_h = 0.20 + num_lines * 0.35
    box = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(box_h))
    box.fill.solid(); box.fill.fore_color.rgb = KEY_MSG_BG; box.line.fill.background()
    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(x), Inches(y), Emu(36576), Inches(box_h))
    bar.fill.solid(); bar.fill.fore_color.rgb = ACCENT; bar.line.fill.background()
    txBox = slide.shapes.add_textbox(Inches(x + 0.15), Inches(y + 0.08),
                                      Inches(w - 0.30), Inches(box_h - 0.16))
    tf = txBox.text_frame; tf.word_wrap = True; tf.auto_size = MSO_AUTO_SIZE.NONE
    p = tf.paragraphs[0]
    marker = p.add_run(); marker.text = "■ "; marker.font.size = Pt(16)
    marker.font.color.rgb = NAVY; marker.font.bold = True; _apply_font(marker)
    parts = re.split(r'(\*\*.*?\*\*)', message)
    for part in parts:
        if not part: continue
        run = p.add_run()
        if part.startswith('**') and part.endswith('**'):
            run.text = part[2:-2]; run.font.bold = True; run.font.color.rgb = NAVY
        else:
            run.text = part; run.font.color.rgb = DARK_GRAY
        run.font.size = Pt(16); _apply_font(run)
    _apply_kinsoku(p); p.line_spacing = 1.5
    return y + box_h + 0.10

def add_bottom_bar_and_footer(slide, page_num=None):
    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(BOTTOM_BAR_Y),
                                  Inches(13.33), Emu(50800))
    bar.fill.solid(); bar.fill.fore_color.rgb = NAVY; bar.line.fill.background()
    line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.5), Inches(6.95),
                                   Inches(12.3), Emu(9525))
    line.fill.solid(); line.fill.fore_color.rgb = BORDER_GRAY; line.line.fill.background()
    txBox = slide.shapes.add_textbox(Inches(0.5), Inches(7.05), Inches(6), Inches(0.25))
    set_text(txBox.text_frame.paragraphs[0], "APOLLO", Pt(8), MEDIUM_GRAY)
    if page_num is not None:
        txBox2 = slide.shapes.add_textbox(Inches(10.5), Inches(7.05), Inches(2.3), Inches(0.25))
        p2 = txBox2.text_frame.paragraphs[0]
        set_text(p2, f"| {page_num}", Pt(8), MEDIUM_GRAY); p2.alignment = PP_ALIGN.RIGHT

def fit_image(slide, image_path, max_x, max_y, max_w, max_h):
    if not os.path.exists(image_path): return None
    img = Image.open(image_path)
    img_w, img_h = img.size
    ratio = img_h / img_w
    if max_w * ratio <= max_h:
        use_w, use_h = max_w, max_w * ratio
    else:
        use_h, use_w = max_h, max_h / ratio
    left = max_x + (max_w - use_w) / 2
    top = max_y + (max_h - use_h) / 2
    pic = slide.shapes.add_picture(image_path, Inches(left), Inches(top),
                                    width=Inches(use_w), height=Inches(use_h))
    img.close()
    return pic

def add_source_label(slide, text, y=6.55):
    txBox = slide.shapes.add_textbox(Inches(0.5), Inches(y), Inches(12.3), Inches(0.35))
    tf = txBox.text_frame; tf.word_wrap = True; tf.auto_size = MSO_AUTO_SIZE.NONE
    set_text(tf.paragraphs[0], f"（出所）{text}", Pt(9), MEDIUM_GRAY)

def add_annotation_block(slide, bullets, x, y, w, h, font_size=14, bg_color=None, has_border=False):
    if bg_color or has_border:
        box = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(h))
        if bg_color:
            box.fill.solid(); box.fill.fore_color.rgb = bg_color
        else: box.fill.background()
        if has_border:
            box.line.color.rgb = BORDER_GRAY; box.line.width = Emu(9525)
        else: box.line.fill.background()
    txBox = slide.shapes.add_textbox(Inches(x + 0.12), Inches(y + 0.08),
                                      Inches(w - 0.24), Inches(h - 0.16))
    tf = txBox.text_frame; tf.word_wrap = True; tf.auto_size = MSO_AUTO_SIZE.NONE
    for i, item in enumerate(bullets):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.space_after = Pt(6)
        marker = p.add_run(); marker.text = "■ "; marker.font.size = Pt(font_size)
        marker.font.color.rgb = NAVY; _apply_font(marker)
        parts = re.split(r'(\*\*.*?\*\*)', item)
        for part in parts:
            if not part: continue
            run = p.add_run()
            if part.startswith('**') and part.endswith('**'):
                run.text = part[2:-2]; run.font.bold = True; run.font.color.rgb = NAVY
            else:
                run.text = part; run.font.color.rgb = DARK_GRAY
            run.font.size = Pt(font_size); _apply_font(run)
        _apply_kinsoku(p); p.line_spacing = 1.5

def add_chart_label(slide, text, x, y, w=3.0, size=14, color=NAVY):
    txBox = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(0.35))
    set_text(txBox.text_frame.paragraphs[0], text, Pt(size), color, bold=True)

# ============ スライドタイプ ============
def add_title_slide(prs, title, subtitle, date, blank):
    slide = prs.slides.add_slide(blank)
    bg = slide.background.fill; bg.solid(); bg.fore_color.rgb = NAVY
    line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(1.2), Inches(1.8),
                                   Inches(2.0), Emu(27432))
    line.fill.solid(); line.fill.fore_color.rgb = ACCENT; line.line.fill.background()
    logo = slide.shapes.add_textbox(Inches(1.2), Inches(1.2), Inches(3), Inches(0.5))
    set_text(logo.text_frame.paragraphs[0], "APOLLO", Pt(14), ACCENT, bold=True)
    txBox = slide.shapes.add_textbox(Inches(1.2), Inches(2.1), Inches(11), Inches(2))
    tf = txBox.text_frame; tf.word_wrap = True
    set_text(tf.paragraphs[0], title, Pt(36), WHITE, bold=True, line_spacing=1.2)
    txBox2 = slide.shapes.add_textbox(Inches(1.2), Inches(4.2), Inches(11), Inches(1))
    set_text(txBox2.text_frame.paragraphs[0], subtitle, Pt(18), RGBColor(0xAA, 0xBB, 0xDD))
    txBox3 = slide.shapes.add_textbox(Inches(1.2), Inches(5.5), Inches(11), Inches(0.5))
    set_text(txBox3.text_frame.paragraphs[0], date, Pt(13), RGBColor(0x88, 0x99, 0xBB))
    bot = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(7.1),
                                  Inches(13.33), Emu(27432))
    bot.fill.solid(); bot.fill.fore_color.rgb = ACCENT; bot.line.fill.background()
    return slide

def add_section_slide(prs, section_num, title, blank, subtitle=None):
    slide = prs.slides.add_slide(blank)
    bg = slide.background.fill; bg.solid(); bg.fore_color.rgb = NAVY
    ghost = slide.shapes.add_textbox(Inches(0.5), Inches(1.0), Inches(5), Inches(3.5))
    p_g = ghost.text_frame.paragraphs[0]; run_g = p_g.add_run()
    run_g.text = f"{section_num:02d}"; run_g.font.size = Pt(180)
    run_g.font.color.rgb = GHOST_NAVY; run_g.font.bold = True; _apply_font(run_g)
    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(1.0), Inches(3.0),
                                  Emu(36576), Inches(2.0))
    bar.fill.solid(); bar.fill.fore_color.rgb = ACCENT; bar.line.fill.background()
    txBox = slide.shapes.add_textbox(Inches(1.3), Inches(3.2), Inches(11), Inches(1.5))
    tf = txBox.text_frame; tf.word_wrap = True
    set_text(tf.paragraphs[0], title, Pt(32), WHITE, bold=True, line_spacing=1.2)
    if subtitle:
        txBox2 = slide.shapes.add_textbox(Inches(1.3), Inches(4.8), Inches(11), Inches(0.8))
        set_text(txBox2.text_frame.paragraphs[0], subtitle, Pt(16), RGBColor(0xCC, 0xDD, 0xEE))
    bot = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(7.1),
                                  Inches(13.33), Emu(27432))
    bot.fill.solid(); bot.fill.fore_color.rgb = ACCENT; bot.line.fill.background()
    return slide

def add_chart_text_slide(prs, title, sub_message, image_path, annotations, blank,
                          caption=None, chart_label=None, source=None, page_num=None,
                          chart_ratio=0.60):
    slide = prs.slides.add_slide(blank)
    sub_y = add_title_shape(slide, title)
    content_y = add_sub_message(slide, sub_message, y=sub_y)
    content_w = 12.3; gap = 0.3
    chart_w = content_w * chart_ratio - gap / 2
    text_w = content_w * (1 - chart_ratio) - gap / 2
    remaining_h = 6.5 - content_y
    chart_x = 0.5; text_x = 0.5 + chart_w + gap
    if chart_label:
        add_chart_label(slide, chart_label, chart_x, content_y, chart_w)
        img_y = content_y + 0.35; img_h = remaining_h - 0.65
    else:
        img_y = content_y; img_h = remaining_h - 0.3
    full_path = image_path if os.path.isabs(image_path) else os.path.join(SNAP, image_path)
    fit_image(slide, full_path, chart_x, img_y, chart_w, img_h)
    if caption:
        txBox = slide.shapes.add_textbox(Inches(chart_x), Inches(content_y + remaining_h - 0.25),
                                          Inches(chart_w), Inches(0.25))
        p = txBox.text_frame.paragraphs[0]
        set_text(p, caption, Pt(10), MEDIUM_GRAY); p.alignment = PP_ALIGN.CENTER
    add_annotation_block(slide, annotations[:5], text_x, content_y, text_w, remaining_h - 0.2)
    if source: add_source_label(slide, source)
    add_bottom_bar_and_footer(slide, page_num)
    return slide

def add_kpi_slide(prs, title, sub_message, kpis, blank, source=None, page_num=None):
    slide = prs.slides.add_slide(blank)
    sub_y = add_title_shape(slide, title)
    content_y = add_sub_message(slide, sub_message, y=sub_y)
    n = len(kpis); available_w = 11.5; start_x = 0.9; gap = 0.2
    if n <= 4: n_cols, n_rows = n, 1
    else: n_cols, n_rows = min(4, (n + 1) // 2), 2
    card_w = (available_w - gap * (n_cols - 1)) / n_cols
    available_h = 6.5 - content_y; row_gap = 0.2
    card_h = min((available_h - row_gap * (n_rows - 1)) / n_rows, 2.8)
    for idx, kpi in enumerate(kpis):
        row = idx // n_cols; col = idx % n_cols
        x = start_x + col * (card_w + gap); y = content_y + row * (card_h + row_gap)
        card = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(x), Inches(y),
                                       Inches(card_w), Inches(card_h))
        card.fill.solid(); card.fill.fore_color.rgb = LIGHT_GRAY
        card.line.color.rgb = BORDER_GRAY; card.line.width = Emu(9525)
        bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(x), Inches(y),
                                      Emu(36576), Inches(card_h))
        bar.fill.solid(); bar.fill.fore_color.rgb = ACCENT; bar.line.fill.background()
        txL = slide.shapes.add_textbox(Inches(x + 0.15), Inches(y + 0.12),
                                        Inches(card_w - 0.3), Inches(0.25))
        set_text(txL.text_frame.paragraphs[0], kpi["label"], Pt(10), MEDIUM_GRAY, bold=True)
        txV = slide.shapes.add_textbox(Inches(x + 0.15), Inches(y + 0.4),
                                        Inches(card_w - 0.3), Inches(0.9))
        p = txV.text_frame.paragraphs[0]
        run = p.add_run(); run.text = kpi["value"]
        # 値の長さで自動縮小
        v_len = len(kpi["value"])
        if v_len <= 5: vs = Pt(32)
        elif v_len <= 8: vs = Pt(22)
        elif v_len <= 12: vs = Pt(16)
        else: vs = Pt(13)
        run.font.size = vs; run.font.bold = True; run.font.color.rgb = NAVY; _apply_font(run)
        txU = slide.shapes.add_textbox(Inches(x + 0.15), Inches(y + card_h - 0.4),
                                        Inches(card_w - 0.3), Inches(0.25))
        set_text(txU.text_frame.paragraphs[0], kpi.get("unit", ""), Pt(9), MEDIUM_GRAY)
    if source: add_source_label(slide, source)
    add_bottom_bar_and_footer(slide, page_num)
    return slide

def add_cards_slide(prs, title, sub_message, cards, blank, source=None, page_num=None):
    slide = prs.slides.add_slide(blank)
    sub_y = add_title_shape(slide, title)
    content_y = add_sub_message(slide, sub_message, y=sub_y)
    n = len(cards); gap = 0.25; total_w = 12.3
    card_w = (total_w - gap * (n - 1)) / n
    card_h = 6.5 - content_y; header_h = 0.45
    colors = [NAVY, BLUE, ACCENT, GREEN_ACCENT, AMBER, RED_ACCENT]
    for i, card in enumerate(cards):
        x = 0.5 + i * (card_w + gap)
        color = card.get("color", colors[i % len(colors)])
        hdr = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(content_y),
                                      Inches(card_w), Inches(header_h))
        hdr.fill.solid(); hdr.fill.fore_color.rgb = color; hdr.line.fill.background()
        txH = slide.shapes.add_textbox(Inches(x + 0.1), Inches(content_y + 0.05),
                                        Inches(card_w - 0.2), Inches(header_h - 0.1))
        p_h = txH.text_frame.paragraphs[0]
        set_text(p_h, card["header"], Pt(13), WHITE, bold=True); p_h.alignment = PP_ALIGN.CENTER
        body_y = content_y + header_h; body_h = card_h - header_h
        bdy = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(x), Inches(body_y),
                                      Inches(card_w), Inches(body_h))
        bdy.fill.solid(); bdy.fill.fore_color.rgb = LIGHT_GRAY
        bdy.line.color.rgb = BORDER_GRAY; bdy.line.width = Emu(9525)
        txB = slide.shapes.add_textbox(Inches(x + 0.12), Inches(body_y + 0.1),
                                        Inches(card_w - 0.24), Inches(body_h - 0.2))
        tf = txB.text_frame; tf.word_wrap = True; tf.auto_size = MSO_AUTO_SIZE.NONE
        body_text = card.get("body", "")
        if isinstance(body_text, list):
            for j, item in enumerate(body_text):
                p = tf.paragraphs[0] if j == 0 else tf.add_paragraph()
                add_rich_runs(p, f"・{item}", base_size=Pt(11), base_color=DARK_GRAY,
                              bold_color=NAVY, line_spacing=1.4)
        else:
            add_rich_runs(tf.paragraphs[0], body_text, base_size=Pt(12),
                          base_color=DARK_GRAY, bold_color=NAVY, line_spacing=1.4)
    if source: add_source_label(slide, source)
    add_bottom_bar_and_footer(slide, page_num)
    return slide

def add_table_slide(prs, title, sub_message, headers, rows, blank,
                    col_widths=None, highlight_rows=None, annotations=None,
                    source=None, page_num=None):
    slide = prs.slides.add_slide(blank)
    sub_y = add_title_shape(slide, title)
    content_y = add_sub_message(slide, sub_message, y=sub_y)
    n_cols = len(headers); n_rows = len(rows) + 1
    if annotations:
        table_w = 7.5; text_x = 8.3; text_w = 4.5
    else:
        table_w = 12.3; text_x = None; text_w = 0
    available_h = 6.4 - content_y
    row_h = min(0.55, max(0.35, available_h / n_rows))
    table_h = row_h * n_rows
    table_shape = slide.shapes.add_table(n_rows, n_cols, Inches(0.5), Inches(content_y),
                                          Inches(table_w), Inches(table_h))
    table = table_shape.table
    if col_widths:
        for i, w in enumerate(col_widths): table.columns[i].width = Inches(w)
    for j, header in enumerate(headers):
        cell = table.cell(0, j); cell.text = ""
        set_text(cell.text_frame.paragraphs[0], header, Pt(13), WHITE, bold=True)
        cell.fill.solid(); cell.fill.fore_color.rgb = NAVY
    highlight_rows = highlight_rows or []
    for i, row in enumerate(rows):
        for j, value in enumerate(row):
            cell = table.cell(i + 1, j); cell.text = ""
            set_text(cell.text_frame.paragraphs[0], str(value), Pt(12), DARK_GRAY)
            if i in highlight_rows:
                cell.fill.solid(); cell.fill.fore_color.rgb = KEY_MSG_BG
            elif i % 2 == 1:
                cell.fill.solid(); cell.fill.fore_color.rgb = LIGHT_GRAY
    if annotations and text_x:
        add_annotation_block(slide, annotations, text_x, content_y, text_w,
                             6.4 - content_y, font_size=13, has_border=True, bg_color=KEY_MSG_BG)
    if source: add_source_label(slide, source)
    add_bottom_bar_and_footer(slide, page_num)
    return slide

def add_progress_bar_slide(prs, title, sub_message, items, blank, source=None, page_num=None):
    slide = prs.slides.add_slide(blank)
    sub_y = add_title_shape(slide, title)
    content_y = add_sub_message(slide, sub_message, y=sub_y)
    n = len(items); available_h = 6.5 - content_y
    bar_gap = 0.1; bar_group_h = (available_h - 0.2) / n
    bar_h = min(0.5, bar_group_h * 0.5); label_h = bar_group_h - bar_h - bar_gap
    bar_max_w = 9.0
    colors = [ACCENT, BLUE, NAVY, GREEN_ACCENT, AMBER, RED_ACCENT]
    for i, item in enumerate(items):
        gy = content_y + i * bar_group_h
        color = item.get("color", colors[i % len(colors)])
        pct = item["value"] / max(item.get("max_value", 100), 1)
        bar_w = bar_max_w * pct
        txL = slide.shapes.add_textbox(Inches(0.5), Inches(gy), Inches(3.0), Inches(label_h))
        set_text(txL.text_frame.paragraphs[0], item["label"], Pt(14), DARK_GRAY, bold=True)
        bg_bar = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(3.5),
                                         Inches(gy + label_h), Inches(bar_max_w), Inches(bar_h))
        bg_bar.fill.solid(); bg_bar.fill.fore_color.rgb = LIGHT_GRAY; bg_bar.line.fill.background()
        if bar_w > 0.1:
            val_bar = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(3.5),
                                              Inches(gy + label_h), Inches(bar_w), Inches(bar_h))
            val_bar.fill.solid(); val_bar.fill.fore_color.rgb = color; val_bar.line.fill.background()
        txP = slide.shapes.add_textbox(Inches(3.5 + bar_w + 0.1), Inches(gy + label_h),
                                        Inches(2.0), Inches(bar_h))
        set_text(txP.text_frame.paragraphs[0],
                 f"{item['value']}{item.get('unit', '%')}", Pt(14), color, bold=True)
    if source: add_source_label(slide, source)
    add_bottom_bar_and_footer(slide, page_num)
    return slide

def add_matrix_slide(prs, title, sub_message, quadrants, blank,
                     x_label="→ 累積件数（規模）", y_label="↑ 直近成長率",
                     source=None, page_num=None):
    slide = prs.slides.add_slide(blank)
    sub_y = add_title_shape(slide, title)
    content_y = add_sub_message(slide, sub_message, y=sub_y)
    mx = 1.0; my = content_y + 0.2; mw = 6.5; mh = 5.8 - content_y + 0.4
    half_w = mw / 2; half_h = mh / 2
    quad_colors = {"TL": ACCENT, "TR": GREEN_ACCENT, "BL": MEDIUM_GRAY, "BR": AMBER}
    positions = {"TL": (mx, my), "TR": (mx + half_w, my),
                 "BL": (mx, my + half_h), "BR": (mx + half_w, my + half_h)}
    for key, pos in positions.items():
        q = quadrants.get(key, {}); color = quad_colors[key]
        qx, qy = pos
        box = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(qx), Inches(qy),
                                      Inches(half_w - 0.05), Inches(half_h - 0.05))
        box.fill.solid(); box.fill.fore_color.rgb = LIGHT_GRAY
        box.line.color.rgb = BORDER_GRAY; box.line.width = Emu(9525)
        txBox = slide.shapes.add_textbox(Inches(qx + 0.1), Inches(qy + 0.1),
                                          Inches(half_w - 0.3), Inches(0.35))
        set_text(txBox.text_frame.paragraphs[0], q.get("title", ""), Pt(13), color, bold=True)
        items = q.get("items", [])
        if items:
            txBox2 = slide.shapes.add_textbox(Inches(qx + 0.15), Inches(qy + 0.5),
                                               Inches(half_w - 0.4), Inches(half_h - 0.7))
            tf = txBox2.text_frame; tf.word_wrap = True
            for j, item in enumerate(items[:6]):
                p = tf.paragraphs[0] if j == 0 else tf.add_paragraph()
                set_text(p, f"・{item}", Pt(10), DARK_GRAY)
    txX = slide.shapes.add_textbox(Inches(mx + mw/2 - 1.0), Inches(my + mh + 0.05),
                                    Inches(2.5), Inches(0.25))
    set_text(txX.text_frame.paragraphs[0], x_label, Pt(10), MEDIUM_GRAY)
    txY = slide.shapes.add_textbox(Inches(mx - 0.8), Inches(my + mh/2 - 0.15),
                                    Inches(0.8), Inches(0.3))
    set_text(txY.text_frame.paragraphs[0], y_label, Pt(10), MEDIUM_GRAY)
    # 右側に注釈
    ann_x = mx + mw + 0.5; ann_w = 13.33 - ann_x - 0.5
    ann_bullets = [
        "**成長リーダー（右上）**: 大規模×高成長。CO₂ 固定化クラスタが筆頭",
        "**新興（左上）**: 小規模×高成長。多面体炭酸・脱硫・石灰泥",
        "**成熟（右下）**: 大規模×低成長。軽質炭酸・ゴム・水性製造",
        "**衰退（左下）**: 縮小。製紙用・高級炭酸・パルプ廃液",
        "**全 35 クラスタ**を 4 象限分類、注目領域は明確",
    ]
    add_annotation_block(slide, ann_bullets, ann_x, my, ann_w, mh,
                          font_size=12, has_border=True, bg_color=KEY_MSG_BG)
    if source: add_source_label(slide, source)
    add_bottom_bar_and_footer(slide, page_num)
    return slide

def add_hypothesis_slide(prs, title, sub_message, hypotheses, blank, source=None, page_num=None):
    slide = prs.slides.add_slide(blank)
    sub_y = add_title_shape(slide, title)
    content_y = add_sub_message(slide, sub_message, y=sub_y)
    headers = ["ID", "仮説", "判定", "根拠"]
    n_rows = len(hypotheses) + 1
    available_h = 6.4 - content_y
    row_h = min(1.0, max(0.50, available_h / n_rows))
    table_h = row_h * n_rows
    VERDICT_MAP = {"confirmed": ("✅ 支持", GREEN_ACCENT),
                   "rejected": ("❌ 棄却", RED_ACCENT),
                   "partially": ("⚠ 部分", AMBER)}
    table_shape = slide.shapes.add_table(n_rows, 4, Inches(0.5), Inches(content_y),
                                          Inches(12.3), Inches(table_h))
    table = table_shape.table
    table.columns[0].width = Inches(0.7)
    table.columns[1].width = Inches(4.0)
    table.columns[2].width = Inches(1.3)
    table.columns[3].width = Inches(6.3)
    for j, header in enumerate(headers):
        cell = table.cell(0, j); cell.text = ""
        set_text(cell.text_frame.paragraphs[0], header, Pt(13), WHITE, bold=True)
        cell.fill.solid(); cell.fill.fore_color.rgb = NAVY
    for i, hyp in enumerate(hypotheses):
        v_label, v_color = VERDICT_MAP.get(hyp.get("verdict", "partially"), ("⚠ 部分", AMBER))
        row_data = [hyp.get("id", ""), hyp.get("hypothesis", ""), v_label, hyp.get("evidence", "")]
        for j, val in enumerate(row_data):
            cell = table.cell(i + 1, j); cell.text = ""
            if j == 2:
                p = cell.text_frame.paragraphs[0]
                set_text(p, val, Pt(12), WHITE, bold=True); p.alignment = PP_ALIGN.CENTER
                cell.fill.solid(); cell.fill.fore_color.rgb = v_color
            else:
                set_text(cell.text_frame.paragraphs[0], str(val), Pt(11), DARK_GRAY)
                if i % 2 == 1:
                    cell.fill.solid(); cell.fill.fore_color.rgb = LIGHT_GRAY
    if source: add_source_label(slide, source)
    add_bottom_bar_and_footer(slide, page_num)
    return slide

def add_stepup_slide(prs, title, sub_message, phases, blank, source=None, page_num=None):
    slide = prs.slides.add_slide(blank)
    sub_y = add_title_shape(slide, title)
    content_y = add_sub_message(slide, sub_message, y=sub_y)
    n = len(phases); gap = 0.2; total_w = 12.3
    bar_w = (total_w - gap * (n - 1)) / n
    base_y = 6.5; max_h = base_y - content_y - 0.2
    colors = [ACCENT, BLUE, NAVY, GREEN_ACCENT]
    for i, phase in enumerate(phases):
        x = 0.5 + i * (bar_w + gap)
        ratio = 0.5 + 0.5 * (i / max(n - 1, 1))
        bar_h = max_h * ratio; y = base_y - bar_h
        color = phase.get("color", colors[i % len(colors)])
        header_h = min(0.5, bar_h * 0.25)
        hdr = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(x), Inches(y),
                                      Inches(bar_w), Inches(header_h))
        hdr.fill.solid(); hdr.fill.fore_color.rgb = color; hdr.line.fill.background()
        txH = slide.shapes.add_textbox(Inches(x + 0.1), Inches(y + 0.05),
                                        Inches(bar_w - 0.2), Inches(header_h - 0.1))
        p_h = txH.text_frame.paragraphs[0]
        set_text(p_h, phase["header"], Pt(14), WHITE, bold=True); p_h.alignment = PP_ALIGN.CENTER
        body_y = y + header_h; body_h = bar_h - header_h
        bdy = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(x), Inches(body_y),
                                      Inches(bar_w), Inches(body_h))
        bdy.fill.solid(); bdy.fill.fore_color.rgb = LIGHT_GRAY
        bdy.line.color.rgb = BORDER_GRAY; bdy.line.width = Emu(9525)
        txB = slide.shapes.add_textbox(Inches(x + 0.1), Inches(body_y + 0.1),
                                        Inches(bar_w - 0.2), Inches(body_h - 0.2))
        tf = txB.text_frame; tf.word_wrap = True; tf.auto_size = MSO_AUTO_SIZE.NONE
        body_text = phase.get("body", "")
        if isinstance(body_text, list):
            for j, item in enumerate(body_text):
                p = tf.paragraphs[0] if j == 0 else tf.add_paragraph()
                add_rich_runs(p, f"・{item}", base_size=Pt(11), base_color=DARK_GRAY,
                              bold_color=NAVY, line_spacing=1.3)
    if source: add_source_label(slide, source)
    add_bottom_bar_and_footer(slide, page_num)
    return slide

def add_timeline_slide(prs, title, sub_message, events, blank, source=None, page_num=None):
    slide = prs.slides.add_slide(blank)
    sub_y = add_title_shape(slide, title)
    content_y = add_sub_message(slide, sub_message, y=sub_y)
    n = len(events); total_w = 11.5; start_x = 0.9
    line_y = content_y + 1.8
    line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(start_x), Inches(line_y),
                                   Inches(total_w), Emu(19050))
    line.fill.solid(); line.fill.fore_color.rgb = NAVY; line.line.fill.background()
    step = total_w / max(n - 1, 1) if n > 1 else total_w
    for i, ev in enumerate(events):
        x = start_x + i * step; color = ev.get("color", ACCENT)
        dot = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(x - 0.17), Inches(line_y - 0.17),
                                      Inches(0.4), Inches(0.4))
        dot.fill.solid(); dot.fill.fore_color.rgb = color; dot.line.fill.background()
        txY = slide.shapes.add_textbox(Inches(x - 0.4), Inches(line_y - 0.65),
                                        Inches(1.0), Inches(0.3))
        p = txY.text_frame.paragraphs[0]
        set_text(p, ev["year"], Pt(11), NAVY, bold=True); p.alignment = PP_ALIGN.CENTER
        if i % 2 == 0:
            ty = line_y + 0.5
        else:
            ty = line_y - 1.5
        txE = slide.shapes.add_textbox(Inches(x - 1.0), Inches(ty),
                                        Inches(2.2), Inches(0.9))
        tf = txE.text_frame; tf.word_wrap = True
        p2 = tf.paragraphs[0]
        set_text(p2, ev["title"], Pt(10), DARK_GRAY); p2.alignment = PP_ALIGN.CENTER
    if source: add_source_label(slide, source)
    add_bottom_bar_and_footer(slide, page_num)
    return slide

def add_toc_slide(prs, title, items, blank, page_num=None):
    slide = prs.slides.add_slide(blank)
    sub_y = add_title_shape(slide, title)
    n = len(items); table_x, table_y, table_w = 1.5, sub_y + 0.1, 10.3
    row_h = min(0.5, (6.4 - table_y) / max(n, 1))
    for i, item in enumerate(items):
        y = table_y + i * row_h
        if i % 2 == 0:
            bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(table_x), Inches(y),
                                         Inches(table_w), Inches(row_h))
            bg.fill.solid(); bg.fill.fore_color.rgb = LIGHT_GRAY; bg.line.fill.background()
        txNum = slide.shapes.add_textbox(Inches(table_x + 0.2), Inches(y + 0.05),
                                          Inches(0.8), Inches(row_h - 0.1))
        set_text(txNum.text_frame.paragraphs[0], f"{item.get('num', i+1)}.", Pt(14), NAVY, bold=True)
        txTitle = slide.shapes.add_textbox(Inches(table_x + 1.2), Inches(y + 0.05),
                                            Inches(7.0), Inches(row_h - 0.1))
        set_text(txTitle.text_frame.paragraphs[0], item["title"], Pt(14), DARK_GRAY, bold=True)
        txPage = slide.shapes.add_textbox(Inches(table_x + 8.5), Inches(y + 0.05),
                                           Inches(1.5), Inches(row_h - 0.1))
        p = txPage.text_frame.paragraphs[0]
        set_text(p, item.get("page", ""), Pt(14), MEDIUM_GRAY); p.alignment = PP_ALIGN.RIGHT
    add_bottom_bar_and_footer(slide, page_num)
    return slide

def add_closing_slide(prs, report_title, blank):
    slide = prs.slides.add_slide(blank)
    bg = slide.background.fill; bg.solid(); bg.fore_color.rgb = NAVY
    txBox = slide.shapes.add_textbox(Inches(1.5), Inches(2.5), Inches(10), Inches(1.5))
    p = txBox.text_frame.paragraphs[0]
    set_text(p, "Thank You", Pt(48), WHITE, bold=True); p.alignment = PP_ALIGN.CENTER
    txBox2 = slide.shapes.add_textbox(Inches(1.5), Inches(4.2), Inches(10), Inches(1))
    p2 = txBox2.text_frame.paragraphs[0]
    set_text(p2, report_title, Pt(16), RGBColor(0xAA, 0xBB, 0xDD)); p2.alignment = PP_ALIGN.CENTER
    txBox3 = slide.shapes.add_textbox(Inches(1.5), Inches(5.5), Inches(10), Inches(0.5))
    p3 = txBox3.text_frame.paragraphs[0]
    set_text(p3, "APOLLO Patent Analytics Platform", Pt(12), RGBColor(0x88, 0x99, 0xBB))
    p3.alignment = PP_ALIGN.CENTER
    bot = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(7.1),
                                  Inches(13.33), Emu(27432))
    bot.fill.solid(); bot.fill.fore_color.rgb = ACCENT; bot.line.fill.background()
    return slide

# ============ スライド構築 ============
prs = Presentation(TEMPLATE)
blank = prs.slide_layouts[6]

# 1. 表紙
add_title_slide(prs, "炭酸カルシウム技術動向分析 2026",
                "炭酸カルシウム関連特許 1,959 件に見る成長領域と出願人動態（1911-2026）",
                "Mission: 最近の注目領域について客観的な分析   |   2026年5月", blank)

# 2. 目次
add_toc_slide(prs, "目次", [
    {"num": 1, "title": "エグゼクティブサマリー", "page": "P3"},
    {"num": 2, "title": "技術環境・ライフサイクル分析（NEBULA）", "page": "P6"},
    {"num": 3, "title": "基本統計・競争構造（ATLAS）", "page": "P9"},
    {"num": 4, "title": "技術ランドスケープ（Saturn V）", "page": "P13"},
    {"num": 5, "title": "出願人動態分析（MEGA）", "page": "P18"},
    {"num": 6, "title": "キーワード共起構造（Explorer）", "page": "P21"},
    {"num": 7, "title": "仮説検証と戦略提言", "page": "P23"},
], blank, page_num=2)

# 3. セクション1: エグゼクティブサマリー
add_section_slide(prs, 1, "エグゼクティブサマリー", blank,
                  subtitle="本分析の核心と主要 KPI")

# 4. KPIダッシュボード
add_kpi_slide(prs, "本母集団の核心指標 ～CO₂固定化クラスタが直近5年集中率39%で成長を牽引",
              "**全 1,959 件**の本母集団から検出された 35 クラスタのうち、CO₂ 固定化が最大の注目領域",
              [
                  {"label": "対象件数", "value": "1,959", "unit": "件（1911-2026）"},
                  {"label": "検出クラスタ数", "value": "35", "unit": "Saturn V"},
                  {"label": "ノイズ率", "value": "11.5", "unit": "% 標準・安定構造"},
                  {"label": "CO₂固定化 直近5年集中率", "value": "39", "unit": "% 全体平均の約5倍"},
                  {"label": "出願人 HHI", "value": "0.063", "unit": "競争的・分散"},
                  {"label": "最高成長 出願人", "value": "住友大阪セメント", "unit": "+13.7% MEGAリーダー"},
              ], blank, source="ATLAS / Saturn V / MEGA PULSE 分析", page_num=3)

# 5. 主要発見カード
add_cards_slide(prs, "主要発見 ～成長軸はCO₂固定化、伝統用途は構造的縮小",
                "本母集団の構造転換を 4 つの観点で要約。複数モジュールが同一結論を裏付け",
                [
                    {"header": "注目領域", "color": ACCENT,
                     "body": "**CO₂ 固定化クラスタ**が99件中39件を直近5年に集中。Saturn V・Explorer・CORE の3手法が一致して中心テーマと検出"},
                    {"header": "成長の担い手", "color": GREEN_ACCENT,
                     "body": "**住友大阪セメント**(+13.7%)と**神島化学**(12/13件が直近)。セメント・建設・大学の産学連携"},
                    {"header": "衰退する軸", "color": RED_ACCENT,
                     "body": "**丸尾117件・オムヤ126件・日本製紙69件**は成熟/衰退象限。製紙特化日本製紙の直近出願は**ゼロ**"},
                    {"header": "政策連動", "color": NAVY,
                     "body": "2022-2023年の出願再上昇は**カーボンニュートラル政策**(2020年宣言)と時系列整合"},
                ], blank, source="統合分析（クロスモジュール P1/P8/P11）", page_num=4)

# 6. 構造転換のサマリー（ナラティブ）
add_chart_text_slide(prs,
    "本母集団は二重ライフサイクルの転換期 ～伝統用途縮小 × CO₂固定化勃興",
    "成熟技術領域における **第二の成長曲線** が立ち上がる構造転換期にある",
    "voyager_ev1_0.png",
    [
        "**伝統用途（填料・顔料・製紙）**は 2004 年ピーク(53件)から縮小局面",
        "**2018-2021年**に年27-30件へ低下、長期成熟化が進行",
        "**2022-2023年に再上昇**(28→45件)、CO₂ 固定化が牽引",
        "**直近2年(2025-26)は未完年**のため絶対値は保守的",
        "二重ライフサイクル：**衰退する伝統用途 × 勃興するCO₂固定化**",
    ], blank, chart_label="本母集団の出願トレンド全景（1911-2026）",
    source="ATLAS 基本統計分析", page_num=5)

# 7. セクション2: NEBULA環境分析
add_section_slide(prs, 2, "技術環境・ライフサイクル分析（NEBULA）", blank,
                  subtitle="本母集団の長期動向と外部政策の影響")

# 8. NEBULA ハイプサイクル
add_chart_text_slide(prs,
    "技術ライフサイクルは第二の成長曲線へ ～2022-23年に出願1.8倍へ再上昇",
    "1946 年以来 80 年の技術領域が、CO₂ 固定化という新軸で再活性化",
    "nebula_hype.png",
    [
        "**第1期**(1946-70年代前半): 萌芽期、年1-13件",
        "**第2期**(70-90年代前半): 拡大期、1991年ピーク**44件**",
        "**第3期**(90後半-10年代前半): 成熟期、2004年**53件**で高原",
        "**第4期**(現在): 転換期、2023年**45件**に再上昇",
        "再上昇は **CO₂ 固定化クラスタ**の急増(2019:6件→2023:18件)が牽引",
    ], blank, chart_label="特許出願トレンドの長期推移",
    source="NEBULA 技術ライフサイクル分析", page_num=6)

# 9. 政策タイムライン
add_timeline_slide(prs,
    "カーボンニュートラル政策が成長軸を駆動 ～時系列整合は明確",
    "外部政策イベントと本母集団の変曲点が高い整合性を示す",
    [
        {"year": "1991", "title": "第2期ピーク\n年44件達成", "color": MEDIUM_GRAY},
        {"year": "2004", "title": "第3期最高水準\n年53件", "color": BLUE},
        {"year": "2018-21", "title": "出願低下\n年27-30件", "color": AMBER},
        {"year": "2020", "title": "CN宣言\n政策ドライバー", "color": GREEN_ACCENT},
        {"year": "2021", "title": "グリーン成長戦略\nカーボンリサイクル重点化", "color": GREEN_ACCENT},
        {"year": "2023", "title": "出願再上昇\n45件・CO₂固定化主導", "color": ACCENT},
    ], blank, source="NEBULA × ATLAS クロス分析", page_num=7)

# 10. セクション3: ATLAS基本統計
add_section_slide(prs, 3, "基本統計・競争構造分析（ATLAS）", blank,
                  subtitle="多様性指標と出願人ランキングで二層構造を可視化")

# 11. ATLAS 出願トレンド
add_chart_text_slide(prs,
    "全期間CAGR +1.0%の平均値は実像を隠す ～衰退と急成長が相殺",
    "本母集団の +1.0% は **衰退する伝統用途** と **急成長する CO₂ 固定化** の相殺結果",
    "atlas_trend_snap.png",
    [
        "全期間 CAGR **+1.0%**（1911-2026 年、増加傾向）",
        "**填料・顔料クラスタ**(成熟・衰退、計732件)は緩やか縮小",
        "**CO₂固定化クラスタ**(99件中39件直近)は直近集中率**39%**",
        "平均値ではなく**クラスタ単位**の成長率の分散に着目が必要",
        "注目領域特定には**動態的視点**が不可欠",
    ], blank, chart_label="出願件数の長期推移（棒グラフ）",
    source="ATLAS 時系列分析", page_num=8)

# 12. ATLAS 多様性指標
add_progress_bar_slide(prs,
    "多様性指標が示す二層構造 ～競争的だが上位に伝統大手が蓄積",
    "**HHI 0.063 (競争的)** × **Gini 0.43 (上位の偏り)** の組み合わせが構造転換を映す",
    [
        {"label": "HHI（集中度）", "value": 6.3, "max_value": 30, "unit": "% 競争的・分散(<10)", "color": GREEN_ACCENT},
        {"label": "Entropy（多様性）", "value": 4.43, "max_value": 5.0, "unit": " 高多様性", "color": ACCENT},
        {"label": "Gini（不平等度）", "value": 43, "max_value": 100, "unit": "% 中程度の不平等", "color": AMBER},
        {"label": "上位5社シェア合計", "value": 22.5, "max_value": 100, "unit": "% 寡占的でない", "color": BLUE},
    ], blank, source="ATLAS 多様性指標（HHI / Entropy / Gini）", page_num=9)

# 13. ATLAS 出願人ランキングテーブル
add_table_slide(prs,
    "件数上位は伝統大手、成長上位は新興プレイヤー ～「上位＝成熟、中位＝成長」の逆転構造",
    "件数ランキングと MEGA 動態象限の **乖離** こそが本母集団の最大の特徴",
    ["順位", "出願人", "件数", "MEGA 象限 / 性格"],
    [
        ["1", "丸尾カルシウム", "117", "成熟・既存（填料・顔料専業）"],
        ["2", "オムヤ・インターナショナル", "126", "衰退・ニッチ（海外填料大手）"],
        ["3", "日本製紙", "69", "衰退・ニッチ（製紙用途・直近ゼロ）"],
        ["4", "白石工業", "61", "成熟・既存（粒子制御専業）"],
        ["5", "奥多摩工業", "57", "衰退・ニッチ（石灰原料）"],
        ["6", "白石中央研究所", "43", "リーダー（粒子制御研究拠点）"],
        ["9", "宇部マテリアルズ", "27", "リーダー（軽質炭酸）"],
        ["10", "住友大阪セメント", "24", "★ リーダー（CO₂固定化、+13.7%）"],
        ["—", "神島化学工業", "13", "★ リーダー（新興・直近12件集中）"],
    ], blank, col_widths=[0.7, 3.5, 1.0, 7.1], highlight_rows=[5, 6, 7, 8],
    annotations=[
        "件数 **1-5位** は伝統的填料・顔料・原料メーカー",
        "件数 **6,9,10位** がリーダー象限（成長を担う）",
        "**住友大阪セメント**(10位)が最高成長率 **+13.7%**",
        "**神島化学**はランキング外から急浮上",
    ], source="ATLAS 基本統計 + MEGA PULSE 象限分類", page_num=10)

# 14. セクション4: Saturn V
add_section_slide(prs, 4, "技術ランドスケープ分析（Saturn V）", blank,
                  subtitle="35クラスタを5超領域に集約、CO₂固定化が成長エンジン")

# 15. Saturn V 全体俯瞰
add_chart_text_slide(prs,
    "35クラスタを5超領域に集約 ～最大はCO₂固定化×粒子制御の連鎖",
    "Saturn V TELESCOPE 分析（SBERT + UMAP + HDBSCAN）で技術空間の骨格を可視化",
    "saturn_main_snap.png",
    [
        "**メガクラスタ2つ**: クラスタ30 基本製造(257件) / クラスタ9 CO₂固定化(99件)",
        "**ノイズ率 11.5%** は標準・安定構造、新興は既存クラスタ内部に集中",
        "**5 超領域**: A CO₂固定化 / B 粒子制御 / C 填料樹脂 / D 資源循環 / E 食品医療",
        "蓄積型(クラスタ30) vs 加速型(クラスタ9) の **二極構造**",
        "意味ベクトル分析が CO₂ 固定化を中心テーマと検出",
    ], blank, chart_label="技術ランドスケープ全体（UMAP 2D 投影）",
    source="Saturn V TELESCOPE 分析", page_num=11)

# 16. 5超領域カード
add_cards_slide(prs,
    "5つの技術超領域 ～CO₂固定化と資源循環が成長軸の核",
    "UMAP 空間の近接関係から、35 クラスタを 5 つの超領域に集約",
    [
        {"header": "🅰 CO₂固定化", "color": GREEN_ACCENT,
         "body": "**約180件・成長軸**\nクラスタ9・6・13・31・33・26。住友大阪セメント・神島化学・建設業・大学の異業種参入"},
        {"header": "🅱 粒子制御", "color": NAVY,
         "body": "**約700件・最大規模**\nクラスタ30・15・16・22・23。白石工業・宇部マテリアルズの専業メーカー主導"},
        {"header": "🅲 填料・樹脂", "color": AMBER,
         "body": "**約240件・成熟/衰退**\nクラスタ25・19・20・12・7。丸尾・オムヤ・日本製紙の伝統用途"},
        {"header": "🅳 資源循環", "color": ACCENT,
         "body": "**約300件・新興連結**\nクラスタ28・34・29・11・14・26。スラグ・貝殻・廃液からのカルシウム回収"},
        {"header": "🅴 機能性", "color": BLUE,
         "body": "**約100件・ニッチ**\nクラスタ3・24・32・13。食品・医療・化粧品。出光・富山大・産学連携"},
    ], blank, source="Saturn V 空間配置分析 + 超領域命名", page_num=12)

# 17. クラスタ動態 4象限マトリクス
add_matrix_slide(prs,
    "クラスタ動態 4象限 ～成長リーダー象限が779件で本母集団の主軸",
    "**成長リーダー9クラスタ・779件** が本母集団の成長を牽引",
    {
        "TL": {"title": "新興（9クラスタ・223件）",
               "items": ["[15] 多面体炭酸 33件", "[6] 脱硫 19件",
                         "[26] 石灰泥 26件", "[10] 炭酸塩 26件",
                         "小規模×相対高成長"]},
        "TR": {"title": "成長リーダー（9クラスタ・779件）",
               "items": ["[9] CO₂固定化 99件 ★", "[30] 基本製造 257件",
                         "[28] 石灰石 57件", "[34] 貝殻 58件",
                         "大規模×相対高成長"]},
        "BL": {"title": "ニッチ/衰退（8クラスタ・212件）",
               "items": ["[7] 製紙用 34件", "[8] 高級炭酸 34件",
                         "[17] 軽質炭酸 27件", "[14] パルプ廃液 15件",
                         "縮小・撤退領域"]},
        "BR": {"title": "成熟（9クラスタ・520件）",
               "items": ["[12] 軽質炭酸 52件", "[19] ゴム 35件",
                         "[27] 水性製造 45件", "[22] 紡錘状 36件",
                         "差別化困難領域"]},
    }, blank, source="Saturn V クラスタ動態マップ", page_num=13)

# 18. 注目領域の代表特許
add_table_slide(prs,
    "注目領域(CO₂固定化)の代表特許 ～セメント・建設・大学の連携が顕著",
    "**クラスタ9（CO₂固定化、99件中39件直近）** から抽出した代表特許",
    ["公開番号", "発明の名称", "出願人", "出願年"],
    [
        ["特開2025-152354", "二酸化炭素の固定化方法及び固定化装置", "住友大阪セメント", "2024"],
        ["特開2026-013631", "硝酸塩の製造方法、燃焼排ガス処理、CO₂製造", "住友大阪セメント", "2024"],
        ["特開2026-059136", "二酸化炭素の回収方法", "神島化学工業", "2024"],
        ["特開2026-083244", "二酸化炭素の回収方法", "神島化学工業", "2026"],
        ["特開2026-001556", "CO₂固定化方法及び装置", "山口大学・前田道路", "2024"],
        ["特開2025-114357", "CO₂固定化方法・装置・粉末", "ヤマハ・京都大学", "2024"],
        ["特開2024-174665", "炭酸ガス固定化×コンクリート製造", "北川鉄工所・鹿島建設", "2023"],
        ["特開2025-104398", "炭酸カルシウム製造×CO₂固定化×土木", "鴻池組・白石工業", "2023"],
        ["特表2025-533803", "鉄鋼スラグ総合利用方法", "ユアンチュウ・テクノロジー", "2023"],
    ], blank, col_widths=[1.8, 5.6, 3.2, 1.7],
    source="特許データセット条件検索（クラスタ9絞込）", page_num=14)

# 19. ノイズ詳細 + 注釈
add_chart_text_slide(prs,
    "ノイズ225件は「過去集中」型 ～新興は既存クラスタ内部の質的転換",
    "**ノイズ率 11.5% は標準・安定構造**。萌芽テーマはノイズではなく既存クラスタの再活性化",
    "saturn_drill_snap.png",
    [
        "ノイズ **225件 (11.5%)**、解釈区分「標準・安定」(5-15%)",
        "時系列パターン「**過去集中**」: 1990s-2010sに分散",
        "ノイズ上位: 丸尾18件・オムヤ14件 = 伝統大手の歴史的個別出願",
        "萌芽キーワード: **バテライト型炭酸カルシウム**(8件)、無機質成形体(8件)",
        "新興は既存クラスタ9内部での質的転換に集中",
    ], blank, chart_label="クラスタ詳細マップ（ドリルダウン）",
    source="Saturn V ノイズ分析", page_num=15)

# 20. セクション5: MEGA
add_section_slide(prs, 5, "出願人動態分析（MEGA PULSE）", blank,
                  subtitle="CAGR×活動量で出願人を4象限分類")

# 21. MEGA リーダー象限カード
add_cards_slide(prs,
    "リーダー象限の5社 ～CO₂固定化×粒子制御の両輪が成長を牽引",
    "上位30出願人中、**リーダー象限はわずか5社**。これらが本母集団の成長を一手に担う",
    [
        {"header": "住友大阪セメント", "color": GREEN_ACCENT,
         "body": "**24件・+13.7% 最高成長**\n2024年に5件加速。CO₂固定化が主戦場、燃焼排ガス処理から炭酸塩化まで一体戦略"},
        {"header": "神島化学工業", "color": GREEN_ACCENT,
         "body": "**13件中12件が直近5年**\n完全新興型。CO₂回収×無機質成形体の垂直統合で建材化"},
        {"header": "白石中央研究所", "color": ACCENT,
         "body": "**43件・+2.9%**\n白石工業グループ研究拠点。粒子制御を継続的に深掘り、伝統技術の高度化"},
        {"header": "宇部マテリアルズ", "color": ACCENT,
         "body": "**27件・+2.2%**\n軽質炭酸カルシウム製造の専業として安定成長"},
    ], blank, source="MEGA PULSE 分析（CAGR × 活動量の4象限）", page_num=16)

# 22. MEGA 衰退対照プログレスバー
add_progress_bar_slide(prs,
    "上位伝統メーカーの直近5年集中度 ～蓄積の大きさが成長を保証しない逆説",
    "件数上位の伝統大手ほど直近活動が縮小。**日本製紙はゼロ**",
    [
        {"label": "神島化学工業 (13件)", "value": 92, "color": GREEN_ACCENT},
        {"label": "住友大阪セメント (19件)", "value": 47, "color": GREEN_ACCENT},
        {"label": "白石中央研究所 (33件)", "value": 15, "color": ACCENT},
        {"label": "丸尾カルシウム (117件)", "value": 3, "color": AMBER},
        {"label": "オムヤ (126件)", "value": 3, "color": AMBER},
        {"label": "日本製紙 (69件)", "value": 0, "color": RED_ACCENT},
    ], blank, source="特許データセットの出願人別 × 直近5年集中度", page_num=17)

# 23. セクション6: Explorer
add_section_slide(prs, 6, "キーワード共起構造（Explorer）", blank,
                  subtitle="テキストマイニングによる注目領域の独立検証")

# 24. Explorer ワードクラウド
add_chart_text_slide(prs,
    "「二酸化炭素」が最高次数中心性 ～SBERTと共起分析が同一結論",
    "意味ベクトル(Saturn V)と共起頻度(Explorer)の **独立2手法** が一致して CO₂ 固定化を中心と検出",
    "explorer_global_wordcloud.png",
    [
        "ノード**58**、エッジ**84**、密度0.051(疎なネットワーク)",
        "最高中心性: **「二酸化炭素」0.1404・頻度53**",
        "次点: 「カルシウム」0.1404・「製造」「炭酸塩」0.1053",
        "**コミュニティ3** が最大(15語)、CO₂固定化群を形成",
        "顔料・塗工コミュニティは構造的に**周辺化**",
    ], blank, chart_label="共起キーワードのワードクラウド",
    source="Explorer 共起ネットワーク分析（Jaccard ≥ 0.03）", page_num=18)

# 25. Explorer トルネード
add_chart_text_slide(prs,
    "コミュニティ別の中心性 ～CO₂固定化群が15ノードで最大",
    "10 コミュニティのうち、CO₂ 固定化群が規模・中心性ともに支配的",
    "exp_tornado_snap.png",
    [
        "**C3 CO₂固定化(15語)**: 固定化方法・回収方法・装置・バテライト",
        "**C1 製造・表面処理(11語)**: 沈降炭酸・貝殻・添加剤・粉砕",
        "**C4 顔料・製紙(9語)**: 塗工紙・複合粒子（中心性低・周辺化）",
        "**C9 表面処理樹脂(6語)**: 樹脂組成物・石灰石",
        "**C6 無機質成形体(2語)**: 神島化学の建材化技術",
    ], blank, chart_label="キーワードの中心性・頻度分析",
    source="Explorer コミュニティ検出 + 中心性分析", page_num=19)

# 26. セクション7: 仮説検証と提言
add_section_slide(prs, 7, "仮説検証と戦略提言", blank,
                  subtitle="3仮説すべてが支持、本分析の結論は頑健")

# 27. 仮説検証
add_hypothesis_slide(prs,
    "仮説検証サマリー ～3仮説すべてが支持、結論の頑健性を確認",
    "NEBULA 環境分析で立てた **3 仮説すべてが Saturn V × MEGA × ATLAS の照合で支持**",
    [
        {"id": "H1", "hypothesis": "カーボンニュートラル政策が本母集団の成長を駆動",
         "verdict": "confirmed",
         "evidence": "2022-23年の出願再上昇(28→45件)が政策イベント(2020年宣言・2021年GGS)と時系列整合。Saturn V クラスタ9 の急増(2019:6→2023:18)と一致"},
        {"id": "H2", "hypothesis": "成長を牽引するのはCO₂多排出産業(セメント・建設)",
         "verdict": "confirmed",
         "evidence": "MEGAリーダー象限の筆頭が住友大阪セメント(+13.7%)・神島化学。建設業(鹿島・大成・鴻池)・大学(山口・京大・富山)の参入も顕著"},
        {"id": "H3", "hypothesis": "伝統的填料・顔料用途は構造的縮小局面",
         "verdict": "confirmed",
         "evidence": "Saturn V で填料・製紙クラスタ(7・8・17・25)が成熟/衰退象限。日本製紙の直近出願ゼロ。Explorerで顔料コミュニティが周辺化"},
    ], blank, source="クロスモジュール統合分析（P1・P8・P11）", page_num=20)

# 28. 戦略ロードマップ（ステップアップ）
add_stepup_slide(prs,
    "戦略ロードマップ ～短期はCO₂固定化集中、中期は産学連携、長期は事業転換",
    "本母集団の成長軸に乗るための **時間軸別アクション**",
    [
        {"header": "短期（1年以内）", "color": ACCENT,
         "body": ["CO₂固定化(B01D×C01F×C04B)へ出願集中",
                  "クラスタ9の99件・住友大阪セメントを起点ベンチマーク"]},
        {"header": "中期（2-3年）", "color": BLUE,
         "body": ["セメント・建設・大学との共同出願立ち上げ",
                  "CO₂由来カルシウムの高機能フィラー化(グリーンフィラー)"]},
        {"header": "中期（2-3年）", "color": NAVY,
         "body": ["資源循環(鉄鋼スラグ・貝殻)×CO₂固定化の二重循環",
                  "産学連携(産学・業際的)で基礎技術を確保"]},
        {"header": "長期（3-5年）", "color": GREEN_ACCENT,
         "body": ["伝統用途のライセンス化と事業転換",
                  "蓄積した粒子制御技術をグリーンフィラーへ転用"]},
    ], blank, source="戦略的提言（本分析の主要発見に基づく）", page_num=21)

# 29. 推奨アクション
add_cards_slide(prs,
    "推奨アクション 5項目 ～優先度高で短期から実行",
    "本分析の視座（注目領域＝CO₂固定化）から導出した具体的提言",
    [
        {"header": "[高] CO₂固定化出願集中", "color": RED_ACCENT,
         "body": "B01D×C01F×C04B のクロス IPC で出願。短期1年以内"},
        {"header": "[高] 産学・業際連携", "color": RED_ACCENT,
         "body": "セメント・鉄鋼・建設・大学との共同出願。短-中期"},
        {"header": "[中] グリーンフィラー化", "color": AMBER,
         "body": "CO₂由来CaCO₃を高機能フィラーへ統合設計。中期2-3年"},
        {"header": "[中] 資源循環の結合", "color": AMBER,
         "body": "スラグ・貝殻×CO₂固定化の二重循環。中期2-3年"},
        {"header": "[低] 伝統用途の選択撤退", "color": GREEN_ACCENT,
         "body": "填料・顔料の絞り込みとライセンス活用。長期3-5年"},
    ], blank, source="戦略的提言（5項目）", page_num=22)

# 30. クロージング
add_closing_slide(prs, "炭酸カルシウム技術動向分析 2026", blank)

# 出力
out_path = os.path.join(SESSION, "reports", "apollo_report_20260527.pptx")
prs.save(out_path)
print(f"PPTX生成完了: {out_path}")
print(f"スライド数: {len(prs.slides)}")
