# -*- coding: utf-8 -*-
"""APOLLO 単体運用 PPTX 生成スクリプト（再生材ペレット 特許分析）。

CAPCOM セッション（session_20260614_145415_3ef247）の分析結果を、
slides_spec.md v6.3 のスライド関数を用いてコンサルティング品質の
単体レポート（Deep Dive 級）として 1 つの .pptx に出力する。

配色: モノトーン＋クリムゾン単一アクセント / 見出し=明朝 / 本文=ゴシック。
"""
import os
import re
import json

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_AUTO_SIZE
from pptx.enum.shapes import MSO_SHAPE, MSO_SHAPE_TYPE
from lxml import etree
from PIL import Image

# ====== パス設定 ======
HERE = os.path.dirname(os.path.abspath(__file__))
ROOT = os.path.join(HERE, "..")
SESSION = os.path.join(ROOT, "session_work", "session_20260614_145415_3ef247")
DATA = os.path.join(SESSION, "data")
SNAP = os.path.join(SESSION, "snapshots")
TEMPLATE = os.path.join(ROOT, "capcom_schema", "templates", "apollo_template.pptx")
OUT_DIR = os.path.join(ROOT, "reports")
OUT = os.path.join(OUT_DIR, "apollo_report_20260614.pptx")
SOURCE = "特許データベース Cyber patent desc（2026-06-14）"

# ====== フォント ======
JA_FONT = "Yu Gothic"
LATIN_FONT = "Century Gothic"
HEADING_JA_FONT = "Yu Mincho"
HEADING_LATIN_FONT = "Yu Mincho"
A_NS = "http://schemas.openxmlformats.org/drawingml/2006/main"

# ====== カラー ======
INK = RGBColor(0x11, 0x11, 0x11)
PAPER = RGBColor(0xFF, 0xFF, 0xFF)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
DARK_SECTION = RGBColor(0x10, 0x10, 0x11)
DARK_GRAY = RGBColor(0x11, 0x11, 0x11)
MEDIUM_GRAY = RGBColor(0x68, 0x68, 0x68)
LIGHT_GRAY = RGBColor(0xF1, 0xF2, 0xF3)
PALE_GRAY = RGBColor(0xF6, 0xF7, 0xF8)
BORDER_GRAY = RGBColor(0xD8, 0xDA, 0xDD)
ACCENT = RGBColor(0xC5, 0x12, 0x12)
RED = ACCENT
RED_ON_DARK = RGBColor(0xFF, 0x70, 0x70)
RED_ACCENT = RGBColor(0xC5, 0x12, 0x12)
DEEP_RED = RGBColor(0x83, 0x10, 0x10)
KEY_MSG_BG = RGBColor(0xF6, 0xF7, 0xF8)
GHOST_ON_DARK = RGBColor(0x4A, 0x4A, 0x52)
# 旧名エイリアス（カテゴリ区別はグレー段階＋赤）
NAVY = INK
BLUE = MEDIUM_GRAY
GREEN_ACCENT = INK
AMBER = RGBColor(0x8F, 0x8F, 0x8F)

MARGIN_L = 0.9
CONTENT_W = 13.33 - MARGIN_L - 0.7
PAGE_NUM_Y = 7.02
PAGE_NUM_X = 12.45
PAGE_NUM_W = 0.55


# =====================================================================
# Section 2: コアユーティリティ
# =====================================================================
def _apply_font(run, heading=False):
    latin = HEADING_LATIN_FONT if heading else LATIN_FONT
    ja = HEADING_JA_FONT if heading else JA_FONT
    run.font.name = latin
    rPr = run._r.get_or_add_rPr()
    rPr.set("lang", "ja-JP")
    rPr.set("altLang", "en-US")
    ea = rPr.find(f"{{{A_NS}}}ea")
    if ea is None:
        ea = etree.SubElement(rPr, f"{{{A_NS}}}ea")
    ea.set("typeface", ja)


def _apply_kinsoku(paragraph):
    pPr = paragraph._p.get_or_add_pPr()
    pPr.set("eaLnBrk", "1")
    pPr.set("hangingPunct", "1")


def _track(run, spc):
    rPr = run._r.get_or_add_rPr()
    rPr.set("spc", str(int(spc)))


def add_rich_runs(paragraph, text, base_size=Pt(14), base_color=DARK_GRAY,
                  bold_color=None, force_bold=False, line_spacing=1.4, heading=False):
    paragraph.clear()
    bold_color = bold_color or base_color
    _apply_kinsoku(paragraph)
    paragraph.line_spacing = line_spacing
    for part in re.split(r"(\*\*.*?\*\*)", text):
        if not part:
            continue
        run = paragraph.add_run()
        if part.startswith("**") and part.endswith("**"):
            run.text = part[2:-2]
            run.font.bold = True
            run.font.color.rgb = bold_color
        else:
            run.text = part
            run.font.bold = force_bold
            run.font.color.rgb = bold_color if force_bold else base_color
        run.font.size = base_size
        _apply_font(run, heading=heading)


def set_text(p, text, size, color, bold=False, line_spacing=None, heading=False):
    p.text = ""
    run = p.add_run()
    run.text = text
    run.font.size = size
    run.font.color.rgb = color
    run.font.bold = bold
    _apply_font(run, heading=heading)
    _apply_kinsoku(p)
    if line_spacing:
        p.line_spacing = line_spacing


def add_title_shape(slide, text, x=MARGIN_L, y=0.55, w=CONTENT_W, label=None):
    text_len = len(text)
    if text_len <= 20:
        font_size, box_h = Pt(30), 0.70
    elif text_len <= 30:
        font_size, box_h = Pt(27), 0.78
    elif text_len <= 44:
        font_size, box_h = Pt(23), 0.92
    else:
        font_size, box_h = Pt(20), 1.05
    eb_y = y + 0.02
    eyebrow = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(x), Inches(eb_y), Inches(0.40), Emu(12700))
    eyebrow.fill.solid()
    eyebrow.fill.fore_color.rgb = ACCENT
    eyebrow.line.fill.background()
    cy = eb_y + 0.12
    if label:
        lab = slide.shapes.add_textbox(Inches(x), Inches(cy), Inches(w), Inches(0.28))
        set_text(lab.text_frame.paragraphs[0], label, Pt(10), ACCENT, bold=True, heading=True)
        cy += 0.30
    txBox = slide.shapes.add_textbox(Inches(x), Inches(cy), Inches(w), Inches(box_h))
    tf = txBox.text_frame
    tf.word_wrap = True
    tf.auto_size = MSO_AUTO_SIZE.NONE
    add_rich_runs(tf.paragraphs[0], text, base_size=font_size, base_color=INK,
                  bold_color=INK, force_bold=True, line_spacing=1.18, heading=True)
    return cy + box_h + 0.06


def add_sub_message(slide, message, x=MARGIN_L, y=None, w=CONTENT_W):
    if y is None:
        y = 1.00
    num_lines = max(1, -(-len(message) // 40))
    box_h = 0.20 + num_lines * 0.35
    box = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(box_h))
    box.fill.solid()
    box.fill.fore_color.rgb = KEY_MSG_BG
    box.line.fill.background()
    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(x), Inches(y), Emu(36576), Inches(box_h))
    bar.fill.solid()
    bar.fill.fore_color.rgb = ACCENT
    bar.line.fill.background()
    txBox = slide.shapes.add_textbox(Inches(x + 0.15), Inches(y + 0.08), Inches(w - 0.30), Inches(box_h - 0.16))
    tf = txBox.text_frame
    tf.word_wrap = True
    tf.auto_size = MSO_AUTO_SIZE.NONE
    p = tf.paragraphs[0]
    marker = p.add_run()
    marker.text = "■ "
    marker.font.size = Pt(15)
    marker.font.color.rgb = ACCENT
    marker.font.bold = True
    _apply_font(marker)
    for part in re.split(r"(\*\*.*?\*\*)", message):
        if not part:
            continue
        run = p.add_run()
        if part.startswith("**") and part.endswith("**"):
            run.text = part[2:-2]
            run.font.bold = True
            run.font.color.rgb = INK
        else:
            run.text = part
            run.font.color.rgb = DARK_GRAY
        run.font.size = Pt(15)
        _apply_font(run)
    _apply_kinsoku(p)
    p.line_spacing = 1.5
    return y + box_h + 0.10


def add_bottom_bar_and_footer(slide, page_num=None):
    if page_num is None:
        return
    txBox = slide.shapes.add_textbox(Inches(PAGE_NUM_X), Inches(PAGE_NUM_Y), Inches(PAGE_NUM_W), Inches(0.25))
    p = txBox.text_frame.paragraphs[0]
    set_text(p, str(page_num), Pt(10), MEDIUM_GRAY)
    p.alignment = PP_ALIGN.RIGHT


def add_source_label(slide, source_text, x=0.5, y=6.62, w=12.3):
    txBox = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(0.35))
    tf = txBox.text_frame
    tf.word_wrap = True
    tf.auto_size = MSO_AUTO_SIZE.NONE
    set_text(tf.paragraphs[0], f"（出所）{source_text}", Pt(9), MEDIUM_GRAY)


def fit_image(slide, image_path, max_x, max_y, max_w, max_h):
    if not os.path.exists(image_path):
        print(f"[WARN] image missing: {image_path}")
        return None
    img = Image.open(image_path)
    img_w, img_h = img.size
    ratio = img_h / img_w
    if max_w * ratio <= max_h:
        use_w, use_h = max_w, max_w * ratio
    else:
        use_h, use_w = max_h, max_h / ratio
    left = max_x + (max_w - use_w) / 2
    top = max_y + (max_h - use_h) / 2
    pic = slide.shapes.add_picture(image_path, Inches(left), Inches(top),
                                   width=Inches(use_w), height=Inches(use_h))
    img.close()
    return pic


def add_annotation_block(slide, bullets, x, y, w, h, font_size=14,
                         has_border=False, bg_color=None):
    if bg_color or has_border:
        box = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(h))
        if bg_color:
            box.fill.solid()
            box.fill.fore_color.rgb = bg_color
        else:
            box.fill.background()
        if has_border:
            box.line.color.rgb = BORDER_GRAY
            box.line.width = Emu(9525)
        else:
            box.line.fill.background()
    txBox = slide.shapes.add_textbox(Inches(x + 0.12), Inches(y + 0.08), Inches(w - 0.24), Inches(h - 0.16))
    tf = txBox.text_frame
    tf.word_wrap = True
    tf.auto_size = MSO_AUTO_SIZE.NONE
    for i, item in enumerate(bullets):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.space_after = Pt(6)
        marker = p.add_run()
        marker.text = "■ "
        marker.font.size = Pt(font_size)
        marker.font.color.rgb = NAVY
        _apply_font(marker)
        for part in re.split(r'(\*\*.*?\*\*)', item):
            if not part:
                continue
            run = p.add_run()
            if part.startswith('**') and part.endswith('**'):
                run.text = part[2:-2]
                run.font.bold = True
                run.font.color.rgb = NAVY
            else:
                run.text = part
                run.font.color.rgb = DARK_GRAY
            run.font.size = Pt(font_size)
            _apply_font(run)
        _apply_kinsoku(p)
        p.line_spacing = 1.45


def add_chart_label(slide, text, x, y, w=3.0, size=14, color=NAVY):
    txBox = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(0.35))
    set_text(txBox.text_frame.paragraphs[0], text, Pt(size), color, bold=True, heading=True)


def add_corner_marks(slide, slide_w_in=13.33, slide_h_in=7.5, leg=0.22, margin=0.14):
    COLOR = ACCENT
    LW = Pt(0.75)

    def _hline(x0, y0, length):
        c = slide.shapes.add_connector(1, Inches(x0), Inches(y0), Inches(x0 + length), Inches(y0))
        c.line.color.rgb = COLOR
        c.line.width = LW

    def _vline(x0, y0, length):
        c = slide.shapes.add_connector(1, Inches(x0), Inches(y0), Inches(x0), Inches(y0 + length))
        c.line.color.rgb = COLOR
        c.line.width = LW
    rx = slide_w_in - margin
    ry = margin
    _hline(rx - leg, ry, leg)
    _vline(rx, ry, leg)
    lx = margin
    ly = slide_h_in - margin
    _hline(lx, ly, leg)
    _vline(lx, ly - leg, leg)


def add_ghost_stat(slide, number_text, label_text, x=7.2, y=1.0, font_size_pt=90, color_hex='ECECEC'):
    r, g, b = int(color_hex[0:2], 16), int(color_hex[2:4], 16), int(color_hex[4:6], 16)
    ghost_color = RGBColor(r, g, b)
    txBig = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(4.5), Inches(1.6))
    p = txBig.text_frame.paragraphs[0]
    p.alignment = PP_ALIGN.LEFT
    run = p.add_run()
    run.text = number_text
    run.font.size = Pt(font_size_pt)
    run.font.bold = True
    run.font.color.rgb = ghost_color
    run.font.name = 'Century Gothic'
    txLbl = slide.shapes.add_textbox(Inches(x + 0.05), Inches(y + 1.45), Inches(4.0), Inches(0.35))
    p2 = txLbl.text_frame.paragraphs[0]
    p2.alignment = PP_ALIGN.LEFT
    run2 = p2.add_run()
    run2.text = label_text
    run2.font.size = Pt(10)
    run2.font.color.rgb = RGBColor(0xBD, 0xBD, 0xBD)
    run2.font.name = 'Yu Gothic'


# ===== shadows / audit =====
def disable_all_shadows(prs):
    def _kill(shp):
        try:
            shp.shadow.inherit = False
        except Exception:
            pass
        try:
            spPr = shp._element.spPr
            if spPr is not None:
                for el in spPr.findall(f"{{{A_NS}}}effectLst"):
                    spPr.remove(el)
                etree.SubElement(spPr, f"{{{A_NS}}}effectLst")
        except Exception:
            pass
        try:
            P_NS = "http://schemas.openxmlformats.org/presentationml/2006/main"
            style = shp._element.find(f"{{{P_NS}}}style")
            if style is not None:
                eff = style.find(f"{{{A_NS}}}effectRef")
                if eff is not None:
                    eff.set("idx", "0")
        except Exception:
            pass

    def _walk(shapes):
        for shp in shapes:
            _kill(shp)
            if getattr(shp, "shape_type", None) == MSO_SHAPE_TYPE.GROUP:
                _walk(shp.shapes)
    for slide in prs.slides:
        _walk(slide.shapes)


def audit_deck(prs):
    issues = []
    for i, sl in enumerate(prs.slides, 1):
        texts = [sh.text_frame.text.strip() for sh in sl.shapes
                 if sh.has_text_frame and sh.text_frame.text.strip()]
        has_pic = any(getattr(sh, "shape_type", None) == MSO_SHAPE_TYPE.PICTURE for sh in sl.shapes)
        if not texts and not has_pic:
            issues.append(f"  スライド{i}: テキストも画像もなし")
        elif len(texts) == 1 and len(texts[0]) < 5 and not has_pic:
            issues.append(f"  スライド{i}: タイトルのみ")
    if issues:
        print(f"[audit_deck] {len(issues)}件の注意:")
        for w in issues:
            print(w)
    else:
        print("[audit_deck] OK — 全スライド診断クリア")
    return issues


# =====================================================================
# Section 3: スライドタイプ
# =====================================================================
def add_title_slide(prs, title, subtitle, date, blank, kicker="TECHNOLOGY INTELLIGENCE REPORT"):
    slide = prs.slides.add_slide(blank)
    bg = slide.background.fill
    bg.solid()
    bg.fore_color.rgb = DARK_SECTION
    ghost = slide.shapes.add_textbox(Inches(0.7), Inches(4.55), Inches(13.0), Inches(2.6))
    gp = ghost.text_frame.paragraphs[0]
    ghost.text_frame.word_wrap = False
    ghost.text_frame.auto_size = MSO_AUTO_SIZE.NONE
    grun = gp.add_run()
    grun.text = "APOLLO"
    grun.font.size = Pt(150)
    grun.font.bold = True
    grun.font.color.rgb = RGBColor(0x1B, 0x1B, 0x1F)
    _apply_font(grun, heading=True)
    _track(grun, 400)
    strip = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(0.30), Inches(7.5))
    strip.fill.solid()
    strip.fill.fore_color.rgb = ACCENT
    strip.line.fill.background()
    kick = slide.shapes.add_textbox(Inches(1.15), Inches(0.95), Inches(11), Inches(0.4))
    kick.text_frame.word_wrap = True
    kick.text_frame.auto_size = MSO_AUTO_SIZE.NONE
    krun = kick.text_frame.paragraphs[0].add_run()
    krun.text = kicker
    krun.font.size = Pt(12)
    krun.font.bold = True
    krun.font.color.rgb = RED_ON_DARK
    _apply_font(krun, heading=True)
    _track(krun, 280)
    wm = slide.shapes.add_textbox(Inches(1.15), Inches(1.42), Inches(8), Inches(0.5))
    wm.text_frame.word_wrap = True
    wm.text_frame.auto_size = MSO_AUTO_SIZE.NONE
    wrun = wm.text_frame.paragraphs[0].add_run()
    wrun.text = "APOLLO"
    wrun.font.size = Pt(20)
    wrun.font.bold = True
    wrun.font.color.rgb = WHITE
    _apply_font(wrun, heading=True)
    _track(wrun, 600)
    tlen = len(title)
    if tlen <= 16:
        t_size, t_y = Pt(48), 2.95
    elif tlen <= 28:
        t_size, t_y = Pt(40), 2.85
    else:
        t_size, t_y = Pt(34), 2.75
    line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(1.18), Inches(t_y - 0.40), Inches(3.4), Inches(0.11))
    line.fill.solid()
    line.fill.fore_color.rgb = ACCENT
    line.line.fill.background()
    txBox = slide.shapes.add_textbox(Inches(1.15), Inches(t_y), Inches(11.4), Inches(2.4))
    txBox.text_frame.word_wrap = True
    txBox.text_frame.auto_size = MSO_AUTO_SIZE.NONE
    set_text(txBox.text_frame.paragraphs[0], title, t_size, WHITE, bold=True, line_spacing=1.16, heading=True)
    txBox2 = slide.shapes.add_textbox(Inches(1.18), Inches(5.55), Inches(11), Inches(0.6))
    txBox2.text_frame.word_wrap = True
    txBox2.text_frame.auto_size = MSO_AUTO_SIZE.NONE
    set_text(txBox2.text_frame.paragraphs[0], subtitle, Pt(15), RGBColor(0xC2, 0xC2, 0xC6), heading=True)
    band = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(6.78), Inches(13.33), Inches(0.46))
    band.fill.solid()
    band.fill.fore_color.rgb = ACCENT
    band.line.fill.background()
    meta = slide.shapes.add_textbox(Inches(1.15), Inches(6.85), Inches(11.0), Inches(0.34))
    meta.text_frame.word_wrap = True
    meta.text_frame.auto_size = MSO_AUTO_SIZE.NONE
    mrun = meta.text_frame.paragraphs[0].add_run()
    mrun.text = date
    mrun.font.size = Pt(11)
    mrun.font.bold = True
    mrun.font.color.rgb = WHITE
    _apply_font(mrun, heading=True)
    _track(mrun, 120)
    return slide


def add_section_slide(prs, section_num, title, blank, subtitle=None):
    slide = prs.slides.add_slide(blank)
    bg = slide.background.fill
    bg.solid()
    bg.fore_color.rgb = DARK_SECTION
    ghost = slide.shapes.add_textbox(Inches(5.2), Inches(1.05), Inches(7.6), Inches(5.4))
    tf_g = ghost.text_frame
    tf_g.word_wrap = False
    tf_g.auto_size = MSO_AUTO_SIZE.NONE
    p_g = tf_g.paragraphs[0]
    p_g.alignment = PP_ALIGN.RIGHT
    run_g = p_g.add_run()
    run_g.text = f"{section_num:02d}"
    run_g.font.size = Pt(300)
    run_g.font.color.rgb = GHOST_ON_DARK
    run_g.font.bold = True
    _apply_font(run_g, heading=True)
    strip = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(0.30), Inches(7.5))
    strip.fill.solid()
    strip.fill.fore_color.rgb = ACCENT
    strip.line.fill.background()
    eb = slide.shapes.add_textbox(Inches(1.15), Inches(2.75), Inches(6), Inches(0.4))
    eb.text_frame.word_wrap = True
    eb.text_frame.auto_size = MSO_AUTO_SIZE.NONE
    erun = eb.text_frame.paragraphs[0].add_run()
    erun.text = f"SECTION {section_num:02d}"
    erun.font.size = Pt(15)
    erun.font.bold = True
    erun.font.color.rgb = RED_ON_DARK
    _apply_font(erun, heading=True)
    _track(erun, 260)
    line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(1.18), Inches(3.32), Inches(2.8), Inches(0.10))
    line.fill.solid()
    line.fill.fore_color.rgb = ACCENT
    line.line.fill.background()
    txBox = slide.shapes.add_textbox(Inches(1.15), Inches(3.55), Inches(8.0), Inches(1.7))
    txBox.text_frame.word_wrap = True
    txBox.text_frame.auto_size = MSO_AUTO_SIZE.NONE
    set_text(txBox.text_frame.paragraphs[0], title, Pt(40), WHITE, bold=True, line_spacing=1.15, heading=True)
    if subtitle:
        txBox2 = slide.shapes.add_textbox(Inches(1.18), Inches(5.35), Inches(7.6), Inches(0.8))
        txBox2.text_frame.word_wrap = True
        txBox2.text_frame.auto_size = MSO_AUTO_SIZE.NONE
        set_text(txBox2.text_frame.paragraphs[0], subtitle, Pt(15), RGBColor(0xC2, 0xC2, 0xC6), heading=True)
    bot = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(7.1), Inches(13.33), Emu(27432))
    bot.fill.solid()
    bot.fill.fore_color.rgb = ACCENT
    bot.line.fill.background()
    return slide


def add_chapter_intro_slide(prs, eyebrow, title, bullets, blank, source=None, page_num=None):
    slide = prs.slides.add_slide(blank)
    col_l_x = MARGIN_L
    col_l_w = 4.6
    col_r_x = MARGIN_L + col_l_w + 0.6
    col_r_w = 13.33 - col_r_x - 0.7
    n = max(1, len(bullets))
    block_h = min(4.6, n * 0.62 + 0.35)
    top_y = max(1.3, (7.0 - block_h) / 2.0)
    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(col_l_x), Inches(top_y + 0.05),
                                 Inches(0.05), Inches(block_h - 0.1))
    bar.fill.solid()
    bar.fill.fore_color.rgb = ACCENT
    bar.line.fill.background()
    eb = slide.shapes.add_textbox(Inches(col_l_x + 0.22), Inches(top_y), Inches(col_l_w), Inches(0.35))
    eb.text_frame.word_wrap = True
    eb.text_frame.auto_size = MSO_AUTO_SIZE.NONE
    set_text(eb.text_frame.paragraphs[0], eyebrow, Pt(11), ACCENT, bold=True, heading=True)
    hd = slide.shapes.add_textbox(Inches(col_l_x + 0.22), Inches(top_y + 0.42), Inches(col_l_w), Inches(block_h - 0.42))
    htf = hd.text_frame
    htf.word_wrap = True
    htf.auto_size = MSO_AUTO_SIZE.NONE
    set_text(htf.paragraphs[0], title, Pt(24), INK, bold=True, line_spacing=1.22, heading=True)
    bd = slide.shapes.add_textbox(Inches(col_r_x), Inches(top_y), Inches(col_r_w), Inches(block_h))
    btf = bd.text_frame
    btf.word_wrap = True
    btf.auto_size = MSO_AUTO_SIZE.NONE
    for j, line in enumerate(bullets):
        p = btf.paragraphs[0] if j == 0 else btf.add_paragraph()
        p.space_after = Pt(8)
        mk = p.add_run()
        mk.text = "▪ "
        mk.font.size = Pt(14)
        mk.font.color.rgb = ACCENT
        mk.font.bold = True
        _apply_font(mk)
        add_rich_runs(p, line, base_size=Pt(14), base_color=DARK_GRAY, bold_color=INK, line_spacing=1.4)
    if source:
        add_source_label(slide, source)
    add_bottom_bar_and_footer(slide, page_num)
    return slide


def add_toc_slide(prs, title, items, blank, page_num=None):
    slide = prs.slides.add_slide(blank)
    sub_y = add_title_shape(slide, title)
    n = len(items)
    table_x, table_y, table_w = 1.5, sub_y + 0.15, 10.3
    row_h = min(0.5, (6.4 - table_y) / max(n, 1))
    for i, item in enumerate(items):
        y = table_y + i * row_h
        if i % 2 == 0:
            bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(table_x), Inches(y), Inches(table_w), Inches(row_h))
            bg.fill.solid()
            bg.fill.fore_color.rgb = LIGHT_GRAY
            bg.line.fill.background()
        txNum = slide.shapes.add_textbox(Inches(table_x + 0.2), Inches(y + 0.05), Inches(0.8), Inches(row_h - 0.1))
        set_text(txNum.text_frame.paragraphs[0], f"{item.get('num', i+1)}.", Pt(14), NAVY, bold=True)
        txTitle = slide.shapes.add_textbox(Inches(table_x + 1.2), Inches(y + 0.05), Inches(7.0), Inches(row_h - 0.1))
        set_text(txTitle.text_frame.paragraphs[0], item["title"], Pt(14), DARK_GRAY, bold=True)
        txPage = slide.shapes.add_textbox(Inches(table_x + 8.5), Inches(y + 0.05), Inches(1.5), Inches(row_h - 0.1))
        p = txPage.text_frame.paragraphs[0]
        set_text(p, item.get("page", ""), Pt(14), MEDIUM_GRAY)
        p.alignment = PP_ALIGN.RIGHT
    add_bottom_bar_and_footer(slide, page_num)
    return slide


def add_chart_text_slide(prs, title, sub_message, image_path, annotations, blank,
                         caption=None, chart_label=None, text_side="right",
                         chart_ratio=0.66, source=None, page_num=None):
    slide = prs.slides.add_slide(blank)
    sub_y = add_title_shape(slide, title)
    content_y = add_sub_message(slide, sub_message, y=sub_y)
    content_w = CONTENT_W + 0.5
    content_x = MARGIN_L
    gap = 0.25
    chart_w = content_w * chart_ratio - gap / 2
    text_w = content_w * (1 - chart_ratio) - gap / 2
    remaining_h = 6.55 - content_y
    if text_side == "right":
        chart_x = content_x
        text_x = content_x + chart_w + gap
    else:
        text_x = content_x
        chart_x = content_x + text_w + gap
    if chart_label:
        add_chart_label(slide, chart_label, chart_x, content_y, chart_w)
        img_y = content_y + 0.35
        img_h = remaining_h - 0.65
    else:
        img_y = content_y
        img_h = remaining_h - 0.3
    full_path = os.path.join(SNAP, image_path) if not os.path.isabs(image_path) else image_path
    fit_image(slide, full_path, max_x=chart_x, max_y=img_y, max_w=chart_w, max_h=img_h)
    if caption:
        txBox = slide.shapes.add_textbox(Inches(chart_x), Inches(content_y + remaining_h - 0.22),
                                         Inches(chart_w), Inches(0.25))
        set_text(txBox.text_frame.paragraphs[0], caption, Pt(10), MEDIUM_GRAY)
        txBox.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
    add_annotation_block(slide, annotations[:5], text_x, content_y, text_w, remaining_h - 0.2)
    if source:
        add_source_label(slide, source)
    add_bottom_bar_and_footer(slide, page_num)
    add_corner_marks(slide)
    return slide


def add_kpi_slide(prs, title, sub_message, kpis, blank, source=None, page_num=None):
    slide = prs.slides.add_slide(blank)
    sub_y = add_title_shape(slide, title)
    content_y = add_sub_message(slide, sub_message, y=sub_y)
    n = len(kpis)
    available_w = 11.5
    start_x = 0.9
    gap = 0.2
    if n <= 4:
        n_cols, n_rows = n, 1
    else:
        n_cols, n_rows = min(4, (n + 1) // 2), 2
    card_w = (available_w - gap * (n_cols - 1)) / n_cols
    available_h = 6.4 - content_y
    row_gap = 0.2
    card_h = min((available_h - row_gap * (n_rows - 1)) / n_rows, 2.5)
    for idx, kpi in enumerate(kpis):
        row = idx // n_cols
        col = idx % n_cols
        x = start_x + col * (card_w + gap)
        y = content_y + row * (card_h + row_gap)
        emph = kpi.get("emphasis", False)
        card = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(x), Inches(y), Inches(card_w), Inches(card_h))
        card.fill.solid()
        card.fill.fore_color.rgb = LIGHT_GRAY
        card.line.color.rgb = BORDER_GRAY
        card.line.width = Emu(9525)
        bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(x), Inches(y), Emu(36576), Inches(card_h))
        bar.fill.solid()
        bar.fill.fore_color.rgb = ACCENT
        bar.line.fill.background()
        txL = slide.shapes.add_textbox(Inches(x + 0.18), Inches(y + 0.14), Inches(card_w - 0.3), Inches(0.25))
        set_text(txL.text_frame.paragraphs[0], kpi["label"], Pt(10), MEDIUM_GRAY, bold=True)
        txV = slide.shapes.add_textbox(Inches(x + 0.18), Inches(y + 0.45), Inches(card_w - 0.3), Inches(0.8))
        p = txV.text_frame.paragraphs[0]
        run = p.add_run()
        run.text = kpi["value"]
        run.font.size = Pt(40)
        run.font.bold = True
        run.font.color.rgb = ACCENT if emph else INK
        _apply_font(run)
        if kpi.get("trend"):
            run2 = p.add_run()
            run2.text = f" {kpi['trend']}"
            run2.font.size = Pt(13)
            run2.font.color.rgb = MEDIUM_GRAY
            _apply_font(run2)
        txU = slide.shapes.add_textbox(Inches(x + 0.18), Inches(y + card_h - 0.42), Inches(card_w - 0.3), Inches(0.3))
        set_text(txU.text_frame.paragraphs[0], kpi.get("unit", ""), Pt(9), MEDIUM_GRAY)
    if source:
        add_source_label(slide, source)
    add_bottom_bar_and_footer(slide, page_num)
    return slide


def add_insight_slide(prs, title, sub_message, layers, blank, label="考察", source=None, page_num=None):
    slide = prs.slides.add_slide(blank)
    sub_y = add_title_shape(slide, title, label=label)
    content_y = add_sub_message(slide, sub_message, y=sub_y) if sub_message else sub_y + 0.1
    n = max(1, len(layers))
    avail_h = 6.7 - content_y
    gap = 0.18
    block_h = (avail_h - gap * (n - 1)) / n
    y = content_y
    for lyr in layers:
        bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(MARGIN_L), Inches(y), Emu(36576), Inches(block_h))
        bar.fill.solid()
        bar.fill.fore_color.rgb = ACCENT
        bar.line.fill.background()
        lb = slide.shapes.add_textbox(Inches(MARGIN_L + 0.18), Inches(y), Inches(2.4), Inches(0.3))
        set_text(lb.text_frame.paragraphs[0], lyr.get("label", ""), Pt(11), ACCENT, bold=True)
        body = slide.shapes.add_textbox(Inches(MARGIN_L + 0.18), Inches(y + 0.30),
                                        Inches(CONTENT_W - 0.18), Inches(block_h - 0.32))
        tf = body.text_frame
        tf.word_wrap = True
        tf.auto_size = MSO_AUTO_SIZE.NONE
        add_rich_runs(tf.paragraphs[0], lyr.get("body", ""), base_size=Pt(13),
                      base_color=DARK_GRAY, bold_color=INK, line_spacing=1.4)
        y += block_h + gap
    if source:
        add_source_label(slide, source)
    add_bottom_bar_and_footer(slide, page_num)
    add_corner_marks(slide)
    return slide


def add_cards_slide(prs, title, sub_message, cards, blank, source=None, page_num=None):
    slide = prs.slides.add_slide(blank)
    sub_y = add_title_shape(slide, title)
    content_y = add_sub_message(slide, sub_message, y=sub_y)
    n = len(cards)
    gap = 0.25
    total_w = 12.3
    card_w = (total_w - gap * (n - 1)) / n
    card_h = 6.4 - content_y
    header_h = 0.5
    for i, card in enumerate(cards):
        x = 0.5 + i * (card_w + gap)
        color = card.get("color", ACCENT)
        hdr = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(content_y), Inches(card_w), Inches(header_h))
        hdr.fill.solid()
        hdr.fill.fore_color.rgb = color
        hdr.line.fill.background()
        txH = slide.shapes.add_textbox(Inches(x + 0.1), Inches(content_y + 0.08), Inches(card_w - 0.2), Inches(header_h - 0.12))
        set_text(txH.text_frame.paragraphs[0], card["header"], Pt(14), WHITE, bold=True)
        txH.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
        body_y = content_y + header_h + 0.1
        body_h = card_h - header_h - 0.1
        bdy = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(x), Inches(body_y), Inches(card_w), Inches(body_h))
        bdy.fill.solid()
        bdy.fill.fore_color.rgb = LIGHT_GRAY
        bdy.line.color.rgb = BORDER_GRAY
        bdy.line.width = Emu(9525)
        txB = slide.shapes.add_textbox(Inches(x + 0.14), Inches(body_y + 0.12), Inches(card_w - 0.28), Inches(body_h - 0.24))
        tf = txB.text_frame
        tf.word_wrap = True
        tf.auto_size = MSO_AUTO_SIZE.NONE
        body_text = card.get("body", "")
        if isinstance(body_text, list):
            for j, item in enumerate(body_text):
                p = tf.paragraphs[0] if j == 0 else tf.add_paragraph()
                p.space_after = Pt(4)
                add_rich_runs(p, f"・{item}", base_size=Pt(12), base_color=DARK_GRAY, bold_color=INK, line_spacing=1.35)
        else:
            add_rich_runs(tf.paragraphs[0], body_text, base_size=Pt(13), base_color=DARK_GRAY, bold_color=INK, line_spacing=1.4)
    if source:
        add_source_label(slide, source)
    add_bottom_bar_and_footer(slide, page_num)
    return slide


def add_table_slide(prs, title, sub_message, headers, rows, blank,
                    col_widths=None, highlight_rows=None, annotations=None,
                    source=None, page_num=None):
    slide = prs.slides.add_slide(blank)
    sub_y = add_title_shape(slide, title)
    content_y = add_sub_message(slide, sub_message, y=sub_y)
    n_cols = len(headers)
    n_rows = len(rows) + 1
    if annotations:
        table_w = 7.6
        text_x = 8.4
        text_w = 4.4
    else:
        table_w = 12.3
        text_x = None
        text_w = 0
    available_table_h = 6.35 - content_y
    row_h = min(0.5, max(0.32, available_table_h / n_rows))
    table_h = row_h * n_rows
    table_shape = slide.shapes.add_table(n_rows, n_cols, Inches(0.5), Inches(content_y), Inches(table_w), Inches(table_h))
    table = table_shape.table
    if col_widths:
        for i, w in enumerate(col_widths):
            table.columns[i].width = Inches(w)
    for j, header in enumerate(headers):
        cell = table.cell(0, j)
        cell.text = ""
        set_text(cell.text_frame.paragraphs[0], header, Pt(12), WHITE, bold=True)
        cell.fill.solid()
        cell.fill.fore_color.rgb = NAVY
    highlight_rows = highlight_rows or []
    for i, row in enumerate(rows):
        for j, value in enumerate(row):
            cell = table.cell(i + 1, j)
            cell.text = ""
            is_hl = i in highlight_rows
            set_text(cell.text_frame.paragraphs[0], str(value), Pt(11),
                     ACCENT if (is_hl and j == 0) else DARK_GRAY, bold=is_hl)
            if is_hl:
                cell.fill.solid()
                cell.fill.fore_color.rgb = KEY_MSG_BG
            elif i % 2 == 1:
                cell.fill.solid()
                cell.fill.fore_color.rgb = LIGHT_GRAY
    if annotations and text_x:
        remaining_h = 6.35 - content_y
        add_annotation_block(slide, annotations, text_x, content_y, text_w, remaining_h,
                             font_size=12, has_border=True, bg_color=KEY_MSG_BG)
    if source:
        add_source_label(slide, source)
    add_bottom_bar_and_footer(slide, page_num)
    add_corner_marks(slide)
    return slide


def add_matrix_slide(prs, title, sub_message, quadrants, blank,
                     x_label="→ 直近の活動量（勢い）", y_label="↑ 累積件数（規模）",
                     annotation_bullets=None, source=None, page_num=None):
    slide = prs.slides.add_slide(blank)
    sub_y = add_title_shape(slide, title)
    content_y = add_sub_message(slide, sub_message, y=sub_y)
    mx = 1.4
    my = content_y + 0.2
    mw = 5.8
    mh = 5.9 - content_y
    half_w = mw / 2
    half_h = mh / 2
    quad_colors = {"TL": MEDIUM_GRAY, "TR": ACCENT, "BL": MEDIUM_GRAY, "BR": INK}
    positions = {"TL": (mx, my), "TR": (mx + half_w, my),
                 "BL": (mx, my + half_h), "BR": (mx + half_w, my + half_h)}
    for key, pos in positions.items():
        q = quadrants.get(key, {})
        color = quad_colors[key]
        qx, qy = pos
        box = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(qx), Inches(qy), Inches(half_w - 0.05), Inches(half_h - 0.05))
        box.fill.solid()
        box.fill.fore_color.rgb = LIGHT_GRAY
        box.line.color.rgb = BORDER_GRAY
        box.line.width = Emu(9525)
        txBox = slide.shapes.add_textbox(Inches(qx + 0.12), Inches(qy + 0.1), Inches(half_w - 0.3), Inches(0.35))
        set_text(txBox.text_frame.paragraphs[0], q.get("title", ""), Pt(13), color, bold=True)
        items = q.get("items", [])
        if items:
            txBox2 = slide.shapes.add_textbox(Inches(qx + 0.15), Inches(qy + 0.5), Inches(half_w - 0.4), Inches(half_h - 0.7))
            tf = txBox2.text_frame
            tf.word_wrap = True
            for j, item in enumerate(items[:5]):
                p = tf.paragraphs[0] if j == 0 else tf.add_paragraph()
                set_text(p, f"・{item}", Pt(10.5), DARK_GRAY)
    txX = slide.shapes.add_textbox(Inches(mx + mw / 2 - 1.0), Inches(my + mh + 0.05), Inches(2.5), Inches(0.25))
    set_text(txX.text_frame.paragraphs[0], x_label, Pt(10), MEDIUM_GRAY)
    txY = slide.shapes.add_textbox(Inches(mx - 0.7), Inches(my + mh / 2 - 0.15), Inches(2.0), Inches(0.3))
    set_text(txY.text_frame.paragraphs[0], y_label, Pt(10), MEDIUM_GRAY)
    ann_x = mx + mw + 0.5
    ann_w = 13.33 - ann_x - 0.6
    if annotation_bullets and ann_w > 2.0:
        add_annotation_block(slide, annotation_bullets, ann_x, my, ann_w, mh,
                             font_size=13, has_border=True, bg_color=KEY_MSG_BG)
    if source:
        add_source_label(slide, source)
    add_bottom_bar_and_footer(slide, page_num)
    add_corner_marks(slide)
    return slide


def add_pyramid_slide(prs, title, sub_message, levels, blank, source=None, page_num=None):
    slide = prs.slides.add_slide(blank)
    sub_y = add_title_shape(slide, title)
    content_y = add_sub_message(slide, sub_message, y=sub_y)
    n = len(levels)
    total_h = 6.4 - content_y - 0.2
    level_h = total_h / n
    base_w = 10.0
    center_x = 6.66
    colors = [ACCENT, DEEP_RED, MEDIUM_GRAY, INK]
    for i, level in enumerate(levels):
        ratio_top = (i + 0.3) / n
        ratio_bot = (i + 1.3) / n
        lw = base_w * (ratio_top + ratio_bot) / 2
        lx = center_x - lw / 2
        ly = content_y + i * level_h
        color = colors[i % len(colors)]
        trap = slide.shapes.add_shape(MSO_SHAPE.TRAPEZOID, Inches(lx), Inches(ly), Inches(lw), Inches(level_h - 0.08))
        trap.fill.solid()
        trap.fill.fore_color.rgb = color
        trap.line.fill.background()
        txBox = slide.shapes.add_textbox(Inches(lx + 0.3), Inches(ly + 0.08), Inches(lw - 0.6), Inches(level_h - 0.18))
        tf = txBox.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        set_text(p, level["title"], Pt(14), WHITE, bold=True)
        if level.get("detail"):
            p2 = tf.add_paragraph()
            p2.alignment = PP_ALIGN.CENTER
            set_text(p2, level["detail"], Pt(11), WHITE)
    if source:
        add_source_label(slide, source)
    add_bottom_bar_and_footer(slide, page_num)
    return slide


def add_hypothesis_slide(prs, title, sub_message, hypotheses, blank, source=None, page_num=None):
    slide = prs.slides.add_slide(blank)
    sub_y = add_title_shape(slide, title)
    content_y = add_sub_message(slide, sub_message, y=sub_y)
    headers = ["ID", "仮説", "判定", "エビデンス"]
    n_rows = len(hypotheses) + 1
    available_h = 6.35 - content_y
    row_h = min(0.85, max(0.45, available_h / n_rows))
    table_h = row_h * n_rows
    table_shape = slide.shapes.add_table(n_rows, 4, Inches(0.5), Inches(content_y), Inches(12.3), Inches(table_h))
    table = table_shape.table
    table.columns[0].width = Inches(0.7)
    table.columns[1].width = Inches(4.3)
    table.columns[2].width = Inches(1.1)
    table.columns[3].width = Inches(6.2)
    VERDICT_MAP = {"confirmed": ("支持", INK), "rejected": ("不支持", MEDIUM_GRAY), "partially": ("一部支持", ACCENT)}
    for j, header in enumerate(headers):
        cell = table.cell(0, j)
        cell.text = ""
        set_text(cell.text_frame.paragraphs[0], header, Pt(12), WHITE, bold=True)
        cell.fill.solid()
        cell.fill.fore_color.rgb = NAVY
    for i, hyp in enumerate(hypotheses):
        vk = hyp.get("verdict", "partially")
        vlabel, vcolor = VERDICT_MAP.get(vk, ("一部支持", ACCENT))
        row_data = [hyp.get("id", ""), hyp.get("hypothesis", ""), vlabel, hyp.get("evidence", "")]
        for j, val in enumerate(row_data):
            cell = table.cell(i + 1, j)
            cell.text = ""
            if j == 2:
                set_text(cell.text_frame.paragraphs[0], val, Pt(12), WHITE, bold=True)
                cell.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
                cell.fill.solid()
                cell.fill.fore_color.rgb = vcolor
            else:
                set_text(cell.text_frame.paragraphs[0], str(val), Pt(11), DARK_GRAY)
                if i % 2 == 1:
                    cell.fill.solid()
                    cell.fill.fore_color.rgb = LIGHT_GRAY
    if source:
        add_source_label(slide, source)
    add_bottom_bar_and_footer(slide, page_num)
    return slide


def add_process_slide(prs, title, sub_message, steps, blank, source=None, page_num=None):
    slide = prs.slides.add_slide(blank)
    sub_y = add_title_shape(slide, title)
    content_y = add_sub_message(slide, sub_message, y=sub_y)
    n = len(steps)
    available_h = 6.4 - content_y
    gap = 0.15
    step_h = (available_h - gap * (n - 1) - 0.5) / n
    step_h = min(step_h, 1.4)
    if n <= 2:
        title_size, desc_size = Pt(16), Pt(14)
    elif n <= 3:
        title_size, desc_size = Pt(14), Pt(13)
    else:
        title_size, desc_size = Pt(13), Pt(12)
    header_w = 2.4
    body_w = 9.6
    colors = [DEEP_RED, ACCENT, MEDIUM_GRAY, INK]
    for i, step in enumerate(steps):
        sy = content_y + i * (step_h + gap + 0.15)
        color = colors[i % len(colors)]
        hdr = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.5), Inches(sy), Inches(header_w), Inches(step_h))
        hdr.fill.solid()
        hdr.fill.fore_color.rgb = color
        hdr.line.fill.background()
        txH = slide.shapes.add_textbox(Inches(0.6), Inches(sy + 0.12), Inches(header_w - 0.2), Inches(step_h - 0.2))
        tf_h = txH.text_frame
        tf_h.word_wrap = True
        p_h = tf_h.paragraphs[0]
        p_h.alignment = PP_ALIGN.CENTER
        set_text(p_h, step.get("phase", f"STEP {i+1}"), Pt(10), WHITE)
        p_t = tf_h.add_paragraph()
        p_t.alignment = PP_ALIGN.CENTER
        set_text(p_t, step["title"], title_size, WHITE, bold=True)
        bdy = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.5 + header_w + 0.1), Inches(sy), Inches(body_w), Inches(step_h))
        bdy.fill.solid()
        bdy.fill.fore_color.rgb = LIGHT_GRAY
        bdy.line.color.rgb = BORDER_GRAY
        bdy.line.width = Emu(9525)
        txB = slide.shapes.add_textbox(Inches(0.5 + header_w + 0.25), Inches(sy + 0.12), Inches(body_w - 0.3), Inches(step_h - 0.2))
        tf_b = txB.text_frame
        tf_b.word_wrap = True
        tf_b.auto_size = MSO_AUTO_SIZE.NONE
        add_rich_runs(tf_b.paragraphs[0], step.get("desc", ""), base_size=desc_size, base_color=DARK_GRAY, bold_color=INK)
        if i < n - 1:
            arrow_y = sy + step_h + 0.02
            arrow = slide.shapes.add_shape(MSO_SHAPE.DOWN_ARROW, Inches(1.3), Inches(arrow_y), Inches(0.5), Inches(gap + 0.05))
            arrow.fill.solid()
            arrow.fill.fore_color.rgb = color
            arrow.line.fill.background()
    if source:
        add_source_label(slide, source)
    add_bottom_bar_and_footer(slide, page_num)
    return slide


def add_recommendation_slide(prs, title, sub_message, recommendations, blank, source=None, page_num=None):
    PRIORITY_COLORS = {"高": ACCENT, "中": MEDIUM_GRAY, "低": INK}
    slide = prs.slides.add_slide(blank)
    sub_y = add_title_shape(slide, title)
    content_y = add_sub_message(slide, sub_message, y=sub_y)
    n = len(recommendations)
    available_h = 6.4 - content_y
    card_h = min(1.3, (available_h - 0.12 * (n - 1)) / n)
    for i, rec in enumerate(recommendations):
        y = content_y + i * (card_h + 0.12)
        p_color = PRIORITY_COLORS.get(rec.get("priority", "中"), MEDIUM_GRAY)
        bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.5), Inches(y), Emu(54864), Inches(card_h))
        bar.fill.solid()
        bar.fill.fore_color.rgb = p_color
        bar.line.fill.background()
        card = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.8), Inches(y), Inches(12.0), Inches(card_h))
        card.fill.solid()
        card.fill.fore_color.rgb = LIGHT_GRAY
        card.line.color.rgb = BORDER_GRAY
        card.line.width = Emu(9525)
        txBox_p = slide.shapes.add_textbox(Inches(1.0), Inches(y + 0.1), Inches(1.0), Inches(0.3))
        set_text(txBox_p.text_frame.paragraphs[0], f"[{rec.get('priority', '中')}]", Pt(10), p_color, bold=True)
        txBox_t = slide.shapes.add_textbox(Inches(2.0), Inches(y + 0.1), Inches(6.5), Inches(0.3))
        set_text(txBox_t.text_frame.paragraphs[0], rec["title"], Pt(16), INK, bold=True)
        if rec.get("timeframe"):
            txBox_tf = slide.shapes.add_textbox(Inches(9.0), Inches(y + 0.1), Inches(3.5), Inches(0.3))
            p_tf = txBox_tf.text_frame.paragraphs[0]
            set_text(p_tf, rec["timeframe"], Pt(13), MEDIUM_GRAY)
            p_tf.alignment = PP_ALIGN.RIGHT
        if rec.get("desc"):
            txBox_d = slide.shapes.add_textbox(Inches(2.0), Inches(y + 0.46), Inches(10.5), Inches(card_h - 0.52))
            tf_d = txBox_d.text_frame
            tf_d.word_wrap = True
            tf_d.auto_size = MSO_AUTO_SIZE.NONE
            add_rich_runs(tf_d.paragraphs[0], rec["desc"], base_size=Pt(12), base_color=DARK_GRAY, bold_color=INK)
    if source:
        add_source_label(slide, source)
    add_bottom_bar_and_footer(slide, page_num)
    return slide


def add_patent_micro_slide(prs, title, sub_message, patents, blank, label="代表特許", source=None, page_num=None):
    slide = prs.slides.add_slide(blank)
    sub_y = add_title_shape(slide, title, label=label)
    content_y = add_sub_message(slide, sub_message, y=sub_y) if sub_message else sub_y + 0.1
    n = max(1, len(patents))
    avail_h = 6.7 - content_y
    gap = 0.14
    row_h = (avail_h - gap * (n - 1)) / n
    y = content_y
    for p in patents:
        card = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(MARGIN_L), Inches(y), Inches(CONTENT_W), Inches(row_h))
        card.fill.solid()
        card.fill.fore_color.rgb = PALE_GRAY
        card.line.fill.background()
        head = slide.shapes.add_textbox(Inches(MARGIN_L + 0.18), Inches(y + 0.08), Inches(CONTENT_W - 0.36), Inches(0.32))
        hp = head.text_frame.paragraphs[0]
        r1 = hp.add_run(); r1.text = p.get("id", "") + "  "; r1.font.size = Pt(12)
        r1.font.bold = True; r1.font.color.rgb = ACCENT; _apply_font(r1)
        r2 = hp.add_run(); r2.text = p.get("name", ""); r2.font.size = Pt(12)
        r2.font.bold = True; r2.font.color.rgb = INK; _apply_font(r2)
        r3 = hp.add_run(); r3.text = "　（" + p.get("applicant", "") + "）"
        r3.font.size = Pt(11); r3.font.color.rgb = MEDIUM_GRAY; _apply_font(r3)
        note = slide.shapes.add_textbox(Inches(MARGIN_L + 0.18), Inches(y + 0.40), Inches(CONTENT_W - 0.36), Inches(row_h - 0.46))
        nt = note.text_frame
        nt.word_wrap = True
        nt.auto_size = MSO_AUTO_SIZE.NONE
        set_text(nt.paragraphs[0], p.get("note", ""), Pt(11), DARK_GRAY, line_spacing=1.3)
        y += row_h + gap
    if source:
        add_source_label(slide, source)
    add_bottom_bar_and_footer(slide, page_num)
    add_corner_marks(slide)
    return slide


def add_applicant_profile_slide(prs, title, sub_message, profiles, blank, label="出願人プロファイル", source=None, page_num=None):
    slide = prs.slides.add_slide(blank)
    sub_y = add_title_shape(slide, title, label=label)
    content_y = add_sub_message(slide, sub_message, y=sub_y) if sub_message else sub_y + 0.1
    n = max(1, len(profiles))
    avail_h = 6.7 - content_y
    gap = 0.25
    card_h = (avail_h - gap * (n - 1)) / n
    y = content_y
    for pr in profiles:
        card = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(MARGIN_L), Inches(y), Inches(CONTENT_W), Inches(card_h))
        card.fill.solid()
        card.fill.fore_color.rgb = LIGHT_GRAY
        card.line.fill.background()
        bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(MARGIN_L), Inches(y), Emu(36576), Inches(card_h))
        bar.fill.solid()
        bar.fill.fore_color.rgb = ACCENT
        bar.line.fill.background()
        head = slide.shapes.add_textbox(Inches(MARGIN_L + 0.2), Inches(y + 0.1), Inches(CONTENT_W - 0.4), Inches(0.32))
        hp = head.text_frame.paragraphs[0]
        rn = hp.add_run(); rn.text = pr.get("name", ""); rn.font.size = Pt(15)
        rn.font.bold = True; rn.font.color.rgb = INK; _apply_font(rn)
        rm = hp.add_run(); rm.text = "　" + pr.get("metrics", "")
        rm.font.size = Pt(11); rm.font.color.rgb = MEDIUM_GRAY; _apply_font(rm)
        body = slide.shapes.add_textbox(Inches(MARGIN_L + 0.2), Inches(y + 0.46), Inches(CONTENT_W - 0.4), Inches(card_h - 0.52))
        tf = body.text_frame
        tf.word_wrap = True
        tf.auto_size = MSO_AUTO_SIZE.NONE
        for j, line in enumerate(pr.get("lines", [])):
            p = tf.paragraphs[0] if j == 0 else tf.add_paragraph()
            rr = p.add_run(); rr.text = "・" + line; rr.font.size = Pt(12)
            rr.font.color.rgb = DARK_GRAY; _apply_font(rr); p.line_spacing = 1.3
        y += card_h + gap
    if source:
        add_source_label(slide, source)
    add_bottom_bar_and_footer(slide, page_num)
    return slide


def add_evidence_slide(prs, title, sub_message, items, blank, label="読んだ原文（要約抜粋）", source=None, page_num=None):
    slide = prs.slides.add_slide(blank)
    sub_y = add_title_shape(slide, title, label=label)
    content_y = add_sub_message(slide, sub_message, y=sub_y) if sub_message else sub_y + 0.1
    n = max(1, len(items))
    gap = 0.16
    h = (6.7 - content_y - gap * (n - 1)) / n
    y = content_y
    for it in items:
        card = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(MARGIN_L), Inches(y), Inches(CONTENT_W), Inches(h))
        card.fill.solid()
        card.fill.fore_color.rgb = PALE_GRAY
        card.line.fill.background()
        bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(MARGIN_L), Inches(y), Emu(36576), Inches(h))
        bar.fill.solid()
        bar.fill.fore_color.rgb = ACCENT
        bar.line.fill.background()
        hd = slide.shapes.add_textbox(Inches(MARGIN_L + 0.18), Inches(y + 0.06), Inches(CONTENT_W - 0.36), Inches(0.3))
        set_text(hd.text_frame.paragraphs[0], it.get("head", ""), Pt(11.5), ACCENT, bold=True)
        qt = slide.shapes.add_textbox(Inches(MARGIN_L + 0.18), Inches(y + 0.38), Inches(CONTENT_W - 0.40), Inches(h - 0.44))
        tf = qt.text_frame
        tf.word_wrap = True
        tf.auto_size = MSO_AUTO_SIZE.NONE
        set_text(tf.paragraphs[0], "「" + it.get("quote", "") + "」", Pt(9.5), DARK_GRAY, line_spacing=1.18)
        y += h + gap
    if source:
        add_source_label(slide, source)
    add_bottom_bar_and_footer(slide, page_num)
    return slide


def add_closing_slide(prs, report_title, blank, subtitle=None):
    slide = prs.slides.add_slide(blank)
    bg = slide.background.fill
    bg.solid()
    bg.fore_color.rgb = DARK_SECTION
    strip = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(0.30), Inches(7.5))
    strip.fill.solid()
    strip.fill.fore_color.rgb = ACCENT
    strip.line.fill.background()
    line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(1.18), Inches(2.55), Inches(3.4), Inches(0.11))
    line.fill.solid()
    line.fill.fore_color.rgb = ACCENT
    line.line.fill.background()
    txBox2 = slide.shapes.add_textbox(Inches(1.15), Inches(2.85), Inches(11.0), Inches(2.2))
    tf2 = txBox2.text_frame
    tf2.word_wrap = True
    set_text(tf2.paragraphs[0], report_title, Pt(30), WHITE, bold=True, line_spacing=1.25, heading=True)
    if subtitle:
        sb = slide.shapes.add_textbox(Inches(1.18), Inches(5.3), Inches(11.0), Inches(0.8))
        sb.text_frame.word_wrap = True
        set_text(sb.text_frame.paragraphs[0], subtitle, Pt(14), RGBColor(0xC2, 0xC2, 0xC6))
    bot = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(7.1), Inches(13.33), Emu(27432))
    bot.fill.solid()
    bot.fill.fore_color.rgb = ACCENT
    bot.line.fill.background()
    return slide


# =====================================================================
# デッキ構築
# =====================================================================
from deck_content import build_deck  # noqa: E402


if __name__ == "__main__":
    build_deck(globals())
