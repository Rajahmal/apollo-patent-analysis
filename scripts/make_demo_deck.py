# -*- coding: utf-8 -*-
"""APOLLO PPTX デモデッキ生成スクリプト（数ページ・新機能ショーケース）。

slides_spec.md v6.0 の関数を転記して、章扉→章導入→考察→クロージング等を
数ページで構成し、新機能（add_chapter_intro_slide / 強化版 disable_all_shadows /
audit_deck / 垂直センタリング narrative）を実機検証する。
"""
import os
import re

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_AUTO_SIZE
from pptx.enum.shapes import MSO_SHAPE, MSO_SHAPE_TYPE
from lxml import etree

# ---- フォント ----
JA_FONT = "Yu Gothic"
LATIN_FONT = "Century Gothic"
HEADING_JA_FONT = "Yu Mincho"
HEADING_LATIN_FONT = "Yu Mincho"
A_NS = "http://schemas.openxmlformats.org/drawingml/2006/main"

# ---- カラー ----
INK = RGBColor(0x11, 0x11, 0x11)
PAPER = RGBColor(0xFF, 0xFF, 0xFF)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
DARK_SECTION = RGBColor(0x10, 0x10, 0x11)
DARK_GRAY = RGBColor(0x11, 0x11, 0x11)
MEDIUM_GRAY = RGBColor(0x68, 0x68, 0x68)
LIGHT_GRAY = RGBColor(0xF1, 0xF2, 0xF3)
PALE_GRAY = RGBColor(0xF6, 0xF7, 0xF8)
BORDER_GRAY = RGBColor(0xD8, 0xDA, 0xDD)
ACCENT = RGBColor(0xC5, 0x12, 0x12)
RED = ACCENT
RED_ON_DARK = RGBColor(0xFF, 0x70, 0x70)
RED_ACCENT = RGBColor(0xC5, 0x12, 0x12)
KEY_MSG_BG = RGBColor(0xF6, 0xF7, 0xF8)
GHOST_ON_DARK = RGBColor(0x4A, 0x4A, 0x52)
NAVY = INK
BLUE = MEDIUM_GRAY
GREEN_ACCENT = INK
AMBER = RGBColor(0x8F, 0x8F, 0x8F)

MARGIN_L = 0.9
CONTENT_W = 13.33 - MARGIN_L - 0.7
PAGE_NUM_Y = 7.02
PAGE_NUM_X = 12.45
PAGE_NUM_W = 0.55

HERE = os.path.dirname(os.path.abspath(__file__))
TEMPLATE = os.path.join(HERE, "..", "capcom_schema", "templates", "apollo_template.pptx")


# ===== フォント・禁則 =====
def _apply_font(run, heading=False):
    latin = HEADING_LATIN_FONT if heading else LATIN_FONT
    ja = HEADING_JA_FONT if heading else JA_FONT
    run.font.name = latin
    rPr = run._r.get_or_add_rPr()
    rPr.set("lang", "ja-JP")
    rPr.set("altLang", "en-US")
    ea = rPr.find(f"{{{A_NS}}}ea")
    if ea is None:
        ea = etree.SubElement(rPr, f"{{{A_NS}}}ea")
    ea.set("typeface", ja)


def _apply_kinsoku(paragraph):
    pPr = paragraph._p.get_or_add_pPr()
    pPr.set("eaLnBrk", "1")
    pPr.set("hangingPunct", "1")


def add_rich_runs(paragraph, text, base_size=Pt(14), base_color=DARK_GRAY,
                  bold_color=None, force_bold=False, line_spacing=1.4, heading=False):
    paragraph.clear()
    bold_color = bold_color or base_color
    _apply_kinsoku(paragraph)
    paragraph.line_spacing = line_spacing
    for part in re.split(r"(\*\*.*?\*\*)", text):
        if not part:
            continue
        run = paragraph.add_run()
        if part.startswith("**") and part.endswith("**"):
            run.text = part[2:-2]
            run.font.bold = True
            run.font.color.rgb = bold_color
        else:
            run.text = part
            run.font.bold = force_bold
            run.font.color.rgb = bold_color if force_bold else base_color
        run.font.size = base_size
        _apply_font(run, heading=heading)


def set_text(p, text, size, color, bold=False, line_spacing=None, heading=False):
    p.text = ""
    run = p.add_run()
    run.text = text
    run.font.size = size
    run.font.color.rgb = color
    run.font.bold = bold
    _apply_font(run, heading=heading)
    _apply_kinsoku(p)
    if line_spacing:
        p.line_spacing = line_spacing


# ===== タイトル・サブメッセージ =====
def add_title_shape(slide, text, x=MARGIN_L, y=0.55, w=CONTENT_W, label=None):
    text_len = len(text)
    if text_len <= 20:
        font_size, box_h = Pt(30), 0.70
    elif text_len <= 30:
        font_size, box_h = Pt(27), 0.78
    elif text_len <= 44:
        font_size, box_h = Pt(23), 0.92
    else:
        font_size, box_h = Pt(20), 1.05

    eb_y = y + 0.02
    eyebrow = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(x), Inches(eb_y), Inches(0.40), Emu(12700))
    eyebrow.fill.solid()
    eyebrow.fill.fore_color.rgb = ACCENT
    eyebrow.line.fill.background()

    cy = eb_y + 0.12
    if label:
        lab = slide.shapes.add_textbox(Inches(x), Inches(cy), Inches(w), Inches(0.28))
        set_text(lab.text_frame.paragraphs[0], label, Pt(10), ACCENT, bold=True, heading=True)
        cy += 0.30

    txBox = slide.shapes.add_textbox(Inches(x), Inches(cy), Inches(w), Inches(box_h))
    tf = txBox.text_frame
    tf.word_wrap = True
    tf.auto_size = MSO_AUTO_SIZE.NONE
    add_rich_runs(tf.paragraphs[0], text, base_size=font_size, base_color=INK,
                  bold_color=INK, force_bold=True, line_spacing=1.18, heading=True)
    return cy + box_h + 0.06


def add_sub_message(slide, message, x=MARGIN_L, y=None, w=CONTENT_W):
    if y is None:
        y = 1.00
    num_lines = max(1, -(-len(message) // 40))
    box_h = 0.20 + num_lines * 0.35
    box = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(box_h))
    box.fill.solid()
    box.fill.fore_color.rgb = KEY_MSG_BG
    box.line.fill.background()
    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(x), Inches(y), Emu(36576), Inches(box_h))
    bar.fill.solid()
    bar.fill.fore_color.rgb = ACCENT
    bar.line.fill.background()
    txBox = slide.shapes.add_textbox(Inches(x + 0.15), Inches(y + 0.08), Inches(w - 0.30), Inches(box_h - 0.16))
    tf = txBox.text_frame
    tf.word_wrap = True
    tf.auto_size = MSO_AUTO_SIZE.NONE
    p = tf.paragraphs[0]
    marker = p.add_run()
    marker.text = "■ "
    marker.font.size = Pt(15)
    marker.font.color.rgb = ACCENT
    marker.font.bold = True
    _apply_font(marker)
    for part in re.split(r"(\*\*.*?\*\*)", message):
        if not part:
            continue
        run = p.add_run()
        if part.startswith("**") and part.endswith("**"):
            run.text = part[2:-2]
            run.font.bold = True
            run.font.color.rgb = INK
        else:
            run.text = part
            run.font.color.rgb = DARK_GRAY
        run.font.size = Pt(15)
        _apply_font(run)
    _apply_kinsoku(p)
    p.line_spacing = 1.5
    return y + box_h + 0.10


def add_bottom_bar_and_footer(slide, page_num=None):
    if page_num is None:
        return
    txBox = slide.shapes.add_textbox(Inches(PAGE_NUM_X), Inches(PAGE_NUM_Y), Inches(PAGE_NUM_W), Inches(0.25))
    p = txBox.text_frame.paragraphs[0]
    set_text(p, str(page_num), Pt(10), MEDIUM_GRAY)
    p.alignment = PP_ALIGN.RIGHT


def add_source_label(slide, source_text, x=0.5, y=6.55, w=12.3):
    txBox = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(0.35))
    tf = txBox.text_frame
    tf.word_wrap = True
    tf.auto_size = MSO_AUTO_SIZE.NONE
    set_text(tf.paragraphs[0], f"（出所）{source_text}", Pt(9), MEDIUM_GRAY)


# ===== shadows / audit =====
def disable_all_shadows(prs):
    def _kill(shp):
        try:
            shp.shadow.inherit = False
        except Exception:
            pass
        try:
            spPr = shp._element.spPr
            if spPr is not None:
                for el in spPr.findall(f"{{{A_NS}}}effectLst"):
                    spPr.remove(el)
                etree.SubElement(spPr, f"{{{A_NS}}}effectLst")
        except Exception:
            pass
        try:
            P_NS = "http://schemas.openxmlformats.org/presentationml/2006/main"
            style = shp._element.find(f"{{{P_NS}}}style")
            if style is not None:
                eff = style.find(f"{{{A_NS}}}effectRef")
                if eff is not None:
                    eff.set("idx", "0")
        except Exception:
            pass

    def _walk(shapes):
        for shp in shapes:
            _kill(shp)
            if getattr(shp, "shape_type", None) == MSO_SHAPE_TYPE.GROUP:
                _walk(shp.shapes)
    for slide in prs.slides:
        _walk(slide.shapes)


def audit_deck(prs):
    issues = []
    for i, sl in enumerate(prs.slides, 1):
        texts = [sh.text_frame.text.strip() for sh in sl.shapes
                 if sh.has_text_frame and sh.text_frame.text.strip()]
        has_pic = any(getattr(sh, "shape_type", None) == MSO_SHAPE_TYPE.PICTURE for sh in sl.shapes)
        if not texts and not has_pic:
            issues.append(f"  スライド{i}: テキストも画像もなし（空白の可能性）")
        elif len(texts) == 1 and len(texts[0]) < 5 and not has_pic:
            issues.append(f"  スライド{i}: タイトルのみ（コンテンツ欠落の可能性）")
        nums = [sh for sh in sl.shapes if sh.has_text_frame
                and sh.left is not None and sh.top is not None
                and sh.left > Inches(11) and sh.top > Inches(6.5)
                and sh.text_frame.text.strip().isdigit()]
        if not nums and i > 1 and texts:
            issues.append(f"  スライド{i}: ページ番号が見当たらない（章扉・表紙なら無視可）")
    if issues:
        print(f"[audit_deck] {len(issues)}件の注意:")
        for w in issues:
            print(w)
    else:
        print("[audit_deck] OK — 全スライド診断クリア")
    return issues


# ===== スライドタイプ =====
def _track(run, spc):
    """run にレタースペーシング（字間）を設定する。spc は 1/100pt 単位。"""
    rPr = run._r.get_or_add_rPr()
    rPr.set("spc", str(int(spc)))


def add_title_slide(prs, title, subtitle, date, blank,
                    kicker="TECHNOLOGY INTELLIGENCE REPORT"):
    slide = prs.slides.add_slide(blank)
    bg = slide.background.fill
    bg.solid()
    bg.fore_color.rgb = DARK_SECTION

    # 巨大ゴースト・ワードマーク（背面・低コントラストで奥行きを作る）
    ghost = slide.shapes.add_textbox(Inches(0.7), Inches(4.55), Inches(13.0), Inches(2.6))
    gtf = ghost.text_frame
    gtf.word_wrap = False
    gtf.auto_size = MSO_AUTO_SIZE.NONE
    gp = gtf.paragraphs[0]
    grun = gp.add_run()
    grun.text = "APOLLO"
    grun.font.size = Pt(150)
    grun.font.bold = True
    grun.font.color.rgb = RGBColor(0x1B, 0x1B, 0x1F)  # 黒よりわずかに明るい墨
    _apply_font(grun, heading=True)
    _track(grun, 400)

    # 左端フル丈クリムゾン・ストリップ（強いアンカー）
    strip = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(0.30), Inches(7.5))
    strip.fill.solid()
    strip.fill.fore_color.rgb = ACCENT
    strip.line.fill.background()

    # キッカー（赤・字間広め・全大文字）
    kick = slide.shapes.add_textbox(Inches(1.15), Inches(0.95), Inches(11), Inches(0.4))
    kick.text_frame.word_wrap = True
    kick.text_frame.auto_size = MSO_AUTO_SIZE.NONE
    krun = kick.text_frame.paragraphs[0].add_run()
    krun.text = kicker
    krun.font.size = Pt(12)
    krun.font.bold = True
    krun.font.color.rgb = RED_ON_DARK
    _apply_font(krun, heading=True)
    _track(krun, 280)

    # APOLLO ワードマーク（白・字間広め）
    wm = slide.shapes.add_textbox(Inches(1.15), Inches(1.42), Inches(8), Inches(0.5))
    wm.text_frame.word_wrap = True
    wm.text_frame.auto_size = MSO_AUTO_SIZE.NONE
    wrun = wm.text_frame.paragraphs[0].add_run()
    wrun.text = "APOLLO"
    wrun.font.size = Pt(20)
    wrun.font.bold = True
    wrun.font.color.rgb = WHITE
    _apply_font(wrun, heading=True)
    _track(wrun, 600)

    tlen = len(title)
    if tlen <= 16:
        t_size, t_y = Pt(48), 2.95
    elif tlen <= 28:
        t_size, t_y = Pt(40), 2.85
    else:
        t_size, t_y = Pt(34), 2.75

    # 太いクリムゾン罫（タイトル直上）
    line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(1.18), Inches(t_y - 0.40), Inches(3.4), Inches(0.11))
    line.fill.solid()
    line.fill.fore_color.rgb = ACCENT
    line.line.fill.background()

    # タイトル（大判 白 明朝）
    txBox = slide.shapes.add_textbox(Inches(1.15), Inches(t_y), Inches(11.4), Inches(2.4))
    tf = txBox.text_frame
    tf.word_wrap = True
    tf.auto_size = MSO_AUTO_SIZE.NONE
    set_text(tf.paragraphs[0], title, t_size, WHITE, bold=True, line_spacing=1.16, heading=True)

    # サブタイトル（淡グレー）
    txBox2 = slide.shapes.add_textbox(Inches(1.18), Inches(5.55), Inches(11), Inches(0.6))
    txBox2.text_frame.word_wrap = True
    txBox2.text_frame.auto_size = MSO_AUTO_SIZE.NONE
    set_text(txBox2.text_frame.paragraphs[0], subtitle, Pt(15), RGBColor(0xC2, 0xC2, 0xC6), heading=True)

    # 下部メタデータ帯（クリムゾンの細帯に白文字を反転で乗せる）
    band = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(6.78), Inches(13.33), Inches(0.46))
    band.fill.solid()
    band.fill.fore_color.rgb = ACCENT
    band.line.fill.background()
    meta = slide.shapes.add_textbox(Inches(1.15), Inches(6.85), Inches(11.0), Inches(0.34))
    meta.text_frame.word_wrap = True
    meta.text_frame.auto_size = MSO_AUTO_SIZE.NONE
    mrun = meta.text_frame.paragraphs[0].add_run()
    mrun.text = date
    mrun.font.size = Pt(11)
    mrun.font.bold = True
    mrun.font.color.rgb = WHITE
    _apply_font(mrun, heading=True)
    _track(mrun, 120)
    return slide


def add_section_slide(prs, section_num, title, blank, subtitle=None):
    slide = prs.slides.add_slide(blank)
    bg = slide.background.fill
    bg.solid()
    bg.fore_color.rgb = DARK_SECTION

    # 巨大ゴースト番号（背面・右側にフル表示。クリップせず奥行きの主役にする）
    ghost = slide.shapes.add_textbox(Inches(5.2), Inches(1.05), Inches(7.6), Inches(5.4))
    tf_g = ghost.text_frame
    tf_g.word_wrap = False
    tf_g.auto_size = MSO_AUTO_SIZE.NONE
    p_g = tf_g.paragraphs[0]
    p_g.alignment = PP_ALIGN.RIGHT
    run_g = p_g.add_run()
    run_g.text = f"{section_num:02d}"
    run_g.font.size = Pt(300)
    run_g.font.color.rgb = GHOST_ON_DARK
    run_g.font.bold = True
    _apply_font(run_g, heading=True)

    # 左端フル丈クリムゾン・ストリップ（表紙と統一）
    strip = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(0.30), Inches(7.5))
    strip.fill.solid()
    strip.fill.fore_color.rgb = ACCENT
    strip.line.fill.background()

    # SECTION ラベル（赤・字間広め）
    eb = slide.shapes.add_textbox(Inches(1.15), Inches(2.75), Inches(6), Inches(0.4))
    eb.text_frame.word_wrap = True
    eb.text_frame.auto_size = MSO_AUTO_SIZE.NONE
    erun = eb.text_frame.paragraphs[0].add_run()
    erun.text = f"SECTION {section_num:02d}"
    erun.font.size = Pt(15)
    erun.font.bold = True
    erun.font.color.rgb = RED_ON_DARK
    _apply_font(erun, heading=True)
    _track(erun, 260)

    # 太いクリムゾン罫
    line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(1.18), Inches(3.32), Inches(2.8), Inches(0.10))
    line.fill.solid()
    line.fill.fore_color.rgb = ACCENT
    line.line.fill.background()

    # セクションタイトル（大判 白 明朝・番号の手前に重ねる）
    txBox = slide.shapes.add_textbox(Inches(1.15), Inches(3.55), Inches(8.0), Inches(1.7))
    tf = txBox.text_frame
    tf.word_wrap = True
    tf.auto_size = MSO_AUTO_SIZE.NONE
    set_text(tf.paragraphs[0], title, Pt(40), WHITE, bold=True, line_spacing=1.15, heading=True)

    if subtitle:
        txBox2 = slide.shapes.add_textbox(Inches(1.18), Inches(5.35), Inches(7.6), Inches(0.8))
        txBox2.text_frame.word_wrap = True
        txBox2.text_frame.auto_size = MSO_AUTO_SIZE.NONE
        set_text(txBox2.text_frame.paragraphs[0], subtitle, Pt(15), RGBColor(0xC2, 0xC2, 0xC6), heading=True)

    bot = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(7.1), Inches(13.33), Emu(27432))
    bot.fill.solid()
    bot.fill.fore_color.rgb = ACCENT
    bot.line.fill.background()
    return slide


def add_chapter_intro_slide(prs, eyebrow, title, bullets, blank, source=None, page_num=None):
    slide = prs.slides.add_slide(blank)
    col_l_x = MARGIN_L
    col_l_w = 4.6
    col_r_x = MARGIN_L + col_l_w + 0.6
    col_r_w = 13.33 - col_r_x - 0.7
    n = max(1, len(bullets))
    block_h = min(4.6, n * 0.52 + 0.35)
    top_y = max(1.3, (7.0 - block_h) / 2.0)
    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(col_l_x), Inches(top_y + 0.05),
                                 Inches(0.05), Inches(block_h - 0.1))
    bar.fill.solid()
    bar.fill.fore_color.rgb = ACCENT
    bar.line.fill.background()
    eb = slide.shapes.add_textbox(Inches(col_l_x + 0.22), Inches(top_y), Inches(col_l_w), Inches(0.35))
    eb.text_frame.word_wrap = True
    eb.text_frame.auto_size = MSO_AUTO_SIZE.NONE
    set_text(eb.text_frame.paragraphs[0], eyebrow, Pt(11), ACCENT, bold=True, heading=True)
    hd = slide.shapes.add_textbox(Inches(col_l_x + 0.22), Inches(top_y + 0.42),
                                  Inches(col_l_w), Inches(block_h - 0.42))
    htf = hd.text_frame
    htf.word_wrap = True
    htf.auto_size = MSO_AUTO_SIZE.NONE
    set_text(htf.paragraphs[0], title, Pt(24), INK, bold=True, line_spacing=1.22, heading=True)
    bd = slide.shapes.add_textbox(Inches(col_r_x), Inches(top_y), Inches(col_r_w), Inches(block_h))
    btf = bd.text_frame
    btf.word_wrap = True
    btf.auto_size = MSO_AUTO_SIZE.NONE
    for j, line in enumerate(bullets):
        p = btf.paragraphs[0] if j == 0 else btf.add_paragraph()
        p.space_after = Pt(8)
        mk = p.add_run()
        mk.text = "▪ "
        mk.font.size = Pt(14)
        mk.font.color.rgb = ACCENT
        mk.font.bold = True
        _apply_font(mk)
        add_rich_runs(p, line, base_size=Pt(14), base_color=DARK_GRAY,
                      bold_color=INK, line_spacing=1.4)
    if source:
        add_source_label(slide, source)
    add_bottom_bar_and_footer(slide, page_num)
    return slide


def add_kpi_slide(prs, title, sub_message, kpis, blank, source=None, page_num=None):
    slide = prs.slides.add_slide(blank)
    sub_y = add_title_shape(slide, title)
    content_y = add_sub_message(slide, sub_message, y=sub_y)
    n = len(kpis)
    available_w = 11.5
    start_x = 0.9
    gap = 0.2
    if n <= 4:
        n_cols, n_rows = n, 1
    else:
        n_cols, n_rows = min(4, (n + 1) // 2), 2
    card_w = (available_w - gap * (n_cols - 1)) / n_cols
    available_h = 6.5 - content_y
    row_gap = 0.2
    card_h = min((available_h - row_gap * (n_rows - 1)) / n_rows, 2.8)
    for idx, kpi in enumerate(kpis):
        row = idx // n_cols
        col = idx % n_cols
        x = start_x + col * (card_w + gap)
        y = content_y + row * (card_h + row_gap)
        card = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(x), Inches(y), Inches(card_w), Inches(card_h))
        card.fill.solid()
        card.fill.fore_color.rgb = LIGHT_GRAY
        card.line.color.rgb = BORDER_GRAY
        card.line.width = Emu(9525)
        bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(x), Inches(y), Emu(36576), Inches(card_h))
        bar.fill.solid()
        bar.fill.fore_color.rgb = ACCENT
        bar.line.fill.background()
        txL = slide.shapes.add_textbox(Inches(x + 0.15), Inches(y + 0.12), Inches(card_w - 0.3), Inches(0.25))
        set_text(txL.text_frame.paragraphs[0], kpi["label"], Pt(10), MEDIUM_GRAY, bold=True)
        txV = slide.shapes.add_textbox(Inches(x + 0.15), Inches(y + 0.4), Inches(card_w - 0.3), Inches(0.7))
        p = txV.text_frame.paragraphs[0]
        run = p.add_run()
        run.text = kpi["value"]
        run.font.size = Pt(32)
        run.font.bold = True
        run.font.color.rgb = NAVY
        _apply_font(run)
        if kpi.get("trend"):
            run2 = p.add_run()
            run2.text = f" {kpi['trend']}"
            run2.font.size = Pt(14)
            run2.font.color.rgb = MEDIUM_GRAY
            _apply_font(run2)
        txU = slide.shapes.add_textbox(Inches(x + 0.15), Inches(y + card_h - 0.4), Inches(card_w - 0.3), Inches(0.25))
        set_text(txU.text_frame.paragraphs[0], kpi.get("unit", ""), Pt(9), MEDIUM_GRAY)
    if source:
        add_source_label(slide, source)
    add_bottom_bar_and_footer(slide, page_num)
    return slide


def add_insight_slide(prs, title, sub_message, layers, blank, label="考察", source=None, page_num=None):
    slide = prs.slides.add_slide(blank)
    sub_y = add_title_shape(slide, title, label=label)
    content_y = add_sub_message(slide, sub_message, y=sub_y) if sub_message else sub_y + 0.1
    n = max(1, len(layers))
    avail_h = 6.7 - content_y
    gap = 0.18
    block_h = (avail_h - gap * (n - 1)) / n
    y = content_y
    for lyr in layers:
        bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(MARGIN_L), Inches(y), Emu(36576), Inches(block_h))
        bar.fill.solid()
        bar.fill.fore_color.rgb = ACCENT
        bar.line.fill.background()
        lb = slide.shapes.add_textbox(Inches(MARGIN_L + 0.18), Inches(y), Inches(2.0), Inches(0.3))
        set_text(lb.text_frame.paragraphs[0], lyr.get("label", ""), Pt(11), ACCENT, bold=True)
        body = slide.shapes.add_textbox(Inches(MARGIN_L + 0.18), Inches(y + 0.30),
                                        Inches(CONTENT_W - 0.18), Inches(block_h - 0.32))
        tf = body.text_frame
        tf.word_wrap = True
        tf.auto_size = MSO_AUTO_SIZE.NONE
        add_rich_runs(tf.paragraphs[0], lyr.get("body", ""), base_size=Pt(13),
                      base_color=DARK_GRAY, bold_color=INK, line_spacing=1.4)
        y += block_h + gap
    if source:
        add_source_label(slide, source)
    add_bottom_bar_and_footer(slide, page_num)
    return slide


def add_cards_slide(prs, title, sub_message, cards, blank, source=None, page_num=None):
    if len(cards) > 4:
        print(f"[WARN] add_cards_slide: cards={len(cards)} > 4")
    slide = prs.slides.add_slide(blank)
    sub_y = add_title_shape(slide, title)
    content_y = add_sub_message(slide, sub_message, y=sub_y)
    n = len(cards)
    gap = 0.25
    total_w = 12.3
    card_w = (total_w - gap * (n - 1)) / n
    card_h = 6.5 - content_y
    header_h = 0.45
    for i, card in enumerate(cards):
        x = 0.5 + i * (card_w + gap)
        color = card.get("color", ACCENT)
        hdr = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(content_y),
                                     Inches(card_w), Inches(header_h))
        hdr.fill.solid()
        hdr.fill.fore_color.rgb = color
        hdr.line.fill.background()
        txH = slide.shapes.add_textbox(Inches(x + 0.1), Inches(content_y + 0.05),
                                       Inches(card_w - 0.2), Inches(header_h - 0.1))
        set_text(txH.text_frame.paragraphs[0], card["header"], Pt(14), WHITE, bold=True)
        txH.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
        body_y = content_y + header_h
        body_h = card_h - header_h
        bdy = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(x), Inches(body_y),
                                     Inches(card_w), Inches(body_h))
        bdy.fill.solid()
        bdy.fill.fore_color.rgb = LIGHT_GRAY
        bdy.line.color.rgb = BORDER_GRAY
        bdy.line.width = Emu(9525)
        txB = slide.shapes.add_textbox(Inches(x + 0.12), Inches(body_y + 0.1),
                                       Inches(card_w - 0.24), Inches(body_h - 0.2))
        tf = txB.text_frame
        tf.word_wrap = True
        tf.auto_size = MSO_AUTO_SIZE.NONE
        body_text = card.get("body", "")
        if isinstance(body_text, list):
            for j, item in enumerate(body_text):
                p = tf.paragraphs[0] if j == 0 else tf.add_paragraph()
                add_rich_runs(p, f"・{item}", base_size=Pt(12), base_color=DARK_GRAY,
                              bold_color=NAVY, line_spacing=1.4)
        else:
            add_rich_runs(tf.paragraphs[0], body_text, base_size=Pt(13),
                          base_color=DARK_GRAY, bold_color=NAVY, line_spacing=1.4)
    if source:
        add_source_label(slide, source)
    add_bottom_bar_and_footer(slide, page_num)
    return slide


def add_narrative_slide(prs, title, sub_message, paragraphs, blank, source=None, page_num=None):
    slide = prs.slides.add_slide(blank)
    sub_y = add_title_shape(slide, title)
    content_y = add_sub_message(slide, sub_message, y=sub_y) if sub_message else sub_y + 0.1
    avail_h = 6.5 - content_y
    n_para = max(1, len(paragraphs))
    est_h = min(avail_h, n_para * 0.55 + 0.2)
    start_y = content_y + max(0.0, (avail_h - est_h) / 2.0)
    box = slide.shapes.add_textbox(Inches(0.5 + 0.12), Inches(start_y + 0.08),
                                   Inches(12.3 - 0.24), Inches(est_h - 0.16))
    tf = box.text_frame
    tf.word_wrap = True
    tf.auto_size = MSO_AUTO_SIZE.NONE
    for i, item in enumerate(paragraphs):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.space_after = Pt(8)
        marker = p.add_run()
        marker.text = "■ "
        marker.font.size = Pt(16)
        marker.font.color.rgb = NAVY
        _apply_font(marker)
        add_rich_runs(p, item, base_size=Pt(16), base_color=DARK_GRAY, bold_color=NAVY, line_spacing=1.5)
        # re-add marker lost by clear(): prepend
    if source:
        add_source_label(slide, source)
    add_bottom_bar_and_footer(slide, page_num)
    return slide


def add_closing_slide(prs, report_title, blank):
    slide = prs.slides.add_slide(blank)
    bg = slide.background.fill
    bg.solid()
    bg.fore_color.rgb = DARK_SECTION
    txBox2 = slide.shapes.add_textbox(Inches(1.2), Inches(3.0), Inches(10.9), Inches(1.6))
    tf2 = txBox2.text_frame
    tf2.word_wrap = True
    set_text(tf2.paragraphs[0], report_title, Pt(32), WHITE, bold=True, line_spacing=1.2, heading=True)
    tf2.paragraphs[0].alignment = PP_ALIGN.CENTER
    bot = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(7.1), Inches(13.33), Emu(27432))
    bot.fill.solid()
    bot.fill.fore_color.rgb = ACCENT
    bot.line.fill.background()
    return slide


# ===== デッキ構築 =====
def build():
    prs = Presentation(TEMPLATE)
    blank = prs.slide_layouts[6]

    add_title_slide(
        prs,
        "全固体電池 特許ランドスケープ",
        "Mission Objective: 主要プレイヤーの技術ポジションと参入余地を特定する",
        "2026-05-27  APOLLO v8.0.0",
        blank,
    )

    add_kpi_slide(
        prs,
        "母集団3,482件・上位5社で全体の46%を占有 ～硫化物系に出願が集中",
        "**硫化物系固体電解質**への集中が顕著。新規参入は酸化物系・界面設計に余地。",
        [
            {"label": "総特許件数", "value": "3,482", "unit": "件", "trend": "↑ CAGR 18%"},
            {"label": "出願人数", "value": "214", "unit": "社", "trend": ""},
            {"label": "上位5社シェア", "value": "46", "unit": "%", "trend": "↑ 集中化"},
            {"label": "HHI", "value": "0.082", "unit": "中程度集中", "trend": ""},
        ],
        blank,
        source="特許データベース XXX（2026-05-27）",
        page_num=2,
    )

    add_section_slide(prs, 3, "クラスタ動態分析", blank, subtitle="技術領域の成長・成熟をマッピングする")

    add_chapter_intro_slide(
        prs,
        "SECTION 03 / クラスタ動態",
        "この章で見ること",
        [
            "**何を**: 8つの技術クラスタを出願件数とCAGRの2軸で4象限に配置する",
            "**なぜ**: 成長領域と成熟領域を区別し、投資の優先順位を判断するため",
            "**着眼**: 「新興」象限に位置する界面設計クラスタの萌芽性を検証する",
            "**裏付け**: クラスタ動態の判定をWeb調査による市場動向で補強する",
        ],
        blank,
        page_num=4,
    )

    add_cards_slide(
        prs,
        "4象限で見ると界面設計が「新興」に位置 ～件数は小さいが成長率は最大",
        "**成長リーダー**（硫化物系）と**新興**（界面設計）の二極で技術投資が進む。",
        [
            {"header": "成長リーダー", "body": ["硫化物系電解質（842件）", "CAGR 24%・大手集中", "コア特許の囲い込み進行"], "color": ACCENT},
            {"header": "新興", "body": ["界面設計（118件）", "CAGR 41%・最速成長", "参入余地が最も大きい"], "color": INK},
            {"header": "成熟", "body": ["酸化物系（596件）", "CAGR 6%・横ばい", "改良発明が中心"], "color": MEDIUM_GRAY},
            {"header": "ニッチ", "body": ["ポリマー系（203件）", "CAGR 9%・限定領域", "特定用途に特化"], "color": AMBER},
        ],
        blank,
        source="特許データベース XXX（2026-05-27）",
        page_num=5,
    )

    add_insight_slide(
        prs,
        "界面設計クラスタは「小さな急成長」 ～3年以内に主戦場化する可能性",
        "件数の小ささを成長率が補い、参入の好機が開いている構造。",
        [
            {"label": "事実", "body": "界面設計クラスタは118件と全体の3.4%に過ぎないが、直近3年のCAGRは41%と全クラスタ最大。出願人数も12社から38社へ急増している。"},
            {"label": "解釈", "body": "硫化物系・酸化物系の電解質材料そのものの開発が一巡し、競争の焦点が「電極と電解質の界面抵抗の低減」という応用フェーズへ移行しつつあることを示す。"},
            {"label": "洞察", "body": "材料系クラスタが大手による特許の囲い込みで参入困難化する一方、界面設計はまだ支配的プレイヤーが不在。技術的難所であるがゆえに、ここを制した者が次の標準を握る構造的好機がある。"},
            {"label": "示唆", "body": "後発企業は飽和した材料系での正面競争を避け、界面設計に経営資源を集中すべき。ただし大手の参入も時間の問題であり、2-3年以内のコア特許出願が分水嶺となる。"},
        ],
        blank,
        source="特許データベース XXX（2026-05-27）",
        page_num=6,
    )

    add_closing_slide(prs, "界面設計こそ次の主戦場 — 早期のコア特許出願が勝敗を分ける", blank)

    disable_all_shadows(prs)
    audit_deck(prs)

    out_dir = os.path.join(HERE, "..", "reports")
    os.makedirs(out_dir, exist_ok=True)
    out = os.path.join(out_dir, "apollo_demo_deck.pptx")
    prs.save(out)
    print("saved:", os.path.abspath(out))


if __name__ == "__main__":
    build()
