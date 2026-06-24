from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_AUTO_SIZE, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE, MSO_SHAPE_TYPE, MSO_CONNECTOR
from pptx.chart.data import CategoryChartData
from pptx.enum.chart import XL_CHART_TYPE, XL_LEGEND_POSITION
from PIL import Image, ImageDraw
from lxml import etree
import os, re, sys, math, tempfile


# --- ベース（墨と紙） ---
INK = RGBColor(0x11, 0x11, 0x11)        # 主要テキスト（旧NAVYの役割）
PAPER = RGBColor(0xFF, 0xFF, 0xFF)      # 紙（明スライド背景）
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
DARK_SECTION = RGBColor(0x10, 0x10, 0x11) # 章扉・ダークセクション背景（黒）
BLACK_BAR = RGBColor(0x0B, 0x0B, 0x0C)  # トップバー等の黒帯

# --- グレー階調 ---
DARK_GRAY = RGBColor(0x11, 0x11, 0x11)  # 本文テキスト（≒INK）
MEDIUM_GRAY = RGBColor(0x68, 0x68, 0x68) # 補足・キャプション・ページ番号
LIGHT_GRAY = RGBColor(0xF1, 0xF2, 0xF3) # カード/ゼブラ既定面
PALE_GRAY = RGBColor(0xF6, 0xF7, 0xF8)  # グレーセクション面
SURFACE_GRAY = RGBColor(0xEC, 0xEE, 0xEF) # やや濃い面
BORDER_GRAY = RGBColor(0xD8, 0xDA, 0xDD) # 罫線・区切り線

# --- クリムゾン・アクセント（単一アクセント） ---
ACCENT = RGBColor(0xC5, 0x12, 0x12)     # アクセントバー・ラベル・強調（=RED）
BRIGHT_RED = RGBColor(0xF4, 0x33, 0x33) # 明るい赤（クリムゾン→明赤グラデの片側）
CRIMSON_DEEP = RGBColor(0x8E, 0x00, 0x14) # 深紅（深→クリムゾングラデの片側）
RED = RGBColor(0xC5, 0x12, 0x12)
BRIGHT_RED = RGBColor(0xE3, 0x33, 0x33) # ハイライト/ホバー相当の明るい赤
DEEP_RED = RGBColor(0x83, 0x10, 0x10)   # チャート最濃バー（"now"）
RED_ON_DARK = RGBColor(0xFF, 0x70, 0x70) # 黒背景上のバッジ文字

# --- 指標色（低彩度の範囲で。緑/橙は使わずグレー+赤で表現） ---
RED_ACCENT = RGBColor(0xC5, 0x12, 0x12)  # マイナス/警告/ネガティブ
POS_INK = RGBColor(0x11, 0x11, 0x11)     # ポジティブは墨で強調（彩色しない）
KEY_MSG_BG = RGBColor(0xF6, 0xF7, 0xF8)  # ■サブメッセージ背景（淡グレー面）

# --- 章扉ゴースト番号（黒背景上）---
GHOST_ON_DARK = RGBColor(0x4A, 0x4A, 0x52) # 黒地で視認できる濃グレー（巨大番号）

# --- チャート配色（グレー段階 + 赤で「今/注目」を示す） ---
BAR_OLD = RGBColor(0xC7, 0xCB, 0xD0)    # 過去
BAR_MID = RGBColor(0x8F, 0x8F, 0x8F)    # 中間
BAR_HOT = RGBColor(0xC5, 0x12, 0x12)    # 注目（赤）
BAR_NOW = RGBColor(0x83, 0x10, 0x10)    # 直近（濃赤）
GRID_LINE = RGBColor(0xE6, 0xE8, 0xEA)  # 目盛線（点線 3 4 推奨）

# レイアウト・グリッド定数（原則1: 余白と整列）
MARGIN_L = 0.9                       # Inches — 全要素を揃える固定左マージン
CONTENT_W = 13.33 - MARGIN_L * 2     # Inches — コンテンツ幅（左右マージンを揃え中央寄せ＝右寄り解消）

# ボトム定数（フッター帯は廃止。右下ページ番号のみ）
PAGE_NUM_Y = 7.02      # Inches — 右下ページ番号のベースライン（スライド内に収める）
PAGE_NUM_X = 12.45     # Inches — 右下ページ番号の左端（右寄せ）
PAGE_NUM_W = 0.55      # Inches — ページ番号ボックス幅（右端 13.0in でスライド内に収める）
# トップバー（明スライド上端の黒帯＋赤下線）
TOPBAR_H = 0.45        # Inches ≒ 44px 相当
TOPBAR_RED_H = 0.04    # Inches ≒ 3px 相当（黒帯下端のクリムゾン線）

# --- 旧名エイリアス（後方互換）---
# 既存15種スライドは旧名(NAVY/BLUE/...)を参照する。配色をモノトーン＋クリムゾン単一
# アクセントへ一括移行するため、旧名を新パレットへ写像する。青/緑/橙は使わず、
# カテゴリ区別はグレー段階＋赤で表現する（設計トーン: 低彩度・高コントラスト）。
NAVY = INK                              # タイトル/強調 → 墨
BLUE = MEDIUM_GRAY                      # 旧・第2系列色 → 中間グレー
GHOST_NAVY = GHOST_ON_DARK              # 章扉ゴースト番号
GREEN_ACCENT = INK                      # 旧・ポジティブ → 墨で強調（彩色しない）
AMBER = BAR_MID                         # 旧・注意 → 中間グレー


from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_AUTO_SIZE, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE, MSO_SHAPE_TYPE, MSO_CONNECTOR
from PIL import Image, ImageDraw
from lxml import etree
import os, re, sys

# --- フォント設定（欧文=Century Gothic / 和文=Yu Gothic） ---
# 環境にCentury Gothicが無い場合に備え、利用側でフォールバックを考慮すること
JA_FONT = "Yu Gothic"
LATIN_FONT = "Century Gothic"
# 見出し・章扉・タイトル用の明朝（和文＝Yu Mincho、欧文も同フォントで統一）
HEADING_JA_FONT = "Yu Mincho"
HEADING_LATIN_FONT = "Yu Mincho"

A_NS = 'http://schemas.openxmlformats.org/drawingml/2006/main'

# テンプレートPPTX

def _apply_font(run, heading=False):
    """runにデュアルフォント（欧文 + 日本語）+ 言語を設定する。
    heading=True で章扉・タイトル・見出し用の明朝（Yu Mincho）にする。"""
    latin = HEADING_LATIN_FONT if heading else LATIN_FONT
    ja = HEADING_JA_FONT if heading else JA_FONT
    run.font.name = latin
    rPr = run._r.get_or_add_rPr()
    rPr.set('lang', 'ja-JP')
    rPr.set('altLang', 'en-US')
    ea = rPr.find(f'{{{A_NS}}}ea')
    if ea is None:
        ea = etree.SubElement(rPr, f'{{{A_NS}}}ea')
    ea.set('typeface', ja)


def _apply_kinsoku(paragraph):
    """段落に日本語禁則処理を設定する"""
    pPr = paragraph._p.get_or_add_pPr()
    pPr.set('eaLnBrk', '1')
    pPr.set('hangingPunct', '1')


def _track(run, spc):
    """run にレタースペーシング（字間）を設定する。spc は 1/100pt 単位。
    キッカー・ワードマーク・SECTIONラベル等を「大きく開いた字間」で締める用途。
    （原則7: 階層差はサイズより太さ・字間・色でつける）"""
    rPr = run._r.get_or_add_rPr()
    rPr.set('spc', str(int(spc)))


def add_rich_runs(paragraph, text, base_size=Pt(14), base_color=DARK_GRAY,
                  bold_color=None, force_bold=False, line_spacing=1.4, heading=False):
    """**太字**マーカー解析 + デュアルフォント + 禁則 + 行間。
    heading=True で見出し明朝（Yu Mincho）を適用する。"""
    paragraph.clear()
    bold_color = bold_color or base_color
    _apply_kinsoku(paragraph)
    paragraph.line_spacing = line_spacing

    parts = re.split(r'(\*\*.*?\*\*)', text)
    for part in parts:
        if not part:
            continue
        if part.startswith('**') and part.endswith('**'):
            run = paragraph.add_run()
            run.text = part[2:-2]
            run.font.size = base_size
            run.font.bold = True
            run.font.color.rgb = bold_color
        else:
            run = paragraph.add_run()
            run.text = part
            run.font.size = base_size
            run.font.bold = force_bold
            run.font.color.rgb = bold_color if force_bold else base_color
        _apply_font(run, heading=heading)


def set_text(p, text, size, color, bold=False, line_spacing=None, heading=False):
    """シンプルテキスト設定（デュアルフォント + 禁則付き）。
    heading=True で見出し明朝（Yu Mincho）を適用する。"""
    p.text = ""
    run = p.add_run()
    run.text = str(text).replace("**", "")   # MD痕跡の保険: 単一スタイル描画では ** を必ず除去
    run.font.size = size
    run.font.color.rgb = color
    run.font.bold = bold
    _apply_font(run, heading=heading)
    _apply_kinsoku(p)
    if line_spacing:
        p.line_spacing = line_spacing


def add_title_shape(slide, text, x=MARGIN_L, y=0.55, w=CONTENT_W, label=None, color=None):
    """スライドタイトル（20-30pt Bold INK + 任意の赤ラベル）

    タイトル＝結論。「～」で副題を付け、**32文字以内の1行**でストーリーを完結。
    全幅下線・左肩の短い赤罫はいずれも廃止（視覚ノイズ削減）。
    label を渡すと、タイトルの上に小さな赤ラベル（EYEBROW）を添える
    （例: "ATLAS / 出願トレンド"、"EXECUTIVE SUMMARY"）。原則3・原則7。
    Returns:
        float: タイトル下端のy座標（サブメッセージの配置基準）
    """
    text_len = len(text)
    # タイトルは32文字以内（1行）を厳守。超過は折返しの恐れがあるため警告する。
    if text_len > 32:
        print(f"[WARN] add_title_shape: タイトル{text_len}文字 > 32。2行折返しの恐れ。32文字以内に短縮推奨: {text[:40]}…")
    if text_len <= 16:
        font_size = Pt(30)
        box_h = 0.62
    elif text_len <= 24:
        font_size = Pt(26)
        box_h = 0.66
    elif text_len <= 32:
        font_size = Pt(22)
        box_h = 0.70
    else:
        font_size = Pt(20)
        box_h = 0.92

    # （改良1）タイトル上の短いクリムゾン EYEBROW 罫は廃止（視覚ノイズのため削除）。
    cy = y + 0.04
    # 任意の赤ラベル（モジュール名・分類）。字間を開け全大文字推奨
    if label:
        lab = slide.shapes.add_textbox(Inches(x), Inches(cy), Inches(w), Inches(0.28))
        set_text(lab.text_frame.paragraphs[0], label, Pt(10), ACCENT, bold=True, heading=True)
        cy += 0.30

    # タイトル本文（墨・字間詰め）
    txBox = slide.shapes.add_textbox(Inches(x), Inches(cy), Inches(w), Inches(box_h))
    tf = txBox.text_frame
    tf.word_wrap = True
    tf.auto_size = MSO_AUTO_SIZE.NONE
    p = tf.paragraphs[0]
    # タイトルは墨・明朝（Yu Mincho）。`～` 以降の副題のみクリムゾンで締める運用も可
    _tc = color if color is not None else INK
    add_rich_runs(p, text, base_size=font_size, base_color=_tc,
                  bold_color=_tc, force_bold=True, line_spacing=1.18, heading=True)

    # （改良3）サブメッセージ（タイトル直下のグレー帯）を少し上へ。返り値の余白を詰める。
    return cy + box_h + 0.00  # サブメッセージの開始y座標を返す


def add_sub_message(slide, message, x=MARGIN_L, y=None, w=CONTENT_W):
    """■マーカー付きサブメッセージ（ボックス囲み、タイトル直下）

    KEY_MSG_BG背景 + 左ACCENTバーのボックスで要点を強調。
    Args:
        y: 開始y座標。Noneの場合はadd_title_shapeの戻り値を使うこと。
    Returns:
        float: ボックス下端のy座標 + マージン（コンテンツ開始位置）
    """
    if y is None:
        y = 1.00
    chars_per_line = 38
    num_lines = max(1, -(-len(message) // chars_per_line))
    box_h = 0.22 + num_lines * 0.38

    # 背景ボックス（角丸・淡グレー面・枠線なし）
    box = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(box_h)
    )
    try: box.adjustments[0] = min(0.5, 0.16 / box_h)   # 角丸半径≈0.16inで一定に
    except Exception: pass
    box.fill.solid()
    box.fill.fore_color.rgb = KEY_MSG_BG
    box.line.fill.background()
    try: box.shadow.inherit = False
    except Exception: pass

    # テキスト（先頭の■記号なし・オブジェクト内で上下中央寄せ・1段大きめ）
    txBox = slide.shapes.add_textbox(
        Inches(x + 0.24), Inches(y), Inches(w - 0.44), Inches(box_h)
    )
    tf = txBox.text_frame
    tf.word_wrap = True
    tf.auto_size = MSO_AUTO_SIZE.NONE
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE          # 上下中央寄せ
    tf.margin_top = 0; tf.margin_bottom = 0
    p = tf.paragraphs[0]
    parts = re.split(r'(\*\*.*?\*\*)', message)
    for part in parts:
        if not part:
            continue
        run = p.add_run()
        if part.startswith('**') and part.endswith('**'):
            run.text = part[2:-2]
            run.font.bold = True
            run.font.color.rgb = INK  # 強調は墨（彩色しない）
        else:
            run.text = part
            run.font.color.rgb = DARK_GRAY
        run.font.size = Pt(16.5)                     # 1段階拡大（旧15pt）
        _apply_font(run)
    _apply_kinsoku(p)
    p.line_spacing = 1.4

    return y + box_h + 0.10


def grad_fill(shape, c1, c2, angle=90):
    """既存図形にクリムゾン→明赤などのグラデーション塗りを適用する（影は継承しない）。
    c1, c2 は RGBColor。アクセントバー・KPI強調・棒グラフの注目バー等に使う。"""
    shape.line.fill.background()
    try: shape.shadow.inherit = False
    except Exception: pass
    shape.fill.gradient()
    gs = shape.fill.gradient_stops
    gs[0].color.rgb = c1; gs[0].position = 0.0
    gs[1].color.rgb = c2; gs[1].position = 1.0
    try: shape.fill.gradient_angle = angle
    except Exception: pass
    return shape


def _grad_text_fill(textbox, c1, c2, ang=5400000):
    """テキスト（全ラン）にグラデーション塗りを適用する。章扉の大番号など。
    c1/c2 は6桁HEX（c1=上, c2=下）。ang は 60000分の1度（5400000=90°＝上→下）。"""
    A = A_NS
    for p in textbox.text_frame.paragraphs:
        for r in p.runs:
            rPr = r._r.get_or_add_rPr()
            sf = rPr.find(f'{{{A}}}solidFill')
            idx = list(rPr).index(sf) if sf is not None else 0
            if sf is not None:
                rPr.remove(sf)
            grad = etree.Element(f'{{{A}}}gradFill')
            gsLst = etree.SubElement(grad, f'{{{A}}}gsLst')
            for pos, col in ((0, c1), (100000, c2)):
                gs = etree.SubElement(gsLst, f'{{{A}}}gs'); gs.set('pos', str(pos))
                etree.SubElement(gs, f'{{{A}}}srgbClr').set('val', col)
            lin = etree.SubElement(grad, f'{{{A}}}lin')
            lin.set('ang', str(ang)); lin.set('scaled', '1')
            rPr.insert(idx, grad)
    return textbox


def grad_rect(slide, x, y, w, h, c1, c2, angle=90, radius=0.0):
    """グラデーション矩形（G1グラデの最小取り込み・影なし）。"""
    shp = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE if radius > 0 else MSO_SHAPE.RECTANGLE,
        Inches(x), Inches(y), Inches(w), Inches(h))
    if radius > 0:
        try: shp.adjustments[0] = radius
        except Exception: pass
    shp.line.fill.background()
    shp.fill.gradient()
    gs = shp.fill.gradient_stops
    gs[0].color.rgb = c1; gs[0].position = 0.0
    gs[1].color.rgb = c2; gs[1].position = 1.0
    try: shp.fill.gradient_angle = angle
    except Exception: pass
    return shp

# 章目印（現在章の進捗）用グローバル。add_section_slide が CURRENT_CHAPTER を更新する。
CURRENT_CHAPTER = 0
TOTAL_CHAPTERS = 0

def _set_fill_alpha(shp, alpha_pct):
    """ソリッド塗りシェイプに不透明度（0-100%）を設定する。色相は保ったまま半透明化する。"""
    sp = shp._element.spPr
    fill = sp.find(f'{{{A_NS}}}solidFill')
    if fill is None:
        return
    srgb = fill.find(f'{{{A_NS}}}srgbClr')
    if srgb is None:
        return
    for a in srgb.findall(f'{{{A_NS}}}alpha'):
        srgb.remove(a)
    al = etree.SubElement(srgb, f'{{{A_NS}}}alpha')
    al.set('val', str(int(alpha_pct * 1000)))   # 0-100% → 0-100000


def add_chapter_marker(slide):
    """章目印（パターン③・改）：左端に縦プログレスゲージ。バーのX中心をスライド左端(x=0)に重ねる。
    全章ぶんのセルに分割し（章ごとに切れ目）、1〜現在章のセルを **同色クリムゾンの半透明** で点灯する。
    現在章は最後に点灯したセルそのもの（突起・幅違いのドットは作らない＝バーと同じ幅）。テキストは書かない。"""
    if not (CURRENT_CHAPTER > 0 and TOTAL_CHAPTERS > 0):
        return
    bw = 0.20                 # バー幅（半分は画面外＝中心線が x=0）
    x = -bw / 2.0
    top, bot = 0.55, 7.0
    Hb = bot - top
    N = int(TOTAL_CHAPTERS)
    gap = 0.045               # 章ごとの切れ目
    seg = (Hb - gap * (N - 1)) / N
    for i in range(N):
        n = i + 1
        y = top + i * (seg + gap)
        cell = slide.shapes.add_shape(
            MSO_SHAPE.PARALLELOGRAM, Inches(x), Inches(y), Inches(bw), Inches(seg))
        try: cell.adjustments[0] = 0.45   # 平行四辺形の傾き
        except Exception: pass
        cell.line.fill.background()
        cell.fill.solid()
        if n <= CURRENT_CHAPTER:
            cell.fill.fore_color.rgb = ACCENT          # いまの色のまま
            _set_fill_alpha(cell, 42)                  # 半透明（約42%）
        else:
            cell.fill.fore_color.rgb = RGBColor(0xE6, 0xE8, 0xEC)  # 未読＝淡グレー


def add_bottom_bar_and_footer(slide, page_num=None):
    """全スライド共通: 右下ページ番号 ＋ 左端の章目印（パターン③）。
    タイトルスライド・セクションスライド・クロージングスライドでは呼ばない。
    """
    add_chapter_marker(slide)
    if page_num is not None and _PENDING_SECTION[0] is not None:
        _SECTION_FIRST_PAGE[_PENDING_SECTION[0]] = page_num
        _PENDING_SECTION[0] = None
    if page_num is None:
        return
    # 右下ページ番号のみ（10pt MEDIUM_GRAY・右寄せ）
    txBox = slide.shapes.add_textbox(
        Inches(PAGE_NUM_X), Inches(PAGE_NUM_Y), Inches(PAGE_NUM_W), Inches(0.25)
    )
    p = txBox.text_frame.paragraphs[0]
    set_text(p, str(page_num), Pt(10), MEDIUM_GRAY)
    p.alignment = PP_ALIGN.RIGHT


def disable_all_shadows(prs):
    """全スライド・全シェイプの継承シャドウを無効化する（オブジェクトに影を付けない方針）。
    **保存直前に必ず呼ぶこと**（prs.save の直前）。グループ内シェイプも再帰的に処理する。"""
    def _kill(shp):
        # 1) 継承シャドウ無効化
        try:
            shp.shadow.inherit = False
        except Exception:
            pass
        # 2) spPr に空の <a:effectLst/> を強制し、プリセット影を完全に消す
        try:
            spPr = shp._element.spPr
            if spPr is not None:
                for el in spPr.findall(f'{{{A_NS}}}effectLst'):
                    spPr.remove(el)
                etree.SubElement(spPr, f'{{{A_NS}}}effectLst')
        except Exception:
            pass
        # 3) スタイルの effectRef（テーマ由来の影）を idx=0 に落とす。
        #    LibreOffice は effectLst だけでは effectRef の影を消さないことがある。
        try:
            P_NS = 'http://schemas.openxmlformats.org/presentationml/2006/main'
            style = shp._element.find(f'{{{P_NS}}}style')
            if style is not None:
                eff = style.find(f'{{{A_NS}}}effectRef')
                if eff is not None:
                    eff.set('idx', '0')
        except Exception:
            pass

    def _walk(shapes):
        for shp in shapes:
            _kill(shp)
            if getattr(shp, "shape_type", None) == MSO_SHAPE_TYPE.GROUP:
                _walk(shp.shapes)
    for slide in prs.slides:
        _walk(slide.shapes)


def audit_deck(prs):
    """生成後の自己診断（ハードエラーにはせず print 警告のみ）。
    保存直前、disable_all_shadows の後に呼ぶ。レイアウト破綻・ページ番号欠損・
    タイトルのみ頁などを検出して、人手の全枚確認を補助する。"""
    issues = []
    for i, sl in enumerate(prs.slides, 1):
        texts = [sh.text_frame.text.strip() for sh in sl.shapes
                 if sh.has_text_frame and sh.text_frame.text.strip()]
        has_pic = any(getattr(sh, "shape_type", None) == MSO_SHAPE_TYPE.PICTURE
                      for sh in sl.shapes)
        if not texts and not has_pic:
            issues.append(f"  スライド{i}: テキストも画像もなし（空白の可能性）")
        elif len(texts) == 1 and len(texts[0]) < 5 and not has_pic:
            issues.append(f"  スライド{i}: タイトルのみ（コンテンツ欠落の可能性）")
        # ページ番号チェック: 右下に小さな数字テキストがあるか（表紙・章扉は除外）
        nums = [sh for sh in sl.shapes if sh.has_text_frame
                and sh.left is not None and sh.top is not None
                and sh.left > Inches(11) and sh.top > Inches(6.5)
                and sh.text_frame.text.strip().isdigit()]
        if not nums and i > 1 and texts:
            issues.append(f"  スライド{i}: ページ番号が見当たらない（章扉・表紙なら無視可）")
    if issues:
        print(f"[audit_deck] {len(issues)}件の注意:")
        for w in issues:
            print(w)
    else:
        print("[audit_deck] OK — 全スライド診断クリア")
    return issues


def fit_image(slide, image_path, max_x, max_y, max_w, max_h,
              halign="center", valign="middle"):
    """画像をアスペクト比保持で指定領域内に配置。

    halign/valign で寄せ方向を指定（既定=中央）。返り値の Picture に実際の描画矩形
    `_fit_rect=(left, top, use_w, use_h)` を付与するので、呼び出し側は実描画右端に
    注釈を密着させる等の精密配置ができる（中央配置の死に余白を解消するため）。
    """
    if not os.path.exists(image_path):
        return None
    img = Image.open(image_path)
    img_w, img_h = img.size
    ratio = img_h / img_w
    if max_w * ratio <= max_h:
        use_w, use_h = max_w, max_w * ratio
    else:
        use_h, use_w = max_h, max_h / ratio
    if halign == "left":
        left = max_x
    elif halign == "right":
        left = max_x + (max_w - use_w)
    else:
        left = max_x + (max_w - use_w) / 2
    if valign == "top":
        top = max_y
    elif valign == "bottom":
        top = max_y + (max_h - use_h)
    else:
        top = max_y + (max_h - use_h) / 2
    pic = slide.shapes.add_picture(
        image_path, Inches(left), Inches(top),
        width=Inches(use_w), height=Inches(use_h)
    )
    img.close()
    pic._fit_rect = (left, top, use_w, use_h)
    return pic


def _text_capacity(width_in, height_in, font_pt, line_spacing=1.4, fudge=1.05):
    """固定フォントで指定ボックスに収まる概算文字数（全角基準・安全側）を返す。

    方針: フォントは縮めない。ここで得た上限を使って content 側の文字数を制限する。
    返り値: (1行あたり字数 cpl, 行数 lines, 許容総文字数)。
    """
    char_w = font_pt / 72.0                      # 全角1文字 ≈ フォントpt
    line_h = font_pt * line_spacing / 72.0
    cpl = max(6, int(width_in / char_w))
    lines = max(1, int(height_in / line_h + 0.02))
    return cpl, lines, int(cpl * lines * fudge)


def _clip_rich(text, max_chars, where=""):
    """**bold** 記法を保ったまま可視文字数で切り詰める。超過時は警告を出す。

    フォントは縮小しない方針のため、箱に収まらない分はここで切り、content 側の
    要約を促す（[WARN] を出力）。bold ペアは閉じたまま維持する。
    """
    if max_chars <= 0 or not text:
        return text
    parts = re.split(r'(\*\*.*?\*\*)', text)
    if sum(len(p[2:-2] if (p.startswith('**') and p.endswith('**')) else p)
           for p in parts if p) <= max_chars:
        return text
    out, used = [], 0
    for part in parts:
        if not part:
            continue
        is_bold = part.startswith('**') and part.endswith('**')
        inner = part[2:-2] if is_bold else part
        if used + len(inner) <= max_chars:
            out.append(part); used += len(inner)
        else:
            remain = max_chars - used
            if remain > 0:
                cut = inner[:remain].rstrip() + "…"
                out.append(f"**{cut}**" if is_bold else cut)
            break
    print(f"[WARN] 文字数上限({max_chars}字)超過のため切り詰め: {where} — content側で要約を推奨")
    return "".join(out)


def add_source_label(slide, source_text, x=0.5, y=6.98, w=11.3):
    """（出所）ラベル（スライド下端ギリギリに配置）"""
    txBox = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(0.35))
    tf = txBox.text_frame
    tf.word_wrap = True
    tf.auto_size = MSO_AUTO_SIZE.NONE
    p = tf.paragraphs[0]
    set_text(p, f"（出所）{source_text}", Pt(9), MEDIUM_GRAY)


def _set_bullet(paragraph, char="•", color="111111", indent_in=0.20):  # （task2）中点は黒
    """（task3）段落をネイティブ箇条書きにする（■マーカー廃止）。
    buClr→buFont→buChar の順で pPr に付与し、ぶら下げインデントを設定。"""
    pPr = paragraph._p.get_or_add_pPr()
    emu = str(int(indent_in * 914400))
    pPr.set('marL', emu); pPr.set('indent', '-' + emu)
    buClr = etree.SubElement(pPr, f'{{{A_NS}}}buClr')
    srgb = etree.SubElement(buClr, f'{{{A_NS}}}srgbClr'); srgb.set('val', color)
    buFont = etree.SubElement(pPr, f'{{{A_NS}}}buFont'); buFont.set('typeface', 'Arial')
    buChar = etree.SubElement(pPr, f'{{{A_NS}}}buChar'); buChar.set('char', char)


def add_annotation_block(slide, bullets, x, y, w, h, font_size=14,
                         has_border=False, bg_color=None):
    """テキスト注釈ブロック（チャート横の分析テキスト）

    （task3）■マーカーは廃止し、ネイティブ箇条書き（クリムゾンの・）でチャートを補足する。
    各bullet: 最大2行、14pt。全体で3-5項目を推奨。
    箱の高さから行数予算を見積もり、収まらない項目は切り詰め/省略する（フォントは縮めない）。
    """
    # （A2）固定フォントのまま箱に収める: 行数予算で項目を取捨し、はみ出しを根絶
    cpl, _, _ = _text_capacity(w - 0.24, h - 0.16, font_size)
    line_h = font_size * 1.4 / 72.0
    budget_lines = max(1, int((h - 0.16) / line_h))
    fitted, used_lines = [], 0.0
    for b in bullets:
        vis = len(re.sub(r'\*\*', '', b))
        need = max(1, -(-vis // cpl))          # 行数(ceil)
        extra = 0.3 if fitted else 0.0          # 項目間 space_after(6pt) を約0.3行で加味
        if used_lines + need + extra <= budget_lines:
            fitted.append(b); used_lines += need + extra
        else:
            remain = int(budget_lines - used_lines - extra)
            if remain >= 1:
                fitted.append(_clip_rich(b, remain * cpl, where="チャート注釈"))
            else:
                print(f"[WARN] チャート注釈が多く{len(bullets) - len(fitted)}項目を省略 — 統合を推奨")
            break
    bullets = fitted

    if bg_color or has_border:
        box = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(h)
        )
        if bg_color:
            box.fill.solid()
            box.fill.fore_color.rgb = bg_color
        else:
            box.fill.background()
        if has_border:
            box.line.color.rgb = BORDER_GRAY
            box.line.width = Emu(9525)
        else:
            box.line.fill.background()

    txBox = slide.shapes.add_textbox(
        Inches(x + 0.12), Inches(y + 0.08),
        Inches(w - 0.24), Inches(h - 0.16)
    )
    tf = txBox.text_frame
    tf.word_wrap = True
    tf.auto_size = MSO_AUTO_SIZE.NONE

    for i, item in enumerate(bullets):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.space_after = Pt(6)
        parts = re.split(r'(\*\*.*?\*\*)', item)
        for part in parts:
            if not part:
                continue
            run = p.add_run()
            if part.startswith('**') and part.endswith('**'):
                run.text = part[2:-2]
                run.font.bold = True
                run.font.color.rgb = NAVY
            else:
                run.text = part
                run.font.color.rgb = DARK_GRAY
            run.font.size = Pt(font_size)
            _apply_font(run)
        _apply_kinsoku(p)
        p.line_spacing = 1.4
        _set_bullet(p, indent_in=0.22)  # （task3）ネイティブ箇条書き化


def add_chart_label(slide, text, x, y, w=3.0, size=14, color=NAVY):
    """チャート小見出し（グラフ上の分類ラベル）— 見出しは明朝"""
    txBox = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(0.35))
    set_text(txBox.text_frame.paragraphs[0], text, Pt(size), color, bold=True, heading=True)


def add_chart_callout(slide, text, x, y, w=2.5,
                      arrow_to_x=None, arrow_to_y=None,
                      bg_color=None, font_size=12, border_color=None):
    """チャート上の吹き出し注釈（画像の上にオーバーレイ配置）"""
    bg_color = bg_color or WHITE
    border_color = border_color or NAVY

    chars_per_line = int(w * 7)
    num_lines = max(1, -(-len(text) // chars_per_line))
    h = 0.15 + num_lines * 0.28

    box = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(h)
    )
    box.fill.solid()
    box.fill.fore_color.rgb = bg_color
    box.line.color.rgb = border_color
    box.line.width = Emu(12700)

    txBox = slide.shapes.add_textbox(
        Inches(x + 0.08), Inches(y + 0.04), Inches(w - 0.16), Inches(h - 0.08)
    )
    tf = txBox.text_frame
    tf.word_wrap = True
    tf.auto_size = MSO_AUTO_SIZE.NONE
    add_rich_runs(tf.paragraphs[0], text, base_size=Pt(font_size),
                  base_color=NAVY, bold_color=NAVY, line_spacing=1.3)

    if arrow_to_x is not None and arrow_to_y is not None:
        cx = x + w / 2
        cy = y + h / 2
        connector = slide.shapes.add_connector(
            1, Inches(cx), Inches(cy), Inches(arrow_to_x), Inches(arrow_to_y)
        )
        connector.line.color.rgb = border_color
        connector.line.width = Emu(12700)
        ln = connector._element.find('.//' + f'{{{A_NS}}}ln')
        if ln is not None:
            tailEnd = etree.SubElement(ln, f'{{{A_NS}}}tailEnd')
            tailEnd.set('type', 'triangle')
            tailEnd.set('w', 'med')
            tailEnd.set('len', 'med')

    return box


def add_highlight_circle(slide, x, y, w=0.5, h=0.5, color=None):
    """チャート上のハイライト丸囲み"""
    color = color or RED_ACCENT
    oval = slide.shapes.add_shape(
        MSO_SHAPE.OVAL, Inches(x), Inches(y), Inches(w), Inches(h)
    )
    oval.fill.background()
    oval.line.color.rgb = color
    oval.line.width = Emu(19050)
    return oval


WORDMARK = ""   # 表紙のプラットフォーム・ワードマーク。空文字なら一切描かない（ブランド名を頁に出さない）

def add_title_slide(prs, title, subtitle, date, blank,
                    kicker="TECHNOLOGY INTELLIGENCE REPORT",
                    umap_points=None, emerging_cids=None):
    """表紙 — 黒背景の大胆なエディトリアル演出（デッキの第一印象を決める頁）

    構成（柱0「インパクト演出」）:
      - 左端フル丈クリムゾン・ストリップ（強いアンカー。章扉と統一）
      - 背面に巨大ゴースト・ワードマーク "APOLLO"（低コントラストで奥行き）
      - キッカー（赤・字間広め・全大文字）+ APOLLO ワードマーク（白・字間広め）
      - タイトル直上の太いクリムゾン罫 + 大判明朝タイトル（最大48pt）
      - 下部クリムゾン帯に日付を白文字反転で乗せる
    kicker はレポート種別ラベル（例 "PATENT LANDSCAPE REPORT"）。"""
    slide = prs.slides.add_slide(blank)
    _hide_master(slide)

    # 白基調: 発明スライド由来の「白×赤斜線」背景画像を全面に敷く（無ければ白ベタ）
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
    if CV_BG_LIGHT_PATH and os.path.exists(CV_BG_LIGHT_PATH):
        slide.shapes.add_picture(CV_BG_LIGHT_PATH, Inches(0), Inches(0), Inches(13.333), Inches(7.5))
    # 左端の極細クリムゾン・ストリップ（赤の背骨＝アンカー）
    _cv_rect(slide, 0, 0, 0.12, 7.5, _CV["crimson"], 100)

    # キッカー（Century Gothic・クリムゾン・タイトルの直上に配置）
    _cv_txt(slide, kicker, 0.97, 2.52, 8.2, 0.36, 12.5, _CV["crimson"], "Century Gothic", False, PP_ALIGN.LEFT, 2.4)
    # ワードマーク（任意・淡グレー）
    if WORDMARK:
        _cv_txt(slide, WORDMARK, 0.97, 1.30, 6.0, 0.4, 11.5, "8A8A8A", "Century Gothic", False, PP_ALIGN.LEFT, 5.0)

    # タイトル（Century Gothic 細字・墨・大きめ・縦中央。長さで級数を自動調整）
    tlen = len(title)
    t_size = 50 if tlen <= 12 else (42 if tlen <= 16 else (34 if tlen <= 22 else 30))
    _cv_txt(slide, title, 0.95, 3.06, 8.3, 1.30, t_size, "1A1A1E", "Century Gothic", False,
            PP_ALIGN.LEFT, 0.4, MSO_ANCHOR.MIDDLE)
    # サブタイトル（対象・件数のみ／濃灰）
    _cv_txt(slide, subtitle, 0.98, 4.62, 7.6, 0.8, 12.5, "55585E", _CV_GO, False, PP_ALIGN.LEFT)

    # 日付（YYYY/MM/DD・Century Gothic・墨グレー・やや上）
    _cv_txt(slide, date, 0.97, 6.55, 9.5, 0.32, 11.5, "55585E", "Century Gothic", False, PP_ALIGN.LEFT, 0.8)
    return slide


# ============================================================
# Crimson Vector V11 ベクター技法（章扉・結論ピラミッド）
# 画像を使わず custGeom(freeform) 4点ポリゴンで連続ピラミッドを描く。
# 各段の頂点を「共通の apex/base を結ぶ直線」上に取り、辺を一直線に連続させる
# （= 台形の積み重ね＝鏡持ち を回避し、1つの大きなピラミッドに見せる）。
# ============================================================
_CV = {"black":"020202","black3":"141414","white":"FFFFFF","ivory":"F6F2EC",
       "crimson":"C7001E","crimson2":"8E0014","crimson3":"4A0008","redDark":"220003","gray":"8A8A8A"}
_CV_SERIF = "Times New Roman"; _CV_MIN = "Yu Mincho"; _CV_GO = "Yu Gothic"
_CV_EMU = 914400
def _cvC(h): return RGBColor(int(h[0:2],16),int(h[2:4],16),int(h[4:6],16))
def _cv_line_alpha(shp, op):
    ln = shp._element.spPr.find(f'{{{A_NS}}}ln')
    if ln is None: return
    f = ln.find(f'{{{A_NS}}}solidFill')
    if f is None: return
    sg = f.find(f'{{{A_NS}}}srgbClr')
    if sg is None: return
    for a in sg.findall(f'{{{A_NS}}}alpha'): sg.remove(a)
    etree.SubElement(sg, f'{{{A_NS}}}alpha').set('val', str(int(op*1000)))
def _cv_rect(s,x,y,w,h,col,op=100,lcol=None,lop=100,lw=0.5):
    sh=s.shapes.add_shape(MSO_SHAPE.RECTANGLE,Inches(x),Inches(y),Inches(w),Inches(h))
    sh.fill.solid(); sh.fill.fore_color.rgb=_cvC(col)
    if op<100: _set_fill_alpha(sh,op)
    if lcol: sh.line.color.rgb=_cvC(lcol); sh.line.width=Pt(lw); _cv_line_alpha(sh,lop)
    else: sh.line.fill.background()
    return sh
def _cv_shape(s,mso,x,y,w,h,col,op=100,lcol=None,lop=100,lw=0.5,rot=0):
    sh=s.shapes.add_shape(mso,Inches(x),Inches(y),Inches(w),Inches(h))
    if rot: sh.rotation=rot
    sh.fill.solid(); sh.fill.fore_color.rgb=_cvC(col)
    if op<100: _set_fill_alpha(sh,op)
    if lcol: sh.line.color.rgb=_cvC(lcol); sh.line.width=Pt(lw); _cv_line_alpha(sh,lop)
    else: sh.line.fill.background()
    return sh
def _cv_ln(s,x,y,dx,dy,col="FFFFFF",op=50,lw=0.5):
    c=s.shapes.add_connector(MSO_CONNECTOR.STRAIGHT,Inches(x),Inches(y),Inches(x+dx),Inches(y+dy))
    c.line.color.rgb=_cvC(col); c.line.width=Pt(lw); _cv_line_alpha(c,op); return c
def _cv_poly(s,pts,col,op=100,lcol="FFFFFF",lop=30,lw=0.6):
    fb=s.shapes.build_freeform(Emu(int(pts[0][0]*_CV_EMU)),Emu(int(pts[0][1]*_CV_EMU)),scale=1)
    fb.add_line_segments([(Emu(int(x*_CV_EMU)),Emu(int(y*_CV_EMU))) for x,y in pts[1:]],close=True)
    sh=fb.convert_to_shape(); sh.fill.solid(); sh.fill.fore_color.rgb=_cvC(col)
    if op<100: _set_fill_alpha(sh,op)
    if lcol is not None: sh.line.color.rgb=_cvC(lcol); sh.line.width=Pt(lw); _cv_line_alpha(sh,lop)
    else: sh.line.fill.background()
    return sh
def _cv_grad_alpha(sh, op):
    grad=sh._element.spPr.find(f'{{{A_NS}}}gradFill')
    if grad is None: return
    for srgb in grad.iter(f'{{{A_NS}}}srgbClr'):
        for a in srgb.findall(f'{{{A_NS}}}alpha'): srgb.remove(a)
        etree.SubElement(srgb,f'{{{A_NS}}}alpha').set('val',str(int(op*1000)))
def _cv_grad_poly(s,pts,c1,c2,angle=90,lcol=None,lop=30,lw=0.6,op=100):
    fb=s.shapes.build_freeform(Emu(int(pts[0][0]*_CV_EMU)),Emu(int(pts[0][1]*_CV_EMU)),scale=1)
    fb.add_line_segments([(Emu(int(x*_CV_EMU)),Emu(int(y*_CV_EMU))) for x,y in pts[1:]],close=True)
    sh=fb.convert_to_shape(); sh.fill.gradient()
    gs=sh.fill.gradient_stops; gs[0].color.rgb=_cvC(c1); gs[0].position=0.0; gs[1].color.rgb=_cvC(c2); gs[1].position=1.0
    try: sh.fill.gradient_angle=angle
    except Exception: pass
    if op<100: _cv_grad_alpha(sh,op)
    if lcol is not None: sh.line.color.rgb=_cvC(lcol); sh.line.width=Pt(lw); _cv_line_alpha(sh,lop)
    else: sh.line.fill.background()
    return sh
def _cv_grad_shape(s,mso,x,y,w,h,c1,c2,angle=90,lcol=None,lop=30,lw=0.6,rot=0,op=100):
    sh=s.shapes.add_shape(mso,Inches(x),Inches(y),Inches(w),Inches(h))
    if rot: sh.rotation=rot
    sh.fill.gradient()
    gs=sh.fill.gradient_stops; gs[0].color.rgb=_cvC(c1); gs[0].position=0.0; gs[1].color.rgb=_cvC(c2); gs[1].position=1.0
    try: sh.fill.gradient_angle=angle
    except Exception: pass
    if op<100: _cv_grad_alpha(sh,op)
    if lcol is not None: sh.line.color.rgb=_cvC(lcol); sh.line.width=Pt(lw); _cv_line_alpha(sh,lop)
    else: sh.line.fill.background()
    return sh
def _cv_txt_two(s,l1,l2,x,y,w,s1,s2,c1="FFFFFF",c2="F6F2EC"):
    tb=s.shapes.add_textbox(Inches(x),Inches(y),Inches(w),Inches(1.2)); tf=tb.text_frame
    tf.word_wrap=True; tf.auto_size=MSO_AUTO_SIZE.NONE
    tf.margin_left=0; tf.margin_right=0; tf.margin_top=0; tf.margin_bottom=0
    def _mk(p,t,sz,col,font,bold):
        r=p.add_run(); r.text=str(t).replace("**",""); r.font.size=Pt(sz); r.font.bold=bold
        r.font.color.rgb=_cvC(col); r.font.name=font
        rPr=r._r.get_or_add_rPr(); rPr.set('lang','ja-JP')
        ea=rPr.find(f'{{{A_NS}}}ea')
        if ea is None: ea=etree.SubElement(rPr,f'{{{A_NS}}}ea')
        ea.set('typeface', _CV_MIN if font in (_CV_MIN,_CV_SERIF) else _CV_GO)
    p1=tf.paragraphs[0]; p1.alignment=PP_ALIGN.CENTER; _mk(p1,l1,s1,c1,_CV_MIN,True)
    p2=tf.add_paragraph(); p2.alignment=PP_ALIGN.CENTER; _mk(p2,l2,s2,c2,_CV_GO,False)
    return tb
def _cv_txt(s,t,x,y,w,h,size,col,font="Yu Gothic",bold=False,align=PP_ALIGN.LEFT,spc=None,anchor=None):
    tb=s.shapes.add_textbox(Inches(x),Inches(y),Inches(w),Inches(h)); tf=tb.text_frame
    tf.word_wrap=True; tf.auto_size=MSO_AUTO_SIZE.NONE
    tf.margin_left=0; tf.margin_right=0; tf.margin_top=0; tf.margin_bottom=0
    if anchor: tf.vertical_anchor=anchor
    p=tf.paragraphs[0]; p.alignment=align
    lines=str(t).replace("**","").split("\n")
    for k,ln_ in enumerate(lines):
        r=p.add_run(); r.text=ln_; r.font.size=Pt(size); r.font.bold=bold
        r.font.color.rgb=_cvC(col); r.font.name=font
        rPr=r._r.get_or_add_rPr(); rPr.set('lang','ja-JP')
        ea=rPr.find(f'{{{A_NS}}}ea')
        if ea is None: ea=etree.SubElement(rPr,f'{{{A_NS}}}ea')
        ea.set('typeface', _CV_MIN if font in (_CV_MIN,_CV_SERIF) else _CV_GO)
        if spc: rPr.set('spc',str(int(spc*100)))
        if k<len(lines)-1: etree.SubElement(p._p,f'{{{A_NS}}}br')
    return tb
def _cv_base_stage(s, red=True, acc=None):
    """黒の舞台＋スキャンライン＋右の建築面（cinematic base）。
    acc=章ごとの差し色セット（赤/紫）。建築面は基本グラデで描く。"""
    acc = acc or _CV
    s.background.fill.solid(); s.background.fill.fore_color.rgb=_cvC(_CV["black"])
    for i in range(18):
        _cv_rect(s,i*0.42,0,0.43,7.5, _CV["black"] if i<7 else _CV["black3"], 100 if i<7 else max(0,30-i))
    if red:
        _cv_grad_shape(s,MSO_SHAPE.PARALLELOGRAM,9.28,-0.15,4.45,7.85, acc["crimson"], acc["crimson3"], 50, acc["crimson"],8,0.2,-11, op=88)
        _cv_grad_shape(s,MSO_SHAPE.PARALLELOGRAM,10.15,-0.35,3.40,7.95, acc["crimson2"], acc["crimson3"], 62, acc["crimson"],16,0.35,-11, op=80)
        _cv_grad_shape(s,MSO_SHAPE.ISOSCELES_TRIANGLE,10.38,0.10,2.55,5.95, acc["crimson"], acc["crimson3"], 70, acc["crimson"],18,0.35,0, op=48)
        _cv_shape(s,MSO_SHAPE.ISOSCELES_TRIANGLE,11.02,0.55,1.65,4.70,_CV["black"],22, acc["crimson"],30,0.3,0)
        _cv_shape(s,MSO_SHAPE.PARALLELOGRAM,8.05,0.0,1.85,7.7,_CV["black"],66,_CV["white"],12,0.15,-17)
    _cv_ln(s,8.0,-0.1,-2.4,7.7,_CV["white"],16,0.4)
    _cv_ln(s,9.95,-0.2,-2.0,7.8, acc["crimson"],18,0.6)
    _cv_ln(s,12.65,0.0,-4.3,7.55,_CV["white"],22,0.45)
    for i in range(24): _cv_ln(s,9.10+i*0.13,0.0,0,7.5, acc["crimson"],max(4,18-i),0.25)
    _cv_shape(s,MSO_SHAPE.TRAPEZOID,6.25,6.46,3.8,0.26, acc["crimson"],15, acc["crimson"],10,0.15,0)
CV_BG_PATH = ""
CV_BG_LIGHT_PATH = ""   # v16 暗赤背景画像のパス（build側で設定）。空なら暗色ベタで代替。
def _hide_master(slide):
    """暗背景スライドでマスター図形（署名ヘアライン等）を非表示にする。"""
    slide._element.set('showMasterSp', '0')


def _cv_bg_image(slide, overlay=72):
    """v16: 暗赤背景画像を全面に敷き、黒の半透明オーバーレイで可読性を確保する。"""
    _hide_master(slide)
    if CV_BG_PATH and os.path.exists(CV_BG_PATH):
        slide.shapes.add_picture(CV_BG_PATH, Inches(0), Inches(0), Inches(13.333), Inches(7.5))
    else:
        slide.background.fill.solid(); slide.background.fill.fore_color.rgb = _cvC("050607")
    ov = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(13.333), Inches(7.5))
    ov.fill.solid(); ov.fill.fore_color.rgb = _cvC("000000"); _set_fill_alpha(ov, 100 - overlay)
    ov.line.fill.background()
def _cv_panel(slide, x, y, w, h, fill="0A0C10", op=48, line="681118", lop=84, lw=1.05, radius=0.06):
    """v16: 文字直下の半透明ラウンド矩形パネル（背景画像の上で可読性を出す）。"""
    sh = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE if radius > 0 else MSO_SHAPE.RECTANGLE,
                                Inches(x), Inches(y), Inches(w), Inches(h))
    if radius > 0:
        try: sh.adjustments[0] = radius
        except Exception: pass
    sh.fill.solid(); sh.fill.fore_color.rgb = _cvC(fill); _set_fill_alpha(sh, op)
    sh.line.color.rgb = _cvC(line); sh.line.width = Pt(lw); _cv_line_alpha(sh, lop)
    return sh
# v16 パレット（赤/黒/白/グレーのみ）
_V16 = {"white":"F6F7F8","muted":"D4D8DE","muted2":"AEB5BF","red":"FF3030",
        "panel":"0A0C10","panel2":"10141A","line":"681118","line2":"8B1A22"}
def _cv_xon(x1,y1,x2,y2,y): return x1+(x2-x1)*((y-y1)/(y2-y1))
def _cv_band(ap,lb,rb,yt,yb):
    return [(_cv_xon(*ap,*lb,yt),yt),(_cv_xon(*ap,*rb,yt),yt),(_cv_xon(*ap,*rb,yb),yb),(_cv_xon(*ap,*lb,yb),yb)]
def _cv_inset(ap,lb,rb,yt,yb,it,ib,tr):
    xlt=_cv_xon(*ap,*lb,yt)+it; xrt=_cv_xon(*ap,*rb,yt)-it; xlb=_cv_xon(*ap,*lb,yb)+ib; xrb=_cv_xon(*ap,*rb,yb)-ib
    return [(xlt+tr,yt+0.10),(xrt-tr,yt+0.10),(xrb-tr,yb-0.12),(xlb+tr,yb-0.12)]
def _cv_lfacet(ap,lb,yt,yb,inn):
    xot=_cv_xon(*ap,*lb,yt); xob=_cv_xon(*ap,*lb,yb)
    return [(xot,yt),(xot+inn,yt+0.08),(xob+inn,yb-0.08),(xob,yb)]
def _cv_rfacet(ap,rb,yt,yb,inn):
    xot=_cv_xon(*ap,*rb,yt); xob=_cv_xon(*ap,*rb,yb)
    return [(xot-inn,yt+0.08),(xot,yt),(xob,yb),(xob-inn,yb-0.08)]
def _cv_insight_box(s,idx,title,body,x,y,w,h):
    _cv_rect(s,x,y,w,h,_CV["black"],92,_CV["crimson"],70,0.9)
    _cv_rect(s,x,y,0.08,h,_CV["crimson"],100)
    _cv_rect(s,x+0.08,y+0.08,w-0.16,h-0.16,_CV["crimson3"],28,_CV["crimson"],14,0.2)
    _cv_txt(s,f"{idx:02d}",x+0.18,y+0.16,0.6,0.34,19,_CV["crimson"],_CV_SERIF,True)
    _cv_txt(s,title,x+0.86,y+0.16,w-1.1,0.32,13.5,_CV["white"],_CV_MIN,True)
    _cv_ln(s,x+0.86,y+0.55,w-1.1,0,_CV["crimson"],78,0.42)
    _cv_txt(s,body,x+0.86,y+0.62,w-1.12,h-0.74,8.6,_CV["ivory"],_CV_GO,False)


def _tint(rgb, f):
    """色を白へ f(0-1) だけ寄せた淡色を返す（KPI強調カード地など）。"""
    return RGBColor(int(rgb[0] + (255 - rgb[0]) * f),
                    int(rgb[1] + (255 - rgb[1]) * f),
                    int(rgb[2] + (255 - rgb[2]) * f))


def _apply_chapter_theme(section_num, total):
    """章ごとに差し色を切り替える: 赤→紫→赤→紫（最初(1)/最後(total)は赤固定）。
    本文が参照する ACCENT / BRIGHT_RED / CRIMSON_DEEP（モジュール全体で共有）も切り替えるため、
    紫章は章扉だけでなく**本文のテーマカラーも紫**（表ヘッダ・KPI強調・考察ラベル・章ラベル等）になる。
    戻り値は章扉の建築面用 hex セット。"""
    global ACCENT, BRIGHT_RED, CRIMSON_DEEP
    is_purple = (section_num % 2 == 0) and section_num != 1 and section_num != int(total or 0)
    if is_purple:
        ACCENT = RGBColor(0xA1, 0x2F, 0xC4)        # 明るめの紫（テーマカラー）
        BRIGHT_RED = RGBColor(0xC3, 0x55, 0xE0)    # グラデ明端（紫）
        CRIMSON_DEEP = RGBColor(0x5E, 0x1C, 0x8A)  # グラデ深端（紫）
        acc = dict(_CV)
        acc.update({"crimson": "A12FC4", "crimson2": "5E1C8A",
                    "crimson3": "2E0B4A", "redDark": "180428"})
        return acc
    ACCENT = RGBColor(0xC5, 0x12, 0x12)            # 既定＝クリムゾン（メインカラー）
    BRIGHT_RED = RGBColor(0xF4, 0x33, 0x33)
    CRIMSON_DEEP = RGBColor(0x8E, 0x00, 0x14)
    return _CV


def add_section_slide(prs, section_num, title, blank, subtitle=None, en=None):
    """章扉（Crimson Vector V11）：黒の舞台＋巨大番号＋ゴースト＋明朝大見出し＋
    右の建築オブジェ（基本グラデ）。差し色は章ごとに赤/紫。CURRENT_CHAPTER を更新する。"""
    global CURRENT_CHAPTER
    CURRENT_CHAPTER = section_num
    _PENDING_SECTION[0] = section_num
    s = prs.slides.add_slide(blank)
    _hide_master(s)

    acc = _apply_chapter_theme(section_num, TOTAL_CHAPTERS)   # 章扉＋以降の本文の差し色を切替
    _cv_base_stage(s, red=True, acc=acc)
    num = f"{section_num:02d}"
    # メイン数字（巨大番号＝主役）。ゴーストは背面＝シャドウとして本体に約2/3重ねる（右下へオフセット）
    _cv_txt(s, num, 1.40,-0.30,6.2,3.6, 220, _CV["black3"], _CV_SERIF, True, PP_ALIGN.LEFT, -4)   # ゴースト（背面）
    _numbox = _cv_txt(s, num, 0.40,-0.55,9.4,3.6, 236, acc["crimson"], _CV_SERIF, True, PP_ALIGN.LEFT, -4)  # 本体（前面）
    _grad_text_fill(_numbox, acc["crimson"], acc["crimson2"])  # 大番号グラデ（章の差し色）
    _cv_ln(s, 0.58,2.85,2.35,0, acc["crimson"], 100, 0.95)   # 横線は数字の下（重ならない位置）へ
    _cv_txt(s, title, 0.56,3.08,8.2,1.18, 40, _CV["white"], _CV_MIN, True, PP_ALIGN.LEFT, 0.2)
    if subtitle:
        _cv_txt(s, "/  "+subtitle, 0.60,4.22,8.2,0.5, 14, _CV["white"], _CV_MIN, True)
    _cv_txt(s, f"SECTION {section_num:02d}", 0.62,4.92,4.2,0.3, 8, _CV["gray"], _CV_GO, False, PP_ALIGN.LEFT, 2.0)
    if en:
        _cv_txt(s, en, 0.62,5.26,4.4,0.6, 7.5, _CV["gray"], _CV_GO, False, PP_ALIGN.LEFT, 1.2)
    # 右の建築オブジェ（1.5倍に拡大・ずらし＋角度/透明度の変化でメリハリ）
    _cv_grad_shape(s, MSO_SHAPE.PARALLELOGRAM, 8.50,-1.33,2.18,9.45, acc["crimson"], acc["crimson3"], 52, acc["crimson"], 40, 0.6, -15, op=94)
    _cv_grad_shape(s, MSO_SHAPE.PARALLELOGRAM, 9.87,-0.06,1.58,7.28, acc["crimson2"], acc["redDark"], 46, acc["crimson"], 22, 0.5, -3, op=64)
    _cv_shape(s, MSO_SHAPE.PARALLELOGRAM, 10.62,-1.03,0.90,9.15, _CV["black"], 70, _CV["white"], 16, 0.55, -11)
    _cv_grad_shape(s, MSO_SHAPE.PARALLELOGRAM, 11.08,0.35,1.50,6.60, acc["crimson"], acc["crimson3"], 66, acc["crimson"], 16, 0.4, 2, op=40)
    _cv_ln(s, 9.55,-1.03,0.0,8.55, _CV["white"], 20, 0.45)
    _cv_ln(s, 11.65,-0.05,0.0,6.9, acc["crimson"], 60, 1.2)
    return s

def add_chart_text_slide(prs, title, sub_message, image_path, annotations, blank,
                         caption=None, chart_label=None, text_side="right",
                         chart_ratio=0.715, source=None, page_num=None):
    """チャート主体 + テキスト注釈 — 主力スライドタイプ

    Args:
        annotations: ["短い注釈1", "短い注釈2", ...] — 各1-2行、最大5項目
        text_side: "right" or "left"
        chart_ratio: チャート側の幅比率（既定0.68）。
            （改良4）チャートを約2割拡大するため既定を 0.60→0.68 に引き上げ、
            縦方向も下端まで使い切る。注釈は3-5項目の短文に絞ること。
    """
    slide = prs.slides.add_slide(blank)
    sub_y = add_title_shape(slide, title)
    content_y = add_sub_message(slide, sub_message, y=sub_y)

    content_w = CONTENT_W  # （改良4）右余白も使ってチャート域を拡張
    content_x = MARGIN_L
    gap = 0.04  # （task2→さらに寄せる）注釈をグラフへ密着。チャートは chart_ratio 0.715 で約5%拡大
    chart_w = content_w * chart_ratio - gap / 2
    text_w = content_w * (1 - chart_ratio) - gap / 2
    remaining_h = 6.7 - content_y  # （改良4）下端まで使い切り縦も拡大

    if text_side == "right":
        chart_x = content_x
        text_x = content_x + chart_w + gap
    else:
        text_x = content_x
        chart_x = content_x + text_w + gap

    # チャート小見出し
    if chart_label:
        add_chart_label(slide, chart_label, chart_x, content_y, chart_w)
        img_y = content_y + 0.35
        img_h = remaining_h - 0.65
    else:
        img_y = content_y
        img_h = remaining_h - 0.3

    # チャート画像（左上揃えで領域に配置）。縦長画像でも左右に死に余白を作らない
    full_path = os.path.join(SNAP, image_path) if not os.path.isabs(image_path) else image_path
    img_halign = "left" if text_side == "right" else "right"
    pic = fit_image(slide, full_path, max_x=chart_x, max_y=img_y,
                    max_w=chart_w, max_h=img_h, halign=img_halign, valign="top")

    # 実際の描画矩形に合わせて注釈を密着させ、左右バランスを締める
    cap_x, cap_w = chart_x, chart_w
    if pic is not None and getattr(pic, "_fit_rect", None):
        real_left, _real_top, real_w, _real_h = pic._fit_rect
        cap_x, cap_w = real_left, real_w
        if text_side == "right":
            text_x = real_left + real_w + 0.30          # 実描画右端＋わずかな間隔
            text_w = max(2.7, content_x + content_w - text_x)
        else:
            text_w = max(2.7, real_left - 0.30 - content_x)

    # キャプション（実描画域の中央に置く）
    if caption:
        txBox = slide.shapes.add_textbox(Inches(cap_x), Inches(content_y + remaining_h - 0.25),
                                          Inches(cap_w), Inches(0.25))
        set_text(txBox.text_frame.paragraphs[0], caption, Pt(10), MEDIUM_GRAY)
        txBox.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER

    # テキスト注釈（14pt、3-5項目）
    add_annotation_block(slide, annotations[:5], text_x, content_y,
                         text_w, remaining_h - 0.2)

    if source:
        add_source_label(slide, source)
    add_corner_marks(slide)  # （A）deco廃止: 明背景の分析頁にエンジンが自動でコーナーマーク
    add_bottom_bar_and_footer(slide, page_num)
    return slide


ICON_FONT = "Material Symbols Outlined"  # （改良6）KPIカードのアイコン用フォント

# Material Symbols の TTF / codepoints の探索先（環境で上書き可）。
# - グローバル MS_FONT_PATH / MS_CP_PATH を定義すればそれを最優先で使う
# - スキル同梱 assets/ → 環境変数 APOLLO_MS_FONT の順に探索
def _ms_font_paths():
    cands_ttf, cands_cp = [], []
    g = globals()
    if g.get("MS_FONT_PATH"): cands_ttf.append(g["MS_FONT_PATH"])
    if g.get("MS_CP_PATH"): cands_cp.append(g["MS_CP_PATH"])
    here = os.path.dirname(g["__file__"]) if g.get("__file__") else "."
    cands_ttf += [
        os.path.join(here, "../../.claude/skills/apollo-pptx/assets/MaterialSymbolsOutlined.ttf"),
        os.environ.get("APOLLO_MS_FONT", ""),
    ]
    cands_cp += [
        os.path.join(here, "../../.claude/skills/apollo-pptx/assets/MaterialSymbolsOutlined.codepoints"),
        os.environ.get("APOLLO_MS_CP", ""),
    ]
    ttf = next((p for p in cands_ttf if p and os.path.exists(p)), None)
    cp = next((p for p in cands_cp if p and os.path.exists(p)), None)
    return ttf, cp

_MS_CP_CACHE = {}
def _ms_codepoint(name):
    if not _MS_CP_CACHE:
        _, cp_path = _ms_font_paths()
        if cp_path:
            for line in open(cp_path, encoding="utf-8"):
                parts = line.split()
                if len(parts) == 2:
                    _MS_CP_CACHE[parts[0]] = int(parts[1], 16)
    return _MS_CP_CACHE.get(name)

def material_icon_png(name, color, px=200):
    """Material Symbols のアイコンを透過PNGにラスタライズしてパスを返す。

    配布先PCにフォントが無くても確実に表示できるよう**画像化して埋め込む**ための助関数。
    フォント/コードポイントが見つからない場合は None（呼び出し側はアイコンを省略）。
    color は (r,g,b) タプル推奨。
    """
    import tempfile
    from PIL import Image, ImageDraw, ImageFont
    ttf, _ = _ms_font_paths()
    cpv = _ms_codepoint(name)
    if not ttf or cpv is None:
        return None
    cdir = os.path.join(tempfile.gettempdir(), "apollo_ms_icons")
    os.makedirs(cdir, exist_ok=True)
    hexc = "%02x%02x%02x" % (color[0], color[1], color[2]) if isinstance(color, (tuple, list)) else str(color)
    out = os.path.join(cdir, f"{name}_{hexc}_{px}.png")
    if not os.path.exists(out):
        rgb = (color[0], color[1], color[2]) if isinstance(color, (tuple, list)) else (
            color >> 16 & 255, color >> 8 & 255, color & 255)
        img = Image.new("RGBA", (px, px), (0, 0, 0, 0))
        d = ImageDraw.Draw(img)
        fnt = ImageFont.truetype(ttf, int(px * 0.82))
        d.text((px / 2, px / 2), chr(cpv), font=fnt, fill=rgb + (255,), anchor="mm")
        img.save(out)
    return out

def _rgb_tuple(c):
    return (c[0], c[1], c[2]) if isinstance(c, (tuple, list)) else (c >> 16 & 255, c >> 8 & 255, c & 255)


def _set_run_font(run, name):
    """run に任意フォント（欧文・和文とも同名）を設定する。アイコンフォント用。"""
    run.font.name = name
    rPr = run._r.get_or_add_rPr()
    rPr.set('lang', 'en-US')
    ea = rPr.find(f'{{{A_NS}}}ea')
    if ea is None:
        ea = etree.SubElement(rPr, f'{{{A_NS}}}ea')
    ea.set('typeface', name)


def add_kpi_slide(prs, title, sub_message, kpis, blank,
                  source=None, page_num=None):
    """KPIダッシュボード — エディトリアル・カード（改良6）

    kpis: [{"label":"総特許件数", "value":"1,176", "unit":"件",
            "trend":"↑", "icon":"description", "emphasis":True}, ...]
      icon     : Material Symbols Outlined のリガチャ名（例 "co2","trending_up",
                 "groups","hub","blur_on","recycling","eco","trending_down"）。
                 環境にフォントが無い場合はリガチャ名がそのまま小さく表示される。
      emphasis : True で「最も言いたい1-2枚」をクリムゾンで強調（赤9:1の原則）。
    設計: **枠線なし・大きな数字(40pt)・上部アクセント帯・アイコン**。
    任意の枚数に対応（n_cols/n_rows を動的計算）。
    """
    slide = prs.slides.add_slide(blank)
    sub_y = add_title_shape(slide, title)
    content_y = add_sub_message(slide, sub_message, y=sub_y)

    n = len(kpis)
    available_w = CONTENT_W
    start_x = MARGIN_L
    gap = 0.24

    # 動的グリッド（任意枚数に対応）
    if n <= 4:
        n_cols, n_rows = n, 1
    else:
        n_cols = min(4, -(-n // 2))           # 5-8→4列, 9→ceil(9/2)=5→4列
        n_rows = -(-n // n_cols)               # 必要行数

    card_w = (available_w - gap * (n_cols - 1)) / n_cols
    available_h = 6.6 - content_y
    row_gap = 0.24
    card_h = (available_h - row_gap * (n_rows - 1)) / n_rows
    card_h = min(card_h, 2.2)

    for idx, kpi in enumerate(kpis):
        row, col = idx // n_cols, idx % n_cols
        x = start_x + col * (card_w + gap)
        y = content_y + row * (card_h + row_gap)
        emph = bool(kpi.get("emphasis"))

        # カード背景（枠線なし）。強調カードは淡クリムゾンのグラデ地で“重要部分”を演出（#7）。
        card = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(x), Inches(y),
                                      Inches(card_w), Inches(card_h))
        if emph:
            grad_fill(card, _tint(ACCENT, 0.93), _tint(ACCENT, 0.82), angle=45)  # テーマ色の淡地
        else:
            card.fill.solid(); card.fill.fore_color.rgb = PALE_GRAY
            card.line.fill.background()
        try: card.shadow.inherit = False
        except Exception: pass

        # 強調カードのみ左端にクリムゾン→明赤のグラデ・アクセントバー（赤9:1・視覚ノイズ最小）
        if emph:
            ab = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(x), Inches(y),
                                        Inches(0.075), Inches(card_h))
            grad_fill(ab, ACCENT, BRIGHT_RED, angle=90)

        # アイコン（Material Symbols をPNG化して右上に配置・配布先のフォント不要）
        if kpi.get("icon"):
            icol = (ACCENT[0], ACCENT[1], ACCENT[2]) if emph else (0x68, 0x68, 0x68)
            ipath = material_icon_png(kpi["icon"], icol, px=200)
            isz = 0.42
            if ipath:
                slide.shapes.add_picture(ipath, Inches(x + card_w - isz - 0.18),
                                         Inches(y + 0.18), width=Inches(isz), height=Inches(isz))
            else:
                # フォント未取得時のフォールバック: アイコンフォント名で文字描画（環境依存）
                ic = slide.shapes.add_textbox(Inches(x + card_w - 0.62), Inches(y + 0.16),
                                              Inches(0.5), Inches(0.5))
                icp = ic.text_frame.paragraphs[0]; icp.alignment = PP_ALIGN.RIGHT
                icr = icp.add_run(); icr.text = kpi["icon"]
                icr.font.size = Pt(20); icr.font.color.rgb = ACCENT if emph else MEDIUM_GRAY
                _set_run_font(icr, ICON_FONT)

        # ラベル（小・グレー）
        txL = slide.shapes.add_textbox(Inches(x + 0.18), Inches(y + 0.22),
                                       Inches(card_w - 0.78), Inches(0.5))
        tfl = txL.text_frame; tfl.word_wrap = True; tfl.auto_size = MSO_AUTO_SIZE.NONE
        set_text(tfl.paragraphs[0], kpi["label"], Pt(10), MEDIUM_GRAY, bold=True, line_spacing=1.05)

        # 値（数字＝特大／符号(+−±)・単位(%等)は数字より2サイズ小さく）
        #   例: "+38%" → "+"と"%"を小さく、"38"を大きく。"鉱物化"等テキストは中サイズ。
        BIG, SMALL = 44, 28        # 数字＝44pt、符号・単位＝28pt（2サイズ小）
        col_v = ACCENT if emph else INK
        m = re.match(r'^\s*([+\-−±▲▼]?)\s*([0-9][0-9.,]*)\s*(.*)$', str(kpi["value"]))
        txV = slide.shapes.add_textbox(Inches(x + 0.16), Inches(y + 0.60),
                                       Inches(card_w - 0.28), Inches(0.85))
        pv = txV.text_frame.paragraphs[0]
        if m:
            lead, num, suf = m.group(1), m.group(2), m.group(3)
            if lead:
                r0 = pv.add_run(); r0.text = lead; r0.font.size = Pt(SMALL); r0.font.bold = True
                r0.font.color.rgb = col_v; _apply_font(r0, heading=True)
            r1 = pv.add_run(); r1.text = num; r1.font.size = Pt(BIG); r1.font.bold = True
            r1.font.color.rgb = col_v; _apply_font(r1, heading=True)
            if suf:
                r2 = pv.add_run(); r2.text = suf; r2.font.size = Pt(SMALL); r2.font.bold = True
                r2.font.color.rgb = col_v; _apply_font(r2, heading=True)
        else:
            rt0 = pv.add_run(); rt0.text = str(kpi["value"]); rt0.font.size = Pt(30)
            rt0.font.bold = True; rt0.font.color.rgb = col_v; _apply_font(rt0, heading=True)
        if kpi.get("trend"):
            tr = kpi["trend"]
            tc = ACCENT if ("-" in tr or "↓" in tr or "DOWN" in tr.upper()) else INK
            rt = pv.add_run(); rt.text = f" {tr}"; rt.font.size = Pt(SMALL - 4)
            rt.font.bold = True; rt.font.color.rgb = tc; _apply_font(rt)

        # 単位（小・グレー）
        txU = slide.shapes.add_textbox(Inches(x + 0.18), Inches(y + card_h - 0.40),
                                       Inches(card_w - 0.32), Inches(0.34))
        tfu = txU.text_frame; tfu.word_wrap = True; tfu.auto_size = MSO_AUTO_SIZE.NONE
        set_text(tfu.paragraphs[0], kpi.get("unit", ""), Pt(9), MEDIUM_GRAY, line_spacing=1.05)

    if source:
        add_source_label(slide, source)
    add_corner_marks(slide)  # （A）deco廃止: 明背景の分析頁にエンジンが自動でコーナーマーク
    add_bottom_bar_and_footer(slide, page_num)
    return slide


def add_cards_slide(prs, title, sub_message, cards, blank,
                    source=None, page_num=None):
    """カード並列表示 — 3-4枚のカードを横並び

    cards: [{"header":"クラスタA", "body":"説明テキスト", "color":NAVY}, ...]
    ヘッダー: 色付き背景 + 白テキスト
    ボディ: LIGHT_GRAY背景 + DARK_GRAYテキスト
    """
    if len(cards) > 4:
        print(f"[WARN] add_cards_slide: cards={len(cards)} > 4。横幅が窮屈になるため2スライドに分割推奨")
    slide = prs.slides.add_slide(blank)
    sub_y = add_title_shape(slide, title)
    content_y = add_sub_message(slide, sub_message, y=sub_y)

    n = len(cards)
    gap = 0.25
    total_w = 12.3
    card_w = (total_w - gap * (n - 1)) / n
    card_h = 6.5 - content_y  # 下端まで使い切る
    header_h = 0.45
    colors = [NAVY, BLUE, ACCENT, GREEN_ACCENT, RED_ACCENT, AMBER]

    for i, card in enumerate(cards):
        x = 0.5 + i * (card_w + gap)
        color = card.get("color", colors[i % len(colors)])

        # ヘッダー（色付き背景 + 白テキスト・長方形）
        hdr = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE, Inches(x), Inches(content_y),
            Inches(card_w), Inches(header_h)
        )
        hdr.fill.solid()
        hdr.fill.fore_color.rgb = color
        hdr.line.fill.background()

        txH = slide.shapes.add_textbox(Inches(x + 0.1), Inches(content_y + 0.05),
                                        Inches(card_w - 0.2), Inches(header_h - 0.1))
        set_text(txH.text_frame.paragraphs[0], card["header"], Pt(14), WHITE, bold=True)
        txH.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER

        # ボディ（LIGHT_GRAY背景）
        body_y = content_y + header_h
        body_h = card_h - header_h
        bdy = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE, Inches(x), Inches(body_y),
            Inches(card_w), Inches(body_h)
        )
        bdy.fill.solid()
        bdy.fill.fore_color.rgb = LIGHT_GRAY
        bdy.line.fill.background()   # （枠線削除）テキストオブジェクトの枠線は付けない

        txB = slide.shapes.add_textbox(Inches(x + 0.12), Inches(body_y + 0.1),
                                        Inches(card_w - 0.24), Inches(body_h - 0.2))
        tf = txB.text_frame
        tf.word_wrap = True
        tf.auto_size = MSO_AUTO_SIZE.NONE

        body_text = card.get("body", "")
        if isinstance(body_text, list):
            for j, item in enumerate(body_text):
                p = tf.paragraphs[0] if j == 0 else tf.add_paragraph()
                add_rich_runs(p, f"・{item}", base_size=Pt(12), base_color=DARK_GRAY,
                              bold_color=NAVY, line_spacing=1.4)
        else:
            add_rich_runs(tf.paragraphs[0], body_text, base_size=Pt(13),
                          base_color=DARK_GRAY, bold_color=NAVY, line_spacing=1.4)

    if source:
        add_source_label(slide, source)
    add_corner_marks(slide)  # （A）deco廃止: 明背景の分析頁にエンジンが自動でコーナーマーク
    add_bottom_bar_and_footer(slide, page_num)
    return slide


def add_process_slide(prs, title, sub_message, steps, blank,
                      source=None, page_num=None):
    """縦STEPプロセスフロー

    steps: [{"title":"データ収集", "desc":"特許DBから1,176件を取得"}, ...]
    2個以下 = 大ボックス、3個 = 中、4個以上 = コンパクト
    """
    slide = prs.slides.add_slide(blank)
    sub_y = add_title_shape(slide, title)
    content_y = add_sub_message(slide, sub_message, y=sub_y)

    n = len(steps)
    available_h = 6.5 - content_y
    gap = 0.15
    step_h = (available_h - gap * (n - 1) - 0.5) / n  # 矢印分0.5確保
    step_h = min(step_h, 1.5)

    # フォントサイズ調整
    if n <= 2:
        title_size, desc_size = Pt(16), Pt(14)
    elif n <= 3:
        title_size, desc_size = Pt(14), Pt(13)
    else:
        title_size, desc_size = Pt(13), Pt(12)

    header_w = 2.2
    body_w = 9.8
    colors = [NAVY, BLUE, ACCENT, GREEN_ACCENT, AMBER]

    for i, step in enumerate(steps):
        sy = content_y + i * (step_h + gap + 0.15)
        color = colors[i % len(colors)]

        # 左ヘッダー（色付き）
        hdr = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE, Inches(0.5), Inches(sy),
            Inches(header_w), Inches(step_h)
        )
        hdr.fill.solid()
        hdr.fill.fore_color.rgb = color
        hdr.line.fill.background()

        txH = slide.shapes.add_textbox(Inches(0.6), Inches(sy + 0.1),
                                        Inches(header_w - 0.2), Inches(step_h - 0.2))
        tf_h = txH.text_frame
        tf_h.word_wrap = True
        p_h = tf_h.paragraphs[0]
        p_h.alignment = PP_ALIGN.CENTER
        set_text(p_h, f"STEP {i+1}", Pt(10), WHITE)
        p_t = tf_h.add_paragraph()
        p_t.alignment = PP_ALIGN.CENTER
        set_text(p_t, step["title"], title_size, WHITE, bold=True)

        # 右ボディ（LIGHT_GRAY背景）
        bdy = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE, Inches(0.5 + header_w + 0.1), Inches(sy),
            Inches(body_w), Inches(step_h)
        )
        bdy.fill.solid()
        bdy.fill.fore_color.rgb = LIGHT_GRAY
        bdy.line.fill.background()   # （枠線削除）テキストオブジェクトの枠線は付けない

        txB = slide.shapes.add_textbox(Inches(0.5 + header_w + 0.25), Inches(sy + 0.1),
                                        Inches(body_w - 0.3), Inches(step_h - 0.2))
        tf_b = txB.text_frame
        tf_b.word_wrap = True
        tf_b.auto_size = MSO_AUTO_SIZE.NONE
        add_rich_runs(tf_b.paragraphs[0], step.get("desc", ""),
                      base_size=desc_size, base_color=DARK_GRAY, bold_color=NAVY)

        # 下矢印（最後のステップ以外）
        if i < n - 1:
            arrow_y = sy + step_h + 0.02
            arrow = slide.shapes.add_shape(
                MSO_SHAPE.DOWN_ARROW, Inches(1.2), Inches(arrow_y),
                Inches(0.5), Inches(gap + 0.05)
            )
            arrow.fill.solid()
            arrow.fill.fore_color.rgb = color
            arrow.line.fill.background()

    if source:
        add_source_label(slide, source)
    add_corner_marks(slide)  # （A）deco廃止: 明背景の分析頁にエンジンが自動でコーナーマーク
    add_bottom_bar_and_footer(slide, page_num)
    return slide


def add_stepup_slide(prs, title, sub_message, phases, blank,
                     source=None, page_num=None):
    """ロードマップ（短期→中期→長期の横並びカード）

    phases: [{"header":"短期", "body":"基盤構築", "color":ACCENT}, ...]
    （改良8）3段の高さは**均一**（旧・階段状を廃止）。3-4段を推奨。
    """
    slide = prs.slides.add_slide(blank)
    sub_y = add_title_shape(slide, title)
    content_y = add_sub_message(slide, sub_message, y=sub_y)

    n = len(phases)
    gap = 0.5            # カード間に右向き三角形2個を置くため広めに確保
    total_w = 12.3
    bar_w = (total_w - gap * (n - 1)) / n
    base_y = 6.5  # ボトムバー直上
    max_h = base_y - content_y - 0.2
    colors = [ACCENT, BLUE, NAVY, GREEN_ACCENT]

    for i, phase in enumerate(phases):
        x = 0.5 + i * (bar_w + gap)
        # （改良8）全段を同じ高さに揃える（フル丈の均一カード）
        bar_h = max_h
        y = base_y - bar_h
        color = phase.get("color", colors[i % len(colors)])

        # ヘッダー部（上部、色付き）
        header_h = 0.5
        hdr = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE, Inches(x), Inches(y),
            Inches(bar_w), Inches(header_h)
        )
        hdr.fill.solid()
        hdr.fill.fore_color.rgb = color
        hdr.line.fill.background()

        txH = slide.shapes.add_textbox(Inches(x + 0.1), Inches(y + 0.05),
                                        Inches(bar_w - 0.2), Inches(header_h - 0.1))
        set_text(txH.text_frame.paragraphs[0], phase["header"], Pt(17), WHITE, bold=True)
        txH.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER

        # ボディ部（下部、LIGHT_GRAY）
        body_y = y + header_h
        body_h = bar_h - header_h
        bdy = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE, Inches(x), Inches(body_y),
            Inches(bar_w), Inches(body_h)
        )
        bdy.fill.solid()
        bdy.fill.fore_color.rgb = LIGHT_GRAY
        bdy.line.fill.background()   # （枠線削除）テキストオブジェクトの枠線は付けない

        txB = slide.shapes.add_textbox(Inches(x + 0.1), Inches(body_y + 0.1),
                                        Inches(bar_w - 0.2), Inches(body_h - 0.2))
        tf = txB.text_frame
        tf.word_wrap = True
        tf.auto_size = MSO_AUTO_SIZE.NONE
        body_text = phase.get("body", "")
        if isinstance(body_text, list):
            for j, item in enumerate(body_text):
                p = tf.paragraphs[0] if j == 0 else tf.add_paragraph()
                add_rich_runs(p, f"・{item}", base_size=Pt(13.5), base_color=DARK_GRAY,
                              bold_color=NAVY, line_spacing=1.35)
        else:
            add_rich_runs(tf.paragraphs[0], body_text, base_size=Pt(14.5),
                          base_color=DARK_GRAY, bold_color=NAVY, line_spacing=1.35)

    # カード間に右向き二等辺三角形を2個ずつ置き、左→右の流れを表現
    tw, th, tgap = 0.18, 0.12, 0.04
    tri_cy = (base_y - max_h) + max_h / 2 - th / 2     # カードの縦中央
    tri_color = RGBColor(0xC9, 0xCD, 0xD2)             # 薄いグレー
    for i in range(n - 1):
        gx0 = 0.5 + i * (bar_w + gap) + bar_w          # 直前カードの右端
        gmid = gx0 + gap / 2.0
        for ccx in (gmid - (th + tgap) / 2.0, gmid + (th + tgap) / 2.0):
            tri = slide.shapes.add_shape(MSO_SHAPE.ISOSCELES_TRIANGLE,
                                         Inches(ccx - tw / 2.0), Inches(tri_cy),
                                         Inches(tw), Inches(th))
            tri.rotation = 90                          # 時計回り90°＝頂点が右
            tri.fill.solid(); tri.fill.fore_color.rgb = tri_color
            tri.line.fill.background()
            try: tri.shadow.inherit = False
            except Exception: pass

    if source:
        add_source_label(slide, source)
    add_corner_marks(slide)  # （A）deco廃止: 明背景の分析頁にエンジンが自動でコーナーマーク
    add_bottom_bar_and_footer(slide, page_num)
    return slide


def add_compare_slide(prs, title, sub_message, left_title, left_items,
                      right_title, right_items, blank,
                      left_color=ACCENT, right_color=RED_ACCENT,
                      source=None, page_num=None):
    """左右比較スライド

    left_items / right_items: 各3-5項目の短い注釈リスト
    """
    slide = prs.slides.add_slide(blank)
    sub_y = add_title_shape(slide, title)
    content_y = add_sub_message(slide, sub_message, y=sub_y)

    col_w = 5.7
    left_x = 0.5
    right_x = 6.9
    remaining_h = 6.5 - content_y
    header_h = 0.4

    # 中央 "VS" マーカー
    vs_box = slide.shapes.add_textbox(Inches(6.1), Inches(content_y + 1.5),
                                       Inches(1.0), Inches(0.5))
    p_vs = vs_box.text_frame.paragraphs[0]
    p_vs.alignment = PP_ALIGN.CENTER
    set_text(p_vs, "VS", Pt(18), MEDIUM_GRAY, bold=True)

    # 中央区切り線
    line = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(6.55), Inches(content_y), Emu(9525), Inches(remaining_h)
    )
    line.fill.solid()
    line.fill.fore_color.rgb = BORDER_GRAY
    line.line.fill.background()

    for side_x, side_title, side_items, side_color in [
        (left_x, left_title, left_items, left_color),
        (right_x, right_title, right_items, right_color),
    ]:
        # カラムヘッダー
        bar = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE, Inches(side_x), Inches(content_y),
            Inches(col_w), Inches(header_h)
        )
        bar.fill.solid()
        bar.fill.fore_color.rgb = side_color
        bar.line.fill.background()

        txBox = slide.shapes.add_textbox(Inches(side_x + 0.1), Inches(content_y + 0.03),
                                          Inches(col_w - 0.2), Inches(header_h - 0.06))
        set_text(txBox.text_frame.paragraphs[0], side_title, Pt(16), WHITE, bold=True)
        txBox.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER

        # 注釈（ヘッダー下から下端まで埋める）
        add_annotation_block(slide, side_items, side_x, content_y + header_h + 0.1,
                             col_w, remaining_h - header_h - 0.2, font_size=14)

    if source:
        add_source_label(slide, source)
    add_corner_marks(slide)  # （A）deco廃止: 明背景の分析頁にエンジンが自動でコーナーマーク
    add_bottom_bar_and_footer(slide, page_num)
    return slide


def add_table_slide(prs, title, sub_message, headers, rows, blank,
                    col_widths=None, highlight_rows=None, annotations=None,
                    source=None, page_num=None):
    """テーブル + オプション注釈テキスト

    highlight_rows: ハイライト行のインデックスリスト
    annotations: テーブル横に注釈テキスト表示
    """
    slide = prs.slides.add_slide(blank)
    sub_y = add_title_shape(slide, title)
    content_y = add_sub_message(slide, sub_message, y=sub_y)

    n_cols = len(headers)
    n_rows = len(rows) + 1

    if annotations:
        table_w = 7.5
        text_x = 8.3
        text_w = 4.5
    else:
        table_w = 12.3
        text_x = None
        text_w = 0

    # 行高を残余スペースに合わせて動的計算
    available_table_h = 6.4 - content_y
    row_h = min(0.55, max(0.35, available_table_h / n_rows))
    table_h = row_h * n_rows

    table_shape = slide.shapes.add_table(
        n_rows, n_cols, Inches(0.5), Inches(content_y), Inches(table_w), Inches(table_h)
    )
    table = table_shape.table
    if col_widths:
        for i, w in enumerate(col_widths):
            table.columns[i].width = Inches(w)

    # Navyヘッダー行
    for j, header in enumerate(headers):
        cell = table.cell(0, j)
        cell.text = ""
        set_text(cell.text_frame.paragraphs[0], header, Pt(13), WHITE, bold=True)
        cell.fill.solid()
        cell.fill.fore_color.rgb = NAVY

    # ゼブラストライプデータ行
    highlight_rows = highlight_rows or []
    for i, row in enumerate(rows):
        for j, value in enumerate(row):
            cell = table.cell(i + 1, j)
            cell.text = ""
            set_text(cell.text_frame.paragraphs[0], str(value), Pt(12), DARK_GRAY)
            if i in highlight_rows:
                cell.fill.solid()
                cell.fill.fore_color.rgb = KEY_MSG_BG
            elif i % 2 == 1:
                cell.fill.solid()
                cell.fill.fore_color.rgb = LIGHT_GRAY

    # 注釈テキスト（テーブル横）
    if annotations and text_x:
        remaining_h = 6.4 - content_y
        add_annotation_block(slide, annotations, text_x, content_y, text_w, remaining_h,
                             font_size=13, has_border=True, bg_color=KEY_MSG_BG)

    if source:
        add_source_label(slide, source)
    add_corner_marks(slide)  # （A）deco廃止: 明背景の分析頁にエンジンが自動でコーナーマーク
    add_bottom_bar_and_footer(slide, page_num)
    return slide


def add_progress_bar_slide(prs, title, sub_message, items, blank,
                           source=None, page_num=None):
    """水平プログレスバー

    items: [{"label":"クラスタA", "value":58, "max_value":100, "color":ACCENT}, ...]
    """
    slide = prs.slides.add_slide(blank)
    sub_y = add_title_shape(slide, title)
    content_y = add_sub_message(slide, sub_message, y=sub_y)

    n = len(items)
    available_h = 6.5 - content_y
    bar_gap = 0.1
    bar_group_h = (available_h - 0.2) / n
    bar_h = min(0.5, bar_group_h * 0.5)
    label_h = bar_group_h - bar_h - bar_gap
    bar_max_w = 9.0
    colors = [ACCENT, BLUE, NAVY, GREEN_ACCENT, AMBER, RED_ACCENT]

    for i, item in enumerate(items):
        gy = content_y + i * bar_group_h
        color = item.get("color", colors[i % len(colors)])
        pct = item["value"] / max(item.get("max_value", 100), 1)
        bar_w = bar_max_w * pct

        # ラベル（左）
        txL = slide.shapes.add_textbox(Inches(0.5), Inches(gy), Inches(3.0), Inches(label_h))
        set_text(txL.text_frame.paragraphs[0], item["label"], Pt(14), DARK_GRAY, bold=True)

        # 背景バー（グレー、全幅）
        bg_bar = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE, Inches(3.5), Inches(gy + label_h),
            Inches(bar_max_w), Inches(bar_h)
        )
        bg_bar.fill.solid()
        bg_bar.fill.fore_color.rgb = LIGHT_GRAY
        bg_bar.line.fill.background()

        # 値バー（色付き）
        if bar_w > 0.1:
            val_bar = slide.shapes.add_shape(
                MSO_SHAPE.ROUNDED_RECTANGLE, Inches(3.5), Inches(gy + label_h),
                Inches(bar_w), Inches(bar_h)
            )
            val_bar.fill.solid()
            val_bar.fill.fore_color.rgb = color
            val_bar.line.fill.background()

        # パーセンテージ（バーの右端）
        txP = slide.shapes.add_textbox(Inches(3.5 + bar_w + 0.1), Inches(gy + label_h),
                                        Inches(1.5), Inches(bar_h))
        set_text(txP.text_frame.paragraphs[0],
                 f"{item['value']}{item.get('unit', '%')}", Pt(14), color, bold=True)

    if source:
        add_source_label(slide, source)
    add_corner_marks(slide)  # （A）deco廃止: 明背景の分析頁にエンジンが自動でコーナーマーク
    add_bottom_bar_and_footer(slide, page_num)
    return slide


def add_triangle_slide(prs, title, sub_message, elements, blank,
                       source=None, page_num=None, center="相互強化"):
    """3要素の収束関係図（技術・市場・政策が中心へ収束＝相互強化）。

    elements: [{"title":"技術","body":"...","color":INK},
               {"title":"市場","body":"...","color":ACCENT},
               {"title":"政策","body":"...","color":MEDIUM_GRAY}]
    3枚の角丸カードを三角配置し、中央ノード(center)へ矢印を収束させる。
    旧・細い灰色直線の三角（チープ）は廃止。center は中心ノードの短い語（例 "鉱物化"）。
    """
    slide = prs.slides.add_slide(blank)
    sub_y = add_title_shape(slide, title)
    content_y = add_sub_message(slide, sub_message, y=sub_y)

    cx = MARGIN_L + CONTENT_W / 2.0
    avail_top = content_y + 0.10
    avail_bot = 6.62
    cardw, cardh, hdr_h = 3.45, 1.68, 0.52
    positions = [
        (cx - cardw / 2.0, avail_top),                          # 上中央
        (MARGIN_L + 0.15, avail_bot - cardh),                   # 左下
        (MARGIN_L + CONTENT_W - cardw - 0.15, avail_bot - cardh),  # 右下
    ]
    colors = [INK, ACCENT, MEDIUM_GRAY]
    node_cx, node_cy, node_r = cx, (avail_top + avail_bot) / 2.0 + 0.05, 0.92

    # 1) 収束矢印（背面）: 各カード中心 → 中央ノード端（矢じりは中央側＝収束を表現）
    for (px, py) in positions:
        ccx, ccy = px + cardw / 2.0, py + cardh / 2.0
        dx, dy = node_cx - ccx, node_cy - ccy
        d = math.hypot(dx, dy) or 1.0
        ex, ey = node_cx - dx / d * (node_r + 0.06), node_cy - dy / d * (node_r + 0.06)
        c = slide.shapes.add_connector(1, Inches(ccx), Inches(ccy), Inches(ex), Inches(ey))
        c.line.color.rgb = ACCENT; c.line.width = Pt(2.4)
        try: c.shadow.inherit = False
        except Exception: pass
        _ln = c._element.find('.//' + f'{{{A_NS}}}ln')
        if _ln is not None:
            _te = etree.SubElement(_ln, f'{{{A_NS}}}tailEnd')
            _te.set('type', 'triangle'); _te.set('w', 'med'); _te.set('len', 'med')

    # 2) 中央ノード（クリムゾンの円＝収束点）
    node = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(node_cx - node_r), Inches(node_cy - node_r),
                                  Inches(node_r * 2), Inches(node_r * 2))
    node.fill.solid(); node.fill.fore_color.rgb = ACCENT; node.line.fill.background()
    try: node.shadow.inherit = False
    except Exception: pass
    ntf = node.text_frame; ntf.word_wrap = True
    try: ntf.vertical_anchor = MSO_ANCHOR.MIDDLE
    except Exception: pass
    np_ = ntf.paragraphs[0]; np_.alignment = PP_ALIGN.CENTER
    nr = np_.add_run(); nr.text = str(center); nr.font.size = Pt(15); nr.font.bold = True
    nr.font.color.rgb = WHITE; _apply_font(nr, heading=True)
    np2 = ntf.add_paragraph(); np2.alignment = PP_ALIGN.CENTER
    nr2 = np2.add_run(); nr2.text = "▲"; nr2.font.size = Pt(12); nr2.font.color.rgb = WHITE; _apply_font(nr2)

    # 3) 3枚の角丸カード（色付きヘッダー＋本文）＝矢印・ノードの前面
    for i, (elem, (px, py)) in enumerate(zip(elements[:3], positions)):
        color = elem.get("color", colors[i % 3])
        card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(px), Inches(py), Inches(cardw), Inches(cardh))
        try: card.adjustments[0] = 0.07
        except Exception: pass
        card.fill.solid(); card.fill.fore_color.rgb = LIGHT_GRAY; card.line.fill.background()
        try: card.shadow.inherit = False
        except Exception: pass
        hdr = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(px), Inches(py), Inches(cardw), Inches(hdr_h))
        try: hdr.adjustments[0] = 0.12
        except Exception: pass
        hdr.fill.solid(); hdr.fill.fore_color.rgb = color; hdr.line.fill.background()
        try: hdr.shadow.inherit = False
        except Exception: pass
        th = slide.shapes.add_textbox(Inches(px), Inches(py + 0.06), Inches(cardw), Inches(hdr_h - 0.10))
        thp = th.text_frame.paragraphs[0]; thp.alignment = PP_ALIGN.CENTER
        set_text(thp, elem.get("title", ""), Pt(15), WHITE, bold=True, heading=True)
        tb = slide.shapes.add_textbox(Inches(px + 0.18), Inches(py + hdr_h + 0.08), Inches(cardw - 0.36), Inches(cardh - hdr_h - 0.16))
        tbt = tb.text_frame; tbt.word_wrap = True; tbt.auto_size = MSO_AUTO_SIZE.NONE
        add_rich_runs(tbt.paragraphs[0], elem.get("body", ""), base_size=Pt(11.5), base_color=DARK_GRAY, bold_color=INK, line_spacing=1.3)

    if source:
        add_source_label(slide, source)
    add_corner_marks(slide)  # （A）deco廃止: 明背景の分析頁にエンジンが自動でコーナーマーク
    add_bottom_bar_and_footer(slide, page_num)
    return slide


def add_pyramid_slide(prs, title, sub_message, levels, blank,
                      notes=None, source=None, page_num=None):
    """結論ピラミッド（白背景版）：右に連続ピラミッド（custGeom・クリムゾン縦グラデ）、左に論拠ボックス。
    白地に合わせ、見出し・論拠は墨/グレー、ピラミッド面はクリムゾン・グラデ＋白文字。
    levels: [{"title":結論, "concrete":結論の具体(2行目), "detail":論拠}]（上→下）。"""
    slide = prs.slides.add_slide(blank)
    slide.background.fill.solid(); slide.background.fill.fore_color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
    sub_y = add_title_shape(slide, title)
    if sub_message:
        add_sub_message(slide, sub_message, y=sub_y)
    # 左：論拠ボックス（白地・淡パネル）
    nt = notes or [{"title": lv.get("title", ""), "body": lv.get("detail", "")} for lv in levels]
    lab = slide.shapes.add_textbox(Inches(0.46), Inches(2.18), Inches(3.6), Inches(0.26))
    set_text(lab.text_frame.paragraphs[0], "結論の論拠", Pt(11), ACCENT, bold=True)
    ys = [2.50, 3.86, 5.22]
    for i in range(min(3, len(nt))):
        x, y, w, h = 0.44, ys[i], 4.5, 1.24
        bx = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(h))
        bx.fill.solid(); bx.fill.fore_color.rgb = PALE_GRAY; bx.line.fill.background()
        bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(x), Inches(y), Emu(64008), Inches(h))
        bar.fill.solid(); bar.fill.fore_color.rgb = ACCENT; bar.line.fill.background()
        nb = slide.shapes.add_textbox(Inches(x + 0.20), Inches(y + 0.12), Inches(0.6), Inches(0.34))
        set_text(nb.text_frame.paragraphs[0], f"{i+1:02d}", Pt(17), ACCENT, bold=True, heading=True)
        tb = slide.shapes.add_textbox(Inches(x + 0.78), Inches(y + 0.13), Inches(w - 0.95), Inches(0.32))
        set_text(tb.text_frame.paragraphs[0], nt[i].get("title", ""), Pt(12), INK, bold=True)
        bd = slide.shapes.add_textbox(Inches(x + 0.22), Inches(y + 0.52), Inches(w - 0.42), Inches(h - 0.62))
        bdt = bd.text_frame; bdt.word_wrap = True; bdt.auto_size = MSO_AUTO_SIZE.NONE
        add_rich_runs(bdt.paragraphs[0], nt[i].get("body", ""), base_size=Pt(10), base_color=DARK_GRAY, bold_color=INK, line_spacing=1.32)
    # 右：連続ピラミッド（白地）
    ap = (9.92, 0.62); lb = (5.56, 6.92); rb = (14.36, 6.92)
    bands = [dict(yt=0.92, yb=2.50, fill=_CV["crimson3"], inner=_CV["crimson2"]),
             dict(yt=2.72, yb=4.40, fill=_CV["crimson2"], inner=_CV["crimson"]),
             dict(yt=4.62, yb=6.70, fill=_CV["crimson"], inner="E23A4E")]
    # 接地の淡い影
    _cv_poly(slide, [(5.2, 6.78), (14.0, 6.78), (13.0, 7.04), (4.7, 7.04)], "E6E8EC", 70, "E6E8EC", 80, 0.2)
    # クリムゾンの稜線
    _cv_ln(slide, *ap, lb[0]-ap[0], lb[1]-ap[1], _CV["crimson"], 60, 1.6)
    _cv_ln(slide, *ap, rb[0]-ap[0], rb[1]-ap[1], _CV["crimson"], 60, 1.6)
    for idx, b in enumerate(bands):
        _cv_grad_poly(slide, _cv_band(ap, lb, rb, b["yt"], b["yb"]), b["inner"], b["fill"], 90, "FFFFFF", 18, 0.9)
        _cv_grad_poly(slide, _cv_lfacet(ap, lb, b["yt"], b["yb"], 0.34 + idx*0.02), _CV["crimson3"], _CV["crimson2"], 90, "FFFFFF", 40, 0.3, op=58)
        _cv_grad_poly(slide, _cv_rfacet(ap, rb, b["yt"], b["yb"], 0.34 + idx*0.02), _CV["crimson3"], _CV["crimson2"], 90, "FFFFFF", 40, 0.3, op=58)
        xl = _cv_xon(*ap, *lb, b["yb"]); xr = _cv_xon(*ap, *rb, b["yb"])
        _cv_ln(slide, xl, b["yb"], xr-xl, 0, "FFFFFF", 88, 1.0)
    # 各段テキスト（2行：結論＋具体）白文字
    pos = [(8.10, 1.92, 3.84, 20, 13), (7.00, 3.28, 6.22, 24, 16), (5.96, 5.10, 8.20, 27, 18)]
    for i, lv in enumerate(levels[:3]):
        x, y, w, s1, s2 = pos[i]
        l1 = str(lv.get("title", "")).split("（")[0]
        l2 = str(lv.get("concrete", lv.get("detail", ""))).split("（")[0]
        _cv_txt_two(slide, l1, l2, x, y, w, s1, s2, c1="FFFFFF", c2="FBE3E6")
    if source:
        add_source_label(slide, source)
    add_corner_marks(slide)  # （A）deco廃止: 明背景の分析頁にエンジンが自動でコーナーマーク
    add_bottom_bar_and_footer(slide, page_num)
    return slide

def add_hypothesis_slide(prs, title, sub_message, hypotheses, blank,
                         source=None, page_num=None):
    """仮説検証テーブル

    hypotheses: [
        {"id":"H1", "hypothesis":"A社は3年以内にシェア首位", "verdict":"partially",
         "evidence":"シェア2位に浮上も首位とのギャップは依然5%"},
        ...
    ]
    verdict: "confirmed" -> OK (緑), "rejected" -> NG (赤),
             "partially" -> 要確認 (黄)
    """
    slide = prs.slides.add_slide(blank)
    sub_y = add_title_shape(slide, title)
    content_y = add_sub_message(slide, sub_message, y=sub_y)

    headers = ["ID", "仮説", "判定", "エビデンス"]
    n_rows = len(hypotheses) + 1
    available_h = 6.4 - content_y
    row_h = min(0.55, max(0.40, available_h / n_rows))
    table_h = row_h * n_rows

    VERDICT_MAP = {
        "confirmed": ("OK", GREEN_ACCENT),
        "rejected": ("NG", RED_ACCENT),
        "partially": ("---", AMBER),
    }

    table_shape = slide.shapes.add_table(
        n_rows, 4, Inches(0.5), Inches(content_y), Inches(12.3), Inches(table_h)
    )
    table = table_shape.table
    # 列幅: ID=0.8, 仮説=4.5, 判定=1.0, エビデンス=6.0
    table.columns[0].width = Inches(0.8)
    table.columns[1].width = Inches(4.5)
    table.columns[2].width = Inches(1.0)
    table.columns[3].width = Inches(6.0)

    # Navyヘッダー
    for j, header in enumerate(headers):
        cell = table.cell(0, j)
        cell.text = ""
        set_text(cell.text_frame.paragraphs[0], header, Pt(13), WHITE, bold=True)
        cell.fill.solid()
        cell.fill.fore_color.rgb = NAVY

    # データ行
    for i, hyp in enumerate(hypotheses):
        verdict_key = hyp.get("verdict", "partially")
        verdict_label, verdict_color = VERDICT_MAP.get(verdict_key, ("---", AMBER))

        row_data = [hyp.get("id", ""), hyp.get("hypothesis", ""),
                    verdict_label, hyp.get("evidence", "")]

        for j, val in enumerate(row_data):
            cell = table.cell(i + 1, j)
            cell.text = ""
            if j == 2:
                # 判定セルは色付き背景
                set_text(cell.text_frame.paragraphs[0], val, Pt(13), WHITE, bold=True)
                cell.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
                cell.fill.solid()
                cell.fill.fore_color.rgb = verdict_color
            else:
                set_text(cell.text_frame.paragraphs[0], str(val), Pt(12), DARK_GRAY)
                if i % 2 == 1:
                    cell.fill.solid()
                    cell.fill.fore_color.rgb = LIGHT_GRAY

    if source:
        add_source_label(slide, source)
    add_corner_marks(slide)  # （A）deco廃止: 明背景の分析頁にエンジンが自動でコーナーマーク
    add_bottom_bar_and_footer(slide, page_num)
    return slide


def add_timeline_slide(prs, title, sub_message, events, blank,
                       source=None, page_num=None):
    """水平タイムライン

    events: [{"year":"2015", "title":"CNF政策支援開始", "color":ACCENT}, ...]
    ラベルは上下交互配置。
    """
    slide = prs.slides.add_slide(blank)
    sub_y = add_title_shape(slide, title)
    content_y = add_sub_message(slide, sub_message, y=sub_y)

    n = len(events)
    total_w = 11.5
    start_x = 0.9
    line_y = content_y + 1.8  # タイムライン中心位置

    # 水平線
    line = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(start_x), Inches(line_y),
        Inches(total_w), Emu(19050)
    )
    line.fill.solid()
    line.fill.fore_color.rgb = NAVY
    line.line.fill.background()

    step = total_w / max(n - 1, 1) if n > 1 else total_w
    for i, ev in enumerate(events):
        x = start_x + i * step
        color = ev.get("color", ACCENT)

        # マーカー円
        dot = slide.shapes.add_shape(
            MSO_SHAPE.OVAL, Inches(x - 0.12), Inches(line_y - 0.12),
            Inches(0.35), Inches(0.35)
        )
        dot.fill.solid()
        dot.fill.fore_color.rgb = color
        dot.line.fill.background()

        # 年ラベル（マーカー内、白テキスト）
        txY = slide.shapes.add_textbox(Inches(x - 0.25), Inches(line_y - 0.4),
                                        Inches(0.8), Inches(0.25))
        p = txY.text_frame.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        set_text(p, ev["year"], Pt(10), NAVY, bold=True)

        # イベントテキスト（交互に上下配置して重なり回避）
        if i % 2 == 0:
            ty = line_y + 0.4
        else:
            ty = line_y - 1.0
        txE = slide.shapes.add_textbox(Inches(x - 0.6), Inches(ty),
                                        Inches(1.5), Inches(0.6))
        tf = txE.text_frame
        tf.word_wrap = True
        p2 = tf.paragraphs[0]
        p2.alignment = PP_ALIGN.CENTER
        set_text(p2, ev["title"], Pt(9), DARK_GRAY)

        # 縦線（マーカーからテキストへ）
        if i % 2 == 0:
            vline_y = line_y + 0.2
            vline_h = 0.2
        else:
            vline_y = line_y - 0.5
            vline_h = 0.4
        vline = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE, Inches(x + 0.03), Inches(vline_y),
            Emu(9525), Inches(vline_h)
        )
        vline.fill.solid()
        vline.fill.fore_color.rgb = color
        vline.line.fill.background()

    if source:
        add_source_label(slide, source)
    add_corner_marks(slide)  # （A）deco廃止: 明背景の分析頁にエンジンが自動でコーナーマーク
    add_bottom_bar_and_footer(slide, page_num)
    return slide


# --- auto-TOC（目次ページ番号の自動解決。手書き禁止＝挿入でズレない） ---
_TOC_PAGE_SLOTS = {}       # section_num -> 目次のページ番号 paragraph
_SECTION_FIRST_PAGE = {}   # section_num -> 本文開始ページ
_PENDING_SECTION = [None]  # 直近の章扉（次の番号付き頁が開始ページ）

def _record_section_page(page_num):
    """add_bottom_bar_and_footer を使わない頁（statement/invention等）の章開始ページ記録用。"""
    if page_num is not None and _PENDING_SECTION[0] is not None:
        _SECTION_FIRST_PAGE[_PENDING_SECTION[0]] = page_num
        _PENDING_SECTION[0] = None


def finalize_toc():
    """ビルド末尾（保存直前）に呼ぶ。章扉の次に現れた番号付き頁を各章の開始ページとして目次へ流し込む。"""
    for num, para in _TOC_PAGE_SLOTS.items():
        pg = _SECTION_FIRST_PAGE.get(num)
        set_text(para, f"P.{pg}" if pg else "—", Pt(14), MEDIUM_GRAY)
        para.alignment = PP_ALIGN.RIGHT


def add_toc_slide(prs, title, items, blank, page_num=None):
    """目次スライド — ヘアライン罫の編集組版目次
    （ゼブラ廃止。番号=クリムゾン明朝2桁、タイトル=左揃え墨、行下に細罫）

    items = [{"num":1, "title":"セクション名", "page":"P5"}, ...]
    """
    slide = prs.slides.add_slide(blank)
    sub_y = add_title_shape(slide, title)

    n = len(items)
    table_x, table_y, table_w = 1.5, sub_y + 0.1, 10.3
    row_h = min(0.5, (6.4 - table_y) / max(n, 1))

    rule = RGBColor(0xDD, 0xDE, 0xDF)
    for i, item in enumerate(items):
        y = table_y + i * row_h
        # 行下の細罫
        ln = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE, Inches(table_x), Inches(y + row_h - 0.012),
            Inches(table_w), Inches(0.012)
        )
        ln.fill.solid()
        ln.fill.fore_color.rgb = rule
        ln.line.fill.background()

        # 番号（クリムゾン・明朝・2桁）
        txNum = slide.shapes.add_textbox(Inches(table_x + 0.1), Inches(y + 0.05),
                                          Inches(0.9), Inches(row_h - 0.1))
        txNum.text_frame.word_wrap = True   # wrap=none による中央寄り描画を防ぐ
        pn = txNum.text_frame.paragraphs[0]
        set_text(pn, f"{int(item.get('num', i+1)):02d}", Pt(15), ACCENT,
                 bold=True, heading=True)
        pn.alignment = PP_ALIGN.LEFT

        # セクション名（左揃え・墨）
        txTitle = slide.shapes.add_textbox(Inches(table_x + 1.2), Inches(y + 0.05),
                                            Inches(7.0), Inches(row_h - 0.1))
        txTitle.text_frame.word_wrap = True
        pt = txTitle.text_frame.paragraphs[0]
        set_text(pt, item["title"], Pt(14), INK, bold=True)
        pt.alignment = PP_ALIGN.LEFT

        # ページ番号
        txPage = slide.shapes.add_textbox(Inches(table_x + 8.5), Inches(y + 0.05),
                                           Inches(1.5), Inches(row_h - 0.1))
        txPage.text_frame.word_wrap = True
        p = txPage.text_frame.paragraphs[0]
        pg = item.get("page", "")
        if pg in (None, "", "auto"):
            _TOC_PAGE_SLOTS[item.get("num", i + 1)] = p
            pg = "…"
        set_text(p, pg, Pt(14), MEDIUM_GRAY)
        p.alignment = PP_ALIGN.RIGHT

    add_corner_marks(slide)  # （A）deco廃止: 明背景の分析頁にエンジンが自動でコーナーマーク
    add_bottom_bar_and_footer(slide, page_num)
    return slide


def add_dual_panel_slide(prs, title, sub_message,
                          left_label, left_image, left_caption,
                          right_label, right_image, right_caption,
                          left_bullets=None, right_bullets=None,
                          blank=None, source=None, page_num=None):
    """2カラムチャート — 2つの可視化を並列比較"""
    slide = prs.slides.add_slide(blank)
    sub_y = add_title_shape(slide, title)
    content_y = add_sub_message(slide, sub_message, y=sub_y)

    col_w = 5.9
    left_x = 0.5
    right_x = 6.9
    remaining_h = 6.5 - content_y

    if left_bullets or right_bullets:
        n_bullets = max(len(left_bullets or []), len(right_bullets or []))
        if n_bullets <= 1:
            chart_h = remaining_h * 0.78
        elif n_bullets <= 2:
            chart_h = remaining_h * 0.68
        else:
            chart_h = remaining_h * 0.58
        text_y = content_y + chart_h + 0.1
        text_h = remaining_h - chart_h - 0.3
    else:
        chart_h = remaining_h - 0.5
        text_y = None
        text_h = 0

    for side_x, label, img_path, caption, bullets in [
        (left_x, left_label, left_image, left_caption, left_bullets),
        (right_x, right_label, right_image, right_caption, right_bullets),
    ]:
        add_chart_label(slide, label, side_x, content_y, col_w)
        full_path = os.path.join(SNAP, img_path) if not os.path.isabs(img_path) else img_path
        fit_image(slide, full_path, max_x=side_x, max_y=content_y + 0.3,
                  max_w=col_w, max_h=chart_h - 0.3)

        if caption:
            txBox = slide.shapes.add_textbox(Inches(side_x), Inches(content_y + chart_h),
                                              Inches(col_w), Inches(0.25))
            set_text(txBox.text_frame.paragraphs[0], caption, Pt(9), MEDIUM_GRAY)

        if bullets and text_y:
            add_annotation_block(slide, bullets, side_x, text_y, col_w, text_h, font_size=13)

    # 中央区切り線
    line = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(6.55), Inches(content_y), Emu(9525), Inches(remaining_h)
    )
    line.fill.solid()
    line.fill.fore_color.rgb = BORDER_GRAY
    line.line.fill.background()

    if source:
        add_source_label(slide, source)
    add_corner_marks(slide)  # （A）deco廃止: 明背景の分析頁にエンジンが自動でコーナーマーク
    add_bottom_bar_and_footer(slide, page_num)
    return slide


def add_narrative_slide(prs, title, sub_message, paragraphs, blank,
                         source=None, page_num=None):
    """テキスト主体 — エグゼクティブサマリーと結論にのみ使用

    全スライドの10%以下に制限すること。
    bullets が少ないときは下半分が間延びしないよう、本文ブロックを
    利用可能領域の上寄り〜中央に垂直センタリングする。
    """
    slide = prs.slides.add_slide(blank)
    sub_y = add_title_shape(slide, title)
    content_y = add_sub_message(slide, sub_message, y=sub_y) if sub_message else sub_y + 0.1

    avail_h = 6.5 - content_y
    # 段落数から概算高さを見積もり、少なければ垂直センタリングして間延びを防ぐ
    # 16pt・行間1.5・1段落あたり概算0.55in（長文は折返しで増えるため上限はavail_h）
    n_para = max(1, len(paragraphs))
    est_h = min(avail_h, n_para * 0.55 + 0.2)
    start_y = content_y + max(0.0, (avail_h - est_h) / 2.0)
    add_annotation_block(slide, paragraphs, 0.5, start_y, 12.3, est_h, font_size=16)

    if source:
        add_source_label(slide, source)
    add_corner_marks(slide)  # （A）deco廃止: 明背景の分析頁にエンジンが自動でコーナーマーク
    add_bottom_bar_and_footer(slide, page_num)
    return slide


def add_chapter_intro_slide(prs, eyebrow, title, bullets, blank,
                            source=None, page_num=None):
    """章扉直後の章導入スライド — 左に大見出し / 右に「何を・なぜ」の箇条書き。

    分析章の冒頭で "この章で何を・どんな意図で見るか" を提示する2カラム頁。
    インパクト演出（柱0）と読みやすさの両立を狙い、左に明朝の大見出し、
    右にGothicの箇条書きを置く。余白を活かして垂直センタリングする。

    Args:
        eyebrow: 章ラベル（例 "SECTION 03 / クラスタ動態"）
        title:   大見出し（明朝・24pt前後）
        bullets: ["この章で見ること1", "意図2", ...]（3-5項目推奨）
    """
    slide = prs.slides.add_slide(blank)

    # 利用領域全体（タイトル帯は使わずフルに2カラム）に対し垂直センタリング
    col_l_x = MARGIN_L
    col_l_w = 4.6
    col_r_x = MARGIN_L + col_l_w + 0.6
    col_r_w = 13.33 - col_r_x - 0.7

    # 左の見出し高さ（タイトル行数に連動）と右の箇条書き高さを別々に見積もり、全体を垂直センタリング
    title_lines = str(title).count("\n") + 1
    left_h = 0.55 + title_lines * 0.56
    n = max(1, len(bullets))
    right_h = min(5.9, n * 1.05 + 0.4)
    content_h = max(left_h, right_h)
    top_y = max(0.85, (7.05 - content_h) / 2.0)

    # 左カラム: EYEBROW（赤）+ 大見出し（明朝）+ クリムゾン縦バー（見出し高さに連動）
    bar = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(col_l_x), Inches(top_y + 0.05),
        Inches(0.06), Inches(left_h - 0.1)
    )
    bar.fill.solid(); bar.fill.fore_color.rgb = ACCENT; bar.line.fill.background()

    eb = slide.shapes.add_textbox(Inches(col_l_x + 0.24), Inches(top_y), Inches(col_l_w), Inches(0.4))
    eb.text_frame.word_wrap = True
    eb.text_frame.auto_size = MSO_AUTO_SIZE.NONE
    set_text(eb.text_frame.paragraphs[0], eyebrow, Pt(14), ACCENT, bold=True, heading=True)

    hd = slide.shapes.add_textbox(Inches(col_l_x + 0.24), Inches(top_y + 0.50),
                                  Inches(col_l_w + 0.2), Inches(left_h))
    htf = hd.text_frame; htf.word_wrap = True; htf.auto_size = MSO_AUTO_SIZE.NONE
    set_text(htf.paragraphs[0], title, Pt(32), INK, bold=True, line_spacing=1.20, heading=True)

    # 右カラム: 箇条書き。「**見出し**：説明」は見出しを太字、コロン後に改行して説明を次行へ（可読性）
    bd = slide.shapes.add_textbox(Inches(col_r_x), Inches(top_y), Inches(col_r_w), Inches(right_h))
    btf = bd.text_frame; btf.word_wrap = True; btf.auto_size = MSO_AUTO_SIZE.NONE
    _hang = str(int(0.30 * 914400))   # 説明行のぶら下げインデント（マーカー幅ぶん）
    for j, line in enumerate(bullets):
        m = re.match(r'\s*\*\*(.+?)\*\*\s*[:：]\s*(.*)', line)
        p = btf.paragraphs[0] if j == 0 else btf.add_paragraph()
        mk = p.add_run(); mk.text = "▪ "
        mk.font.size = Pt(18); mk.font.color.rgb = ACCENT; mk.font.bold = True; _apply_font(mk)
        if m:
            head, desc = m.group(1), m.group(2)
            rh = p.add_run(); rh.text = head
            rh.font.size = Pt(18); rh.font.bold = True; rh.font.color.rgb = INK; _apply_font(rh)
            p.space_after = Pt(2); p.line_spacing = 1.25
            p2 = btf.add_paragraph(); p2.space_after = Pt(13)
            p2._p.get_or_add_pPr().set('marL', _hang)
            add_rich_runs(p2, desc, base_size=Pt(15), base_color=DARK_GRAY,
                          bold_color=INK, line_spacing=1.3)
        else:
            add_rich_runs(p, line, base_size=Pt(18), base_color=DARK_GRAY,
                          bold_color=INK, line_spacing=1.4)
            p.space_after = Pt(11)

    if source:
        add_source_label(slide, source)
    add_corner_marks(slide)  # （A）deco廃止: 明背景の分析頁にエンジンが自動でコーナーマーク
    add_bottom_bar_and_footer(slide, page_num)
    return slide


def add_image_slide(prs, title, sub_message, image_path, blank,
                    caption=None, chart_label=None, source=None, page_num=None):
    """チャート全画面 — 画像が主役のスライド"""
    slide = prs.slides.add_slide(blank)
    sub_y = add_title_shape(slide, title)
    content_y = add_sub_message(slide, sub_message, y=sub_y)

    if chart_label:
        add_chart_label(slide, chart_label, 0.5, content_y, 12.3)
        img_y = content_y + 0.3
    else:
        img_y = content_y

    full_path = os.path.join(SNAP, image_path) if not os.path.isabs(image_path) else image_path
    fit_image(slide, full_path, max_x=0.5, max_y=img_y, max_w=12.3, max_h=6.4 - img_y)

    if caption:
        txBox = slide.shapes.add_textbox(Inches(0.5), Inches(6.45), Inches(12.3), Inches(0.25))
        set_text(txBox.text_frame.paragraphs[0], caption, Pt(10), MEDIUM_GRAY)
        txBox.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER

    if source:
        add_source_label(slide, source)
    add_corner_marks(slide)  # （A）deco廃止: 明背景の分析頁にエンジンが自動でコーナーマーク
    add_bottom_bar_and_footer(slide, page_num)
    return slide


def add_recommendation_slide(prs, title, sub_message, recommendations, blank,
                              source=None, page_num=None):
    """推奨アクション — 優先度バー付き

    recommendations: [{"priority":"高","title":"出願強化","timeframe":"短期","desc":"..."},...]
    """
    PRIORITY_COLORS = {"高": RED_ACCENT, "中": AMBER, "低": GREEN_ACCENT}
    slide = prs.slides.add_slide(blank)
    sub_y = add_title_shape(slide, title)
    content_y = add_sub_message(slide, sub_message, y=sub_y)

    n = len(recommendations)
    available_h = 6.5 - content_y
    card_h = min(1.3, (available_h - 0.1 * (n - 1)) / n)

    for i, rec in enumerate(recommendations):
        y = content_y + i * (card_h + 0.1)
        p_color = PRIORITY_COLORS.get(rec.get("priority", "中"), AMBER)

        bar = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE, Inches(0.5), Inches(y), Emu(54864), Inches(card_h)
        )
        bar.fill.solid()
        bar.fill.fore_color.rgb = p_color
        bar.line.fill.background()

        card = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE, Inches(0.8), Inches(y), Inches(12.0), Inches(card_h)
        )
        card.fill.solid()
        card.fill.fore_color.rgb = LIGHT_GRAY
        card.line.fill.background()   # （枠線削除）提言カードの枠線を撤去

        txBox_p = slide.shapes.add_textbox(Inches(1.0), Inches(y + 0.08), Inches(0.8), Inches(0.3))
        set_text(txBox_p.text_frame.paragraphs[0], f"[{rec.get('priority', '中')}]", Pt(10), p_color, bold=True)

        txBox_t = slide.shapes.add_textbox(Inches(1.9), Inches(y + 0.08), Inches(5.5), Inches(0.3))
        set_text(txBox_t.text_frame.paragraphs[0], rec["title"], Pt(16), NAVY, bold=True)

        if rec.get("timeframe"):
            txBox_tf = slide.shapes.add_textbox(Inches(9.0), Inches(y + 0.08), Inches(3.5), Inches(0.3))
            p_tf = txBox_tf.text_frame.paragraphs[0]
            set_text(p_tf, rec["timeframe"], Pt(13), MEDIUM_GRAY)
            p_tf.alignment = PP_ALIGN.RIGHT

        if rec.get("desc"):
            txBox_d = slide.shapes.add_textbox(Inches(1.9), Inches(y + 0.4), Inches(10.5), Inches(card_h - 0.5))
            tf_d = txBox_d.text_frame
            tf_d.word_wrap = True
            tf_d.auto_size = MSO_AUTO_SIZE.NONE
            add_rich_runs(tf_d.paragraphs[0], rec["desc"], base_size=Pt(13),
                          base_color=DARK_GRAY, bold_color=NAVY)

    if source:
        add_source_label(slide, source)
    add_corner_marks(slide)  # （A）deco廃止: 明背景の分析頁にエンジンが自動でコーナーマーク
    add_bottom_bar_and_footer(slide, page_num)
    return slide


def add_matrix_slide(prs, title, sub_message, quadrants, blank,
                     x_label="→ 成長率", y_label="↑ 規模",
                     source=None, page_num=None):
    """2x2マトリクス（4象限）

    quadrants: {"TL":{"title":"新興","items":["A社"]}, "TR":..., "BL":..., "BR":...}
    """
    slide = prs.slides.add_slide(blank)
    sub_y = add_title_shape(slide, title)
    content_y = add_sub_message(slide, sub_message, y=sub_y)

    mx = 1.5
    my = content_y + 0.2
    mw = 5.5
    mh = 6.0 - content_y
    half_w = mw / 2
    half_h = mh / 2
    quad_colors = {"TL": ACCENT, "TR": GREEN_ACCENT, "BL": MEDIUM_GRAY, "BR": RED_ACCENT}
    positions = {"TL": (mx, my), "TR": (mx + half_w, my),
                 "BL": (mx, my + half_h), "BR": (mx + half_w, my + half_h)}

    for key, pos in positions.items():
        q = quadrants.get(key, {})
        color = quad_colors[key]
        qx, qy = pos

        box = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE, Inches(qx), Inches(qy),
            Inches(half_w - 0.05), Inches(half_h - 0.05)
        )
        box.fill.solid()
        box.fill.fore_color.rgb = LIGHT_GRAY
        box.line.fill.background()   # （枠線削除）2x2マトリクスの枠線を撤去

        txBox = slide.shapes.add_textbox(
            Inches(qx + 0.1), Inches(qy + 0.1),
            Inches(half_w - 0.3), Inches(0.35)
        )
        set_text(txBox.text_frame.paragraphs[0], q.get("title", ""), Pt(13), color, bold=True)

        items = q.get("items", [])
        if items:
            txBox2 = slide.shapes.add_textbox(
                Inches(qx + 0.15), Inches(qy + 0.5),
                Inches(half_w - 0.4), Inches(half_h - 0.7)
            )
            tf = txBox2.text_frame
            tf.word_wrap = True
            for j, item in enumerate(items[:5]):
                p = tf.paragraphs[0] if j == 0 else tf.add_paragraph()
                set_text(p, f"・{item}", Pt(10), DARK_GRAY)

    # 右側に注釈用スペース（マトリクスの解説）
    ann_x = mx + mw + 0.5
    ann_w = 13.33 - ann_x - 0.5
    if ann_w > 2.0:
        txAnn = slide.shapes.add_textbox(Inches(ann_x), Inches(my),
                                          Inches(ann_w), Inches(mh))
        tf_ann = txAnn.text_frame
        tf_ann.word_wrap = True

    # 軸ラベル
    txX = slide.shapes.add_textbox(Inches(mx + mw/2 - 0.5), Inches(my + mh + 0.05),
                                    Inches(1.5), Inches(0.25))
    set_text(txX.text_frame.paragraphs[0], x_label, Pt(10), MEDIUM_GRAY)
    txY = slide.shapes.add_textbox(Inches(mx - 0.6), Inches(my + mh/2 - 0.15),
                                    Inches(0.5), Inches(0.3))
    set_text(txY.text_frame.paragraphs[0], y_label, Pt(10), MEDIUM_GRAY)

    if source:
        add_source_label(slide, source)
    add_corner_marks(slide)  # （A）deco廃止: 明背景の分析頁にエンジンが自動でコーナーマーク
    add_bottom_bar_and_footer(slide, page_num)
    return slide


def add_insight_slide(prs, title, sub_message, layers, blank,
                      label="考察", source=None, page_num=None):
    """本格考察の専用頁。各トピックを最低3段落で論じる受け皿。

    layers: [{"label":"事実", "body":"段落テキスト（複数文可）"},
             {"label":"解釈", "body":"..."},
             {"label":"洞察", "body":"..."},
             {"label":"示唆", "body":"..."}]
      4層モデル（事実→解釈→洞察→示唆）推奨。2-4ブロック。
      各 body は段落（80-220字）。ラベルは赤の小見出し、本文は墨。
    """
    if len(layers) > 4:
        print(f"[WARN] add_insight_slide: layers={len(layers)} > 4。1ブロックが薄くなるため2スライドに分割推奨")
    slide = prs.slides.add_slide(blank)
    sub_y = add_title_shape(slide, title, label=label)
    content_y = add_sub_message(slide, sub_message, y=sub_y) if sub_message else sub_y + 0.1

    n = max(1, len(layers))
    avail_h = 6.72 - content_y
    gap = 0.14                     # ブロック間を詰めて本文行数を確保
    block_h = (avail_h - gap * (n - 1)) / n
    # （A1）固定13pt・行間1.3のまま、本文ボックスに収まる文字数を算出して切り詰める
    body_off = 0.27                # ラベル行ぶんのオフセット
    body_w = CONTENT_W
    body_h = block_h - body_off
    line_sp = 1.3
    _, _, body_cap = _text_capacity(body_w, body_h, 13, line_spacing=line_sp)
    y = content_y
    for lyr in layers:
        # （赤棒削減）左の赤バーは廃止。赤ラベル文字で各層を識別する
        # ラベル（赤・小・字間広め）
        lb = slide.shapes.add_textbox(Inches(MARGIN_L), Inches(y),
                                      Inches(2.0), Inches(0.3))
        set_text(lb.text_frame.paragraphs[0], lyr.get("label", ""), Pt(11), ACCENT, bold=True)
        # 本文段落（墨・読みやすい行間）— 箱に収まる文字数に制限（フォントは縮めない）
        body = slide.shapes.add_textbox(Inches(MARGIN_L), Inches(y + body_off),
                                        Inches(body_w), Inches(body_h))
        tf = body.text_frame; tf.word_wrap = True; tf.auto_size = MSO_AUTO_SIZE.NONE
        # 考察は深く書きたい場合を優先: 切り詰めず全文表示（容量超過は注意喚起のみ・重なり許容）
        body_txt = lyr.get("body", "")
        _vis = len(re.sub(r'\*\*', '', body_txt))
        if _vis > body_cap:
            print(f"[INFO] 考察『{lyr.get('label', '')}』は目安{body_cap}字超（{_vis}字）"
                  f"。全文表示（次ブロックと重なる場合あり）— 必要なら2スライドに分割")
        add_rich_runs(tf.paragraphs[0], body_txt, base_size=Pt(13),
                      base_color=DARK_GRAY, bold_color=INK, line_spacing=line_sp)
        y += block_h + gap

    if source:
        add_source_label(slide, source)
    add_corner_marks(slide)  # （A）deco廃止: 明背景の分析頁にエンジンが自動でコーナーマーク
    add_bottom_bar_and_footer(slide, page_num)
    return slide


def _render_umap_scatter(points, kinds, ring_cid, out_path, w=2040, h=860):
    """意味マップ（UMAP散布図）をPILで描く。日本語注釈はpptx側で重ねるため画像は点と
    破線リングのみ（フォント非依存）。
    points: [(x, y, cid)] / kinds: {cid: 'emerging'|'mature'|'growth'} / ring_cid: 囲むクラスタID。
    """
    ACC = (0xC5, 0x12, 0x12); MATURE = (0x4D, 0x50, 0x55)
    GROW = (0x90, 0x94, 0x99); FAINT = (0xCE, 0xD2, 0xD7)
    KCOL = {"emerging": ACC, "mature": MATURE, "growth": GROW}
    KRAD = {"emerging": 11, "mature": 7, "growth": 7}
    KORD = {"emerging": 4, "mature": 3, "growth": 2}
    SS = 2
    cw, ch = w * SS, h * SS
    im = Image.new("RGBA", (cw, ch), (255, 255, 255, 0))
    dr = ImageDraw.Draw(im)
    xs = [p[0] for p in points]; ys = [p[1] for p in points]
    xmin, xmax = min(xs), max(xs); ymin, ymax = min(ys), max(ys)
    padx = (xmax - xmin) * 0.07 or 1.0; pady = (ymax - ymin) * 0.10 or 1.0
    xmin -= padx; xmax += padx; ymin -= pady; ymax += pady
    def tx(x): return (x - xmin) / (xmax - xmin) * cw
    def ty(y): return (1 - (y - ymin) / (ymax - ymin)) * ch
    for p in sorted(points, key=lambda p: KORD.get(kinds.get(p[2]), 0)):
        x, y, cid = p; k = kinds.get(cid)
        c = KCOL.get(k, FAINT); r = KRAD.get(k, 5) * SS
        cx, cy = tx(x), ty(y)
        a = 235 if k == "emerging" else (185 if k == "mature" else 150)
        dr.ellipse([cx - r, cy - r, cx + r, cy + r], fill=c + (a,))
    if ring_cid is not None:
        cc = [(tx(x), ty(y)) for x, y, cid in points if cid == ring_cid]
        if cc:
            ox = sum(p[0] for p in cc) / len(cc); oy = sum(p[1] for p in cc) / len(cc)
            rr = max(((p[0] - ox) ** 2 + (p[1] - oy) ** 2) ** 0.5 for p in cc) + 30 * SS
            for ang in range(0, 360, 9):
                dr.arc([ox - rr, oy - rr, ox + rr, oy + rr], ang, ang + 5,
                       fill=ACC + (255,), width=3 * SS)
    im = im.resize((w, h), Image.LANCZOS)
    im.save(out_path)
    return w / float(h)   # アスペクト比（w/h）


# =====================================================================
# 統計的予測（軽量・依存追加なし）: 公開ラグ補正 → モデル選択 → 95%予測区間
# =====================================================================
def publication_lag_adjust(counts_by_year, asof_year, visibility=(0.55, 0.85)):
    """直近年の **公開ラグ** を補正する。出願→公開に約18か月かかるため直近年は過小計上。
    visibility=(直近年の可視率, 1年前の可視率)。観測値を可視率で割り戻して実勢を推定する。"""
    out = dict(counts_by_year)
    for k, frac in enumerate(visibility):
        y = asof_year - k
        if y in out and frac and frac > 0:
            out[y] = out[y] / float(frac)
    return out


def _ols(xs, ys):
    """単回帰 y=a+bx。(a, b, 残差SE, R², Σx, Σx², n) を返す。"""
    n = len(xs); sx = sum(xs); sy = sum(ys)
    sxx = sum(x * x for x in xs); sxy = sum(x * y for x, y in zip(xs, ys))
    den = (n * sxx - sx * sx) or 1e-9
    b = (n * sxy - sx * sy) / den; a = (sy - b * sx) / n
    sse = sum((y - (a + b * x)) ** 2 for x, y in zip(xs, ys))
    se = math.sqrt(sse / max(1, n - 2))
    sst = sum((y - sy / n) ** 2 for y in ys) or 1e-9
    return a, b, se, 1 - sse / sst, sx, sxx, n


def _pred_interval(x0, a, b, se, sx, sxx, n, z=1.96):
    """回帰の95%予測区間（点推定, 下限, 上限）。"""
    xbar = sx / n; sxx_c = (sxx - n * xbar * xbar) or 1e-9
    half = z * se * math.sqrt(1 + 1.0 / n + (x0 - xbar) ** 2 / sxx_c)
    yh = a + b * x0
    return yh, yh - half, yh + half


def _fit_logistic(years, ys):
    """ロジスティック（S字飽和）をK格子探索＋ロジット線形化で当てる。(SSE,K,a,b) or None。"""
    ymax = max(ys) or 1.0; best = None; K = ymax * 1.2
    while K <= ymax * 6 + 1e-9:
        pts = [(x, math.log(v / (K - v))) for x, v in zip(years, ys) if 0 < v < K]
        if len(pts) >= 3:
            a, b, se, r2, sx, sxx, n = _ols([p[0] for p in pts], [p[1] for p in pts])
            sse = sum((v - K / (1 + math.exp(-(a + b * x)))) ** 2 for x, v in zip(years, ys))
            if best is None or sse < best[0]:
                best = (sse, K, a, b)
        K += ymax * 0.3
    return best


def fit_forecast(years, counts, horizon=2, asof_year=None, lag=True,
                 visibility=(0.55, 0.85)):
    """件数時系列に対し **線形/指数/ロジスティック** を当て、AICで選択して将来を予測する。
    公開ラグ補正→モデル選択→95%予測区間まで一気通貫の軽量フォーキャスタ。
    返り値: {'model','r2','corrected','observed','fit','pi','forecast_years'}。
    """
    years = list(years); counts = list(counts)
    asof_year = asof_year or max(years)
    cby = dict(zip(years, counts))
    corrected = publication_lag_adjust(cby, asof_year, visibility) if lag else dict(cby)
    ys = [corrected[y] for y in years]
    cand = {}
    # 線形
    a, b, se, r2, sx, sxx, n = _ols(years, ys)
    sse = sum((corrected[y] - (a + b * y)) ** 2 for y in years)
    cand["linear"] = (r2, len(years) * math.log(max(1e-9, sse / len(years))) + 4,
                      ("linear", (a, b, se, sx, sxx, n)))
    # 指数（log線形）
    pos = [(x, math.log(v)) for x, v in zip(years, ys) if v > 0]
    if len(pos) >= 3:
        a2, b2, se2, r22, sx2, sxx2, n2 = _ols([p[0] for p in pos], [p[1] for p in pos])
        sse2 = sum((corrected[y] - math.exp(a2 + b2 * y)) ** 2 for y in years)
        cand["exp"] = (r22, len(years) * math.log(max(1e-9, sse2 / len(years))) + 4,
                       ("exp", (a2, b2, se2, sx2, sxx2, n2)))
    # ロジスティック
    lg = _fit_logistic(years, ys)
    if lg:
        sse3, K, a3, b3 = lg
        cand["logistic"] = (None, len(years) * math.log(max(1e-9, sse3 / len(years))) + 6,
                            ("logistic", (K, a3, b3)))
    name = min(cand, key=lambda k: cand[k][1])      # AIC最小
    kind, prm = cand[name][2]
    r2_sel = cand[name][0]
    fyears = years + [asof_year + i for i in range(1, horizon + 1)]
    fit = {}; pi = {}
    for x in fyears:
        if kind == "linear":
            a, b, se, sx, sxx, n = prm
            yh, lo, hi = _pred_interval(x, a, b, se, sx, sxx, n)
        elif kind == "exp":
            a, b, se, sx, sxx, n = prm
            yh, lo, hi = _pred_interval(x, a, b, se, sx, sxx, n)
            yh, lo, hi = math.exp(yh), math.exp(lo), math.exp(hi)
        else:  # logistic
            K, a, b = prm
            yh = K / (1 + math.exp(-(a + b * x)))
            lo, hi = yh * 0.7, yh * 1.3   # S字は区間を素朴に±30%
        fit[x] = max(0.0, yh); pi[x] = (max(0.0, lo), max(0.0, hi))
    return {"model": name, "r2": r2_sel, "corrected": corrected, "observed": cby,
            "fit": fit, "pi": pi, "forecast_years": fyears[len(years):]}


def _render_forecast_chart(fits, out_path, w=2000, h=900):
    """予測チャートをPILで描画（観測=実線+点／予測=点線／95%区間=半透明リボン）。
    fits: [{'name','color':(r,g,b),'observed','corrected','fit','pi','forecast_years'}]
    日本語の軸ラベル・凡例はpptx側で重ねる。返り値: (xmin, xmax, ymax)。"""
    SS = 2; pad = 0.05
    cw, ch = w * SS, h * SS
    im = Image.new("RGBA", (cw, ch), (255, 255, 255, 0)); dr = ImageDraw.Draw(im)
    allyears = sorted(set().union(*[set(f["fit"].keys()) | set(f["observed"].keys()) for f in fits]))
    xmin, xmax = min(allyears), max(allyears)
    ymax = max(max(f["observed"].values() or [0]) for f in fits)
    ymax = max(ymax, max(hi for f in fits for (lo, hi) in f["pi"].values())) * 1.12 or 1.0
    def tx(x): return (pad + (x - xmin) / (xmax - xmin) * (1 - 2 * pad)) * cw
    def ty(v): return (1 - pad - v / ymax * (1 - 2 * pad)) * ch
    # 横グリッド
    for g in range(5):
        gy = ty(ymax * g / 4.0)
        dr.line([(tx(xmin), gy), (tx(xmax), gy)], fill=(0xE6, 0xE8, 0xEC, 255), width=2)
    for f in fits:
        c = f["color"]; oy = sorted(f["observed"].keys())
        last_obs = max(oy)
        # 予測リボン（最後の観測年以降）
        fy = [y for y in sorted(f["fit"].keys()) if y >= last_obs]
        if len(fy) >= 2:
            top = [(tx(y), ty(f["pi"][y][1])) for y in fy]
            bot = [(tx(y), ty(f["pi"][y][0])) for y in fy][::-1]
            dr.polygon(top + bot, fill=c + (40,))
        # 観測実線＋点
        pts = [(tx(y), ty(f["observed"][y])) for y in oy]
        if len(pts) >= 2: dr.line(pts, fill=c + (255,), width=5 * SS // 2)
        for y in oy:
            x0, y0 = tx(y), ty(f["observed"][y]); r = 7
            corr = f.get("corrected", {}).get(y)
            if corr is not None and abs(corr - f["observed"][y]) > 1e-6:
                # ラグ補正点（白抜き＋実勢位置に小マーカー）
                dr.ellipse([x0 - r, y0 - r, x0 + r, y0 + r], outline=c + (255,), width=3)
                yc = ty(corr); dr.ellipse([x0 - r, yc - r, x0 + r, yc + r], fill=c + (130,))
            else:
                dr.ellipse([x0 - r, y0 - r, x0 + r, y0 + r], fill=c + (255,))
        # 予測点線
        fline = [(tx(y), ty(f["fit"][y])) for y in fy]
        for i in range(len(fline) - 1):
            (xa, ya), (xb, yb) = fline[i], fline[i + 1]
            steps = 14
            for s in range(steps):
                if s % 2: continue
                t0 = s / steps; t1 = (s + 1) / steps
                dr.line([(xa + (xb - xa) * t0, ya + (yb - ya) * t0),
                         (xa + (xb - xa) * t1, ya + (yb - ya) * t1)], fill=c + (255,), width=4)
    im = im.resize((w, h), Image.LANCZOS); im.save(out_path)
    return xmin, xmax, ymax


def add_forecast_slide(prs, title, sub_message, series, blank,
                       label="統計予測", horizon=2, source=None, page_num=None,
                       footnote=None, visibility=(0.55, 0.85)):
    """統計予測スライド。件数時系列に軽量MLフォーキャスト（モデル選択＋公開ラグ補正＋95%区間）を当て、
    チャート（観測実線／予測点線／予測リボン）＋数値テーブルで「予測と裏付け」を可視化する。

    series: [{"name":"直接炭酸塩化","data":{2019:0,...,2024:8},"kind":"primary"}]
      kind=primary(クリムゾン)/secondary(濃灰)。2系列まで推奨。
    データ点が4年未満の系列はチャートを省略しテーブルのみ（弱い系列も数値で必ず示す）。
    """
    PCOL = {"primary": ACCENT, "secondary": RGBColor(0x4D, 0x50, 0x55)}
    PCOL_T = {"primary": (0xC5, 0x12, 0x12), "secondary": (0x4D, 0x50, 0x55)}
    slide = prs.slides.add_slide(blank)
    sub_y = add_title_shape(slide, title, label=label)
    content_y = add_sub_message(slide, sub_message, y=sub_y) if sub_message else sub_y + 0.1

    fits = []
    for s in series:
        d = {int(k): float(v) for k, v in s["data"].items()}
        yrs = sorted(d)
        if len(yrs) < 4:
            fits.append({"name": s["name"], "kind": s.get("kind", "primary"), "sparse": True,
                         "observed": d, "fit": {}, "pi": {}, "forecast_years": [], "corrected": {},
                         "model": "—", "r2": None,
                         "color": PCOL_T.get(s.get("kind", "primary"))})
            continue
        fc = fit_forecast(yrs, [d[y] for y in yrs], horizon=horizon, visibility=visibility)
        fc.update({"name": s["name"], "kind": s.get("kind", "primary"),
                   "color": PCOL_T.get(s.get("kind", "primary")), "sparse": False})
        fits.append(fc)

    drawable = [f for f in fits if not f["sparse"]]
    table_top = content_y + 0.10
    if drawable:
        tmp = os.path.join(tempfile.gettempdir(), "apollo_forecast.png")
        CW_PX, CH_PX = 2200, 760
        xmin, xmax, ymax = _render_forecast_chart(drawable, tmp, w=CW_PX, h=CH_PX)
        ratio = float(CH_PX) / CW_PX
        ch_h = min(2.5, (6.05 if footnote else 6.45) - content_y - 2.0)
        ch_w = CONTENT_W
        if ch_w * ratio > ch_h: ch_w = ch_h / ratio
        ch_h_act = ch_w * ratio
        cx = MARGIN_L + (CONTENT_W - ch_w) / 2.0; cy = content_y + 0.05
        slide.shapes.add_picture(tmp, Inches(cx), Inches(cy), Inches(ch_w), Inches(ch_h_act))
        # X軸の年ラベル（pptx重ね）
        nlab = 6
        for i in range(nlab):
            yv = round(xmin + (xmax - xmin) * i / (nlab - 1))
            fx = 0.05 + (yv - xmin) / (xmax - xmin) * 0.90
            lx = cx + fx * ch_w
            tb = slide.shapes.add_textbox(Inches(lx - 0.3), Inches(cy + ch_h_act + 0.01),
                                          Inches(0.6), Inches(0.22))
            p = tb.text_frame.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
            set_text(p, str(yv), Pt(9), MEDIUM_GRAY)
        # 凡例
        lgx = cx + 0.1
        for f in drawable:
            sw = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(lgx), Inches(cy + 0.04),
                                        Inches(0.22), Inches(0.12))
            sw.fill.solid(); sw.fill.fore_color.rgb = PCOL.get(f["kind"]); sw.line.fill.background()
            lt = slide.shapes.add_textbox(Inches(lgx + 0.28), Inches(cy - 0.04),
                                          Inches(2.4), Inches(0.25))
            set_text(lt.text_frame.paragraphs[0], f["name"], Pt(10), INK, bold=True)
            lgx += 0.3 + min(2.4, 0.3 + len(f["name"]) * 0.12)
        table_top = cy + ch_h_act + 0.28

    # 数値テーブル（モデル・R²・直近補正・翌年予測±95%CI）＝表フォールバック兼裏付け
    headers = ["系列", "モデル", "R²", "直近(ラグ補正)", f"翌年予測 (95%CI)"]
    colw = [3.0, 1.5, 1.0, 2.6, CONTENT_W - 8.1]
    rh = 0.42
    # ヘッダ
    hx = MARGIN_L
    hdr = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(MARGIN_L), Inches(table_top),
                                 Inches(CONTENT_W), Inches(rh))
    hdr.fill.solid(); hdr.fill.fore_color.rgb = ACCENT; hdr.line.fill.background()
    for j, htxt in enumerate(headers):
        tb = slide.shapes.add_textbox(Inches(hx + 0.08), Inches(table_top + 0.06),
                                      Inches(colw[j] - 0.12), Inches(rh - 0.1))
        set_text(tb.text_frame.paragraphs[0], htxt, Pt(11), RGBColor(0xFF, 0xFF, 0xFF), bold=True)
        hx += colw[j]
    # 行
    for ridx, f in enumerate(fits):
        ry = table_top + rh * (ridx + 1)
        bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(MARGIN_L), Inches(ry),
                                    Inches(CONTENT_W), Inches(rh))
        bg.fill.solid(); bg.fill.fore_color.rgb = PALE_GRAY if ridx % 2 == 0 else RGBColor(0xFF, 0xFF, 0xFF)
        bg.line.fill.background()
        if f["sparse"]:
            yrs = sorted(f["observed"]); latest = f["observed"][max(yrs)] if yrs else 0
            cells = [f["name"], "データ薄", "—", f"{latest:.0f}", "予測保留（点線提示）"]
        else:
            fy = f["forecast_years"][0]
            lo, hi = f["pi"][fy]; cor = f["corrected"][max(f["corrected"])]
            mdl = {"linear": "線形", "exp": "指数", "logistic": "ロジスティック"}.get(f["model"], f["model"])
            cells = [f["name"], mdl, f"{f['r2']:.2f}" if f["r2"] is not None else "—",
                     f"{cor:.0f}", f"{f['fit'][fy]:.0f}  [{lo:.0f}–{hi:.0f}]"]
        cxp = MARGIN_L
        for j, ctxt in enumerate(cells):
            tb = slide.shapes.add_textbox(Inches(cxp + 0.08), Inches(ry + 0.07),
                                          Inches(colw[j] - 0.12), Inches(rh - 0.12))
            col = PCOL.get(f["kind"]) if j == 0 else INK
            set_text(tb.text_frame.paragraphs[0], ctxt, Pt(10.5), col, bold=(j == 0))
            cxp += colw[j]

    if footnote:
        fy0 = table_top + rh * (len(fits) + 1) + 0.12
        fbar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(MARGIN_L), Inches(fy0),
                                      Emu(36576), Inches(0.34))
        fbar.fill.solid(); fbar.fill.fore_color.rgb = ACCENT; fbar.line.fill.background()
        ft = slide.shapes.add_textbox(Inches(MARGIN_L + 0.18), Inches(fy0 - 0.02),
                                      Inches(CONTENT_W - 0.18), Inches(0.5))
        ftt = ft.text_frame; ftt.word_wrap = True; ftt.auto_size = MSO_AUTO_SIZE.NONE
        set_text(ftt.paragraphs[0], footnote, Pt(11.5), INK, bold=True, line_spacing=1.2)

    if source:
        add_source_label(slide, source)
    add_corner_marks(slide)  # （A）deco廃止: 明背景の分析頁にエンジンが自動でコーナーマーク
    add_bottom_bar_and_footer(slide, page_num)
    return slide


def fit_lifecycle(years, annual, asof_year=None, lag=True, visibility=(0.55, 0.85)):
    """技術ライフサイクル：**累積件数にロジスティック回帰**を当て、飽和上限K・変曲年・成熟度を推定し、
    直近成長率（公開ラグ補正）と併せて局面（萌芽/成長/成熟/衰退）を判定する。
    返り値: {'cum','cum_fit','K','t0_year','maturity','cagr','phase','next_year','confidence','total'}。
    """
    years = list(years); annual = list(annual)
    asof_year = asof_year or max(years)
    ann = dict(zip(years, annual))
    if lag:
        for k, fr in enumerate(visibility):
            y = asof_year - k
            if y in ann and fr: ann[y] = ann[y] / float(fr)
    yrs = sorted(ann); cum = {}; s = 0.0
    for y in yrs:
        s += ann[y]; cum[y] = s
    last = cum[yrs[-1]]
    lg = _fit_logistic(yrs, [cum[y] for y in yrs])   # 累積にロジスティック
    if lg:
        sse, K, a, b = lg
        t0 = (-a / b) if b else yrs[len(yrs) // 2]
        cum_fit = {y: K / (1 + math.exp(-(a + b * y))) for y in yrs + [asof_year + 1, asof_year + 2]}
        maturity = last / K if K else 1.0
    else:
        K = last; a = 0.0; b = 0.0; t0 = yrs[len(yrs) // 2]; cum_fit = dict(cum); maturity = 1.0
    a3 = sum(ann[y] for y in yrs[-6:-3]) or 0.5
    b3 = sum(ann[y] for y in yrs[-3:]) or 0.5
    cagr = (b3 / a3) ** (1 / 3.0) - 1 if a3 > 0 else 0.0
    if cagr <= -0.05: phase = "衰退"
    elif cagr >= 0.20 and maturity < 0.6: phase = "萌芽/成長"
    elif cagr >= 0.05: phase = "成長"
    else: phase = "成熟"
    ny = max(0.0, cum_fit.get(asof_year + 1, last) - cum_fit.get(asof_year, last))
    conf = "高" if last >= 60 else ("中" if last >= 20 else "低")
    return {"cum": cum, "cum_fit": cum_fit, "K": K, "a": a, "b": b, "t0_year": t0,
            "maturity": maturity, "cagr": cagr, "phase": phase, "next_year": ny,
            "confidence": conf, "total": last, "years": yrs}


def add_lifecycle_slide(prs, title, sub_message, series, blank,
                        label="技術ライフサイクル", asof_year=None, source=None,
                        page_num=None, footnote=None):
    """全クラスタ＋分析トピックを **累積ロジスティック回帰** で一括評価する一覧スライド。
    各系列の局面（萌芽/成長/成熟/衰退）・飽和上限K・直近成長率・翌年予測・信頼度を1表にまとめる。

    series: [{"name":"[2] ガス分離・回収","data":{2006:4,...,2024:7},"group":"クラスタ"}]
      group は任意（"クラスタ"/"トピック" 等で区切り見出し代わりの色分けに）。
    """
    PHCOL = {"萌芽/成長": ACCENT, "成長": RGBColor(0x9A, 0x1A, 0x1A),
             "成熟": INK, "衰退": MEDIUM_GRAY}
    slide = prs.slides.add_slide(blank)
    sub_y = add_title_shape(slide, title, label=label)
    content_y = add_sub_message(slide, sub_message, y=sub_y) if sub_message else sub_y + 0.1

    rows = []
    for s in series:
        d = {int(k): float(v) for k, v in s["data"].items()}
        yrs = sorted(d)
        fc = fit_lifecycle(yrs, [d[y] for y in yrs], asof_year=asof_year)
        rows.append((s["name"], s.get("group", ""), fc))

    headers = ["系列", "局面", "累計", "変曲年(推定)", "直近CAGR", "翌年予測", "信頼度"]
    colw = [3.1, 1.5, 0.95, 1.6, 1.4, 1.35, CONTENT_W - 9.9]
    n = len(rows)
    avail = (6.15 if footnote else 6.5) - content_y - 0.10
    rh = min(0.42, avail / (n + 1))
    top = content_y + 0.10
    # ヘッダ
    hx = MARGIN_L
    hdr = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(MARGIN_L), Inches(top),
                                 Inches(CONTENT_W), Inches(rh))
    hdr.fill.solid(); hdr.fill.fore_color.rgb = ACCENT; hdr.line.fill.background()
    for j, h in enumerate(headers):
        tb = slide.shapes.add_textbox(Inches(hx + 0.07), Inches(top + (rh - 0.24) / 2),
                                      Inches(colw[j] - 0.1), Inches(0.26))
        set_text(tb.text_frame.paragraphs[0], h, Pt(10.5), RGBColor(0xFF, 0xFF, 0xFF), bold=True)
        hx += colw[j]
    # 行
    for ridx, (name, grp, fc) in enumerate(rows):
        ry = top + rh * (ridx + 1)
        bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(MARGIN_L), Inches(ry),
                                    Inches(CONTENT_W), Inches(rh))
        bg.fill.solid(); bg.fill.fore_color.rgb = PALE_GRAY if ridx % 2 == 0 else RGBColor(0xFF, 0xFF, 0xFF)
        bg.line.fill.background()
        # 変曲年：薄い系列は伏せ、未来の変曲は「≈YYYY(後)」で前変曲/後変曲を区別
        if fc["confidence"] == "低":
            tdisp = "—"
        else:
            t0 = fc["t0_year"]
            tdisp = f"{t0:.0f}" if t0 <= (asof_year or 2024) else f"≈{t0:.0f}(後)"
        cells = [name, fc["phase"], f"{fc['total']:.0f}", tdisp,
                 f"{fc['cagr']*100:+.0f}%", f"{fc['next_year']:.0f}", fc["confidence"]]
        cxp = MARGIN_L
        for j, ctxt in enumerate(cells):
            tb = slide.shapes.add_textbox(Inches(cxp + 0.07), Inches(ry + (rh - 0.22) / 2),
                                          Inches(colw[j] - 0.1), Inches(0.24))
            if j == 0:
                col = INK; bold = True
            elif j == 1:
                col = PHCOL.get(ctxt, INK); bold = True
            else:
                col = INK; bold = False
            set_text(tb.text_frame.paragraphs[0], ctxt, Pt(10), col, bold=bold)
            cxp += colw[j]

    if footnote:
        fy0 = top + rh * (n + 1) + 0.12
        fbar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(MARGIN_L), Inches(fy0),
                                      Emu(36576), Inches(0.34))
        fbar.fill.solid(); fbar.fill.fore_color.rgb = ACCENT; fbar.line.fill.background()
        ft = slide.shapes.add_textbox(Inches(MARGIN_L + 0.18), Inches(fy0 - 0.02),
                                      Inches(CONTENT_W - 0.18), Inches(0.5))
        ftt = ft.text_frame; ftt.word_wrap = True; ftt.auto_size = MSO_AUTO_SIZE.NONE
        set_text(ftt.paragraphs[0], footnote, Pt(11), INK, bold=True, line_spacing=1.2)

    if source:
        add_source_label(slide, source)
    add_corner_marks(slide)  # （A）deco廃止: 明背景の分析頁にエンジンが自動でコーナーマーク
    add_bottom_bar_and_footer(slide, page_num)
    return slide


def _render_lifecycle_curves(curves, out_path, asof_year, w=2200, h=980):
    """累積ロジスティックS字曲線をPILで描画。各系列の観測累積点＋ロジスティック回帰曲線
    （実績区間=実線／将来=点線）＋変曲点マーカーを重ねる。日本語注釈はpptx側。
    curves: [{'name','color','cum':{y:v},'a','b','K','t0_year','last_year'}]。
    返り値: (xmin, xmax, ymax)。"""
    SS = 2; pad = 0.06
    cw, chh = w * SS, h * SS
    im = Image.new("RGBA", (cw, chh), (255, 255, 255, 0)); dr = ImageDraw.Draw(im)
    allyears = []
    for c in curves: allyears += list(c["cum"].keys())
    xmin = min(allyears); xmax = max(asof_year + 2, max(allyears))
    # Y軸は実観測累積の最大を基準（外挿Kには引っ張られない）。将来外挿が突き抜けても上方でクリップ。
    ymax = (max(max(c["cum"].values() or [0]) for c in curves) * 1.25) or 1.0
    def tx(x): return (pad + (x - xmin) / (xmax - xmin) * (1 - 2 * pad)) * cw
    def ty(v): return (1 - pad - min(v, ymax) / ymax * (1 - 2 * pad)) * chh
    for g in range(5):
        gy = ty(ymax * g / 4.0)
        dr.line([(tx(xmin), gy), (tx(xmax), gy)], fill=(0xE6, 0xE8, 0xEC, 255), width=2)
    for c in curves:
        col = c["color"]; a = c["a"]; b = c["b"]; K = c["K"]; ly = c["last_year"]
        # 連続S字曲線（細刻み）
        prev = None; step = (xmax - xmin) / 240.0; x = xmin
        while x <= xmax + 1e-6:
            v = K / (1 + math.exp(-(a + b * x))) if b else 0
            p = (tx(x), ty(v))
            if prev is not None:
                if x <= ly:
                    dr.line([prev, p], fill=col + (255,), width=5)            # 実績区間=実線
                else:
                    seg = int((x - xmin) / step)
                    if seg % 2 == 0:
                        dr.line([prev, p], fill=col + (255,), width=4)        # 将来=点線
            prev = p; x += step
        # 観測累積点
        for y in sorted(c["cum"].keys()):
            x0, y0 = tx(y), ty(c["cum"][y]); r = 6
            dr.ellipse([x0 - r, y0 - r, x0 + r, y0 + r], fill=col + (235,))
        # 変曲点マーカー（◇）
        t0 = c["t0_year"]
        if xmin <= t0 <= xmax:
            vy = ty(K / 2.0); xx = tx(t0); d = 11
            dr.polygon([(xx, vy - d), (xx + d, vy), (xx, vy + d), (xx - d, vy)],
                       outline=col + (255,), fill=(255, 255, 255, 230))
            dr.polygon([(xx, vy - d), (xx + d, vy), (xx, vy + d), (xx - d, vy)], outline=col + (255,))
    im = im.resize((w, h), Image.LANCZOS); im.save(out_path)
    return xmin, xmax, ymax


def add_lifecycle_curve_slide(prs, title, sub_message, series, blank,
                              label="技術ライフサイクル", asof_year=None, source=None,
                              page_num=None, footnote=None):
    """技術ライフサイクル（S字曲線）スライド。指定系列の **累積件数にロジスティック回帰** を当て、
    観測点＋S字フィット曲線（実績=実線／将来=点線）＋変曲点（◇）を重ねて描く。一覧表（add_lifecycle_slide）
    と対で使い、『成熟＝寝たS字／萌芽＝立ち上がり途中』を視覚化する。3-5系列推奨。"""
    # 設計原則（赤は主役1点）: 先頭系列だけクリムゾン、残りはグレー階調で描く
    PAL = [ACCENT, RGBColor(0x3A, 0x3D, 0x42), RGBColor(0x7A, 0x7E, 0x84),
           RGBColor(0xB0, 0xB4, 0xBA), RGBColor(0x55, 0x58, 0x5E)]
    PAL_T = [(0xC5, 0x12, 0x12), (0x3A, 0x3D, 0x42), (0x7A, 0x7E, 0x84),
             (0xB0, 0xB4, 0xBA), (0x55, 0x58, 0x5E)]
    asof_year = asof_year or 2024
    slide = prs.slides.add_slide(blank)
    sub_y = add_title_shape(slide, title, label=label)
    content_y = add_sub_message(slide, sub_message, y=sub_y) if sub_message else sub_y + 0.1

    curves = []
    for i, s in enumerate(series[:5]):
        d = {int(k): float(v) for k, v in s["data"].items()}
        fc = fit_lifecycle(sorted(d), [d[y] for y in sorted(d)], asof_year=asof_year)
        curves.append({"name": s["name"], "color": PAL_T[i % len(PAL_T)], "a": fc["a"],
                       "b": fc["b"], "K": fc["K"], "t0_year": fc["t0_year"],
                       "last_year": max(fc["years"]), "cum": fc["cum"], "phase": fc["phase"]})

    tmp = os.path.join(tempfile.gettempdir(), "apollo_lifecycle_curve.png")
    CW_PX, CH_PX = 2200, 980
    xmin, xmax, ymax = _render_lifecycle_curves(curves, tmp, asof_year, w=CW_PX, h=CH_PX)
    ratio = float(CH_PX) / CW_PX
    ch_h = min(2.6, (6.4 - content_y - 0.9))
    ch_w = CONTENT_W
    if ch_w * ratio > ch_h: ch_w = ch_h / ratio
    ch_h_act = ch_w * ratio
    cx = MARGIN_L + (CONTENT_W - ch_w) / 2.0; cy = content_y + 0.05
    slide.shapes.add_picture(tmp, Inches(cx), Inches(cy), Inches(ch_w), Inches(ch_h_act))
    # X軸年ラベル
    for i in range(6):
        yv = round(xmin + (xmax - xmin) * i / 5)
        fx = 0.06 + (yv - xmin) / (xmax - xmin) * 0.88
        tb = slide.shapes.add_textbox(Inches(cx + fx * ch_w - 0.3), Inches(cy + ch_h_act + 0.02),
                                      Inches(0.6), Inches(0.22))
        p = tb.text_frame.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
        set_text(p, str(yv), Pt(9), MEDIUM_GRAY)
    # 凡例（系列名＋局面）
    lgx = cx + 0.1; lgy = cy + ch_h_act + 0.28
    for i, c in enumerate(curves):
        sw = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(lgx), Inches(lgy + 0.03),
                                    Inches(0.30), Inches(0.10))
        sw.fill.solid(); sw.fill.fore_color.rgb = PAL[i % len(PAL)]; sw.line.fill.background()
        nm = f"{c['name']}（{c['phase']}）"
        lt = slide.shapes.add_textbox(Inches(lgx + 0.36), Inches(lgy - 0.04),
                                      Inches(3.6), Inches(0.26))
        set_text(lt.text_frame.paragraphs[0], nm, Pt(10), INK, bold=True)
        lgx += 0.4 + min(3.5, 0.5 + len(nm) * 0.11)
        if lgx > MARGIN_L + CONTENT_W - 2.5:
            lgx = cx + 0.1; lgy += 0.30

    if footnote:
        fy0 = lgy + 0.40
        fbar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(MARGIN_L), Inches(fy0),
                                      Emu(36576), Inches(0.34))
        fbar.fill.solid(); fbar.fill.fore_color.rgb = ACCENT; fbar.line.fill.background()
        ft = slide.shapes.add_textbox(Inches(MARGIN_L + 0.18), Inches(fy0 - 0.02),
                                      Inches(CONTENT_W - 0.18), Inches(0.5))
        ftt = ft.text_frame; ftt.word_wrap = True; ftt.auto_size = MSO_AUTO_SIZE.NONE
        set_text(ftt.paragraphs[0], footnote, Pt(11), INK, bold=True, line_spacing=1.2)

    if source:
        add_source_label(slide, source)
    add_corner_marks(slide)  # （A）deco廃止: 明背景の分析頁にエンジンが自動でコーナーマーク
    add_bottom_bar_and_footer(slide, page_num)
    return slide


def add_query_logic_slide(prs, title, sub_message, rows, intent, blank,
                          final_logic=None, source=None, page_num=None):
    """検索式（母集団論理式）スライド：左に検索条件の表、右に設計意図【要約】。
    rows: [(no, target, cond), ...] / intent: 要約の箇条書き（list[str]）。"""
    slide = prs.slides.add_slide(blank)
    sub_y = add_title_shape(slide, title, label="母集団設計 / 検索式")
    content_y = add_sub_message(slide, sub_message, y=sub_y) if sub_message else sub_y + 0.1
    tx0 = MARGIN_L; tw = 7.35; colw = [0.45, 1.45, tw - 0.45 - 1.45]
    n = len(rows); top = content_y + 0.05
    bot = 6.45 if final_logic else 6.7
    rh = min(0.36, (bot - top) / (n + 1))
    hd = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(tx0), Inches(top), Inches(tw), Inches(rh))
    hd.fill.solid(); hd.fill.fore_color.rgb = ACCENT; hd.line.fill.background()
    for j, h in enumerate(["#", "対象", "検索条件"]):
        cx = tx0 + sum(colw[:j])
        tb = slide.shapes.add_textbox(Inches(cx + 0.06), Inches(top + (rh - 0.22) / 2), Inches(colw[j] - 0.1), Inches(0.24))
        set_text(tb.text_frame.paragraphs[0], h, Pt(10), RGBColor(0xFF, 0xFF, 0xFF), bold=True)
    for i, row in enumerate(rows):
        no, tgt, cond = row
        ry = top + rh * (i + 1)
        bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(tx0), Inches(ry), Inches(tw), Inches(rh))
        bg.fill.solid(); bg.fill.fore_color.rgb = PALE_GRAY if i % 2 == 0 else RGBColor(0xFF, 0xFF, 0xFF); bg.line.fill.background()
        for j, ct in enumerate([str(no), tgt, cond]):
            cx = tx0 + sum(colw[:j])
            tb = slide.shapes.add_textbox(Inches(cx + 0.06), Inches(ry + 0.02), Inches(colw[j] - 0.1), Inches(rh - 0.04))
            tf = tb.text_frame; tf.word_wrap = True; tf.auto_size = MSO_AUTO_SIZE.NONE
            col = ACCENT if j == 0 else INK
            set_text(tf.paragraphs[0], ct, Pt(8.0 if j == 2 else 9), col, bold=(j == 0), line_spacing=1.0)
    # 右：設計意図【要約】パネル
    px = tx0 + tw + 0.25; pw = 13.33 - px - 0.7; ph = bot - top
    panel = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(px), Inches(top), Inches(pw), Inches(ph))
    panel.fill.solid(); panel.fill.fore_color.rgb = PALE_GRAY; panel.line.fill.background()
    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(px), Inches(top), Emu(36576), Inches(ph))
    bar.fill.solid(); bar.fill.fore_color.rgb = ACCENT; bar.line.fill.background()
    t1 = slide.shapes.add_textbox(Inches(px + 0.2), Inches(top + 0.14), Inches(pw - 0.35), Inches(0.3))
    set_text(t1.text_frame.paragraphs[0], "設計意図【要約】", Pt(12), ACCENT, bold=True)
    body = slide.shapes.add_textbox(Inches(px + 0.2), Inches(top + 0.56), Inches(pw - 0.36), Inches(ph - 0.7))
    tf = body.text_frame; tf.word_wrap = True; tf.auto_size = MSO_AUTO_SIZE.NONE
    for k, b in enumerate(intent):
        p = tf.paragraphs[0] if k == 0 else tf.add_paragraph()
        add_rich_runs(p, "■ " + b, base_size=Pt(9.5), base_color=DARK_GRAY, bold_color=INK, line_spacing=1.3)
        p.space_after = Pt(5)
    if final_logic:
        fb = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(MARGIN_L), Inches(bot + 0.12), Emu(36576), Inches(0.34))
        fb.fill.solid(); fb.fill.fore_color.rgb = ACCENT; fb.line.fill.background()
        ft = slide.shapes.add_textbox(Inches(MARGIN_L + 0.18), Inches(bot + 0.10), Inches(CONTENT_W - 0.18), Inches(0.4))
        set_text(ft.text_frame.paragraphs[0], "最終論理式： " + final_logic, Pt(10.5), INK, bold=True)
    if source:
        add_source_label(slide, source)
    add_corner_marks(slide)  # （A）deco廃止: 明背景の分析頁にエンジンが自動でコーナーマーク
    add_bottom_bar_and_footer(slide, page_num)
    return slide


def add_statement_slide(prs, title, sub, rows, blank, page_num=None, emphasize_last=True):
    """暗赤背景の3点ステートメント（総括・クロージング共通）。タイトルは他スライドと整合（add_title_shape・白）。
    各行は白塗り半透明パネル＋赤ラベル＋黒本文。rows=[(label, text) x3]。"""
    s = prs.slides.add_slide(blank)
    _cv_bg_image(s, overlay=70)
    add_title_shape(s, title, color=_cvC(_V16["white"]))
    if sub:
        st = s.shapes.add_textbox(Inches(MARGIN_L), Inches(1.30), Inches(CONTENT_W), Inches(0.3))
        set_text(st.text_frame.paragraphs[0], sub, Pt(12.5), _cvC(_V16["muted"]))
    for i, (lab, val) in enumerate(rows[:3]):
        y = 1.86 + i * 1.66
        pan = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.76), Inches(y), Inches(11.82), Inches(1.42))
        try: pan.adjustments[0] = 0.05
        except Exception: pass
        pan.fill.solid(); pan.fill.fore_color.rgb = RGBColor(0xFF, 0xFF, 0xFF); _set_fill_alpha(pan, 80)  # （task4）透明度20%＝不透明度80%へ
        if emphasize_last and i == 2:
            pan.line.color.rgb = ACCENT; pan.line.width = Pt(1.6)
        else:
            pan.line.color.rgb = RGBColor(0xFF, 0xFF, 0xFF); pan.line.width = Pt(0.75); _cv_line_alpha(pan, 55)
        lbx = s.shapes.add_textbox(Inches(1.04), Inches(y), Inches(2.95), Inches(1.42))  # （task4）パネルのy中心へ
        lbx.text_frame.word_wrap = True; lbx.text_frame.auto_size = MSO_AUTO_SIZE.NONE
        lbx.text_frame.vertical_anchor = MSO_ANCHOR.MIDDLE
        set_text(lbx.text_frame.paragraphs[0], lab, Pt(19.5), ACCENT, bold=True)  # 見出しワードを拡大＋縦中央
        vb = s.shapes.add_textbox(Inches(4.12), Inches(y + 0.18), Inches(8.18), Inches(1.06))
        vt = vb.text_frame; vt.word_wrap = True; vt.auto_size = MSO_AUTO_SIZE.NONE
        vt.vertical_anchor = MSO_ANCHOR.MIDDLE
        add_rich_runs(vt.paragraphs[0], val, base_size=Pt(15), base_color=RGBColor(0x14, 0x14, 0x16), bold_color=ACCENT, line_spacing=1.28)
    _record_section_page(page_num)
    if page_num is not None:
        set_text(s.shapes.add_textbox(Inches(12.45), Inches(7.02), Inches(0.55), Inches(0.25)).text_frame.paragraphs[0],
                 str(page_num), Pt(10), _cvC(_V16["muted2"]))
    return s


def add_closing_slide(prs, finding, answer, strategy, blank, page_num=None, title=None, sub=None):
    """クロージング（add_statement_slide を利用）。『何が分かったか』でなく『だから何をするか』で締める。"""
    return add_statement_slide(
        prs,
        title or "クロージング — 分析結果を、出願と事業戦略へ",
        sub or "「分析による発見」「仮説のアンサー」「とるべき事業戦略」の3本立てで締める。",
        [("分析による発見", finding), ("仮説のアンサー", answer), ("とるべき事業戦略", strategy)],
        blank, page_num=page_num)

def add_pest_slide(prs, title, sub_message, pest, blank, source=None, page_num=None):
    """簡易PEST分析（2x2）。各象限に Material Symbol アイコン＋見出し＋箇条書き。
    pest: {"P":{"label","icon","points":[..]},"E":..,"S":..,"T":..}（政治/経済/社会/技術）"""
    slide = prs.slides.add_slide(blank)
    sub_y = add_title_shape(slide, title, label="PEST 分析")
    content_y = add_sub_message(slide, sub_message, y=sub_y) if sub_message else sub_y + 0.1
    top = content_y + 0.08
    gap = 0.30
    cw = (CONTENT_W - gap) / 2.0
    ch = (6.5 - top - gap) / 2.0
    order = [("P", 0, 0), ("E", 1, 0), ("S", 0, 1), ("T", 1, 1)]
    for key, cx_i, cy_i in order:
        d = pest.get(key, {})
        x = MARGIN_L + cx_i * (cw + gap)
        y = top + cy_i * (ch + gap)
        card = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(x), Inches(y), Inches(cw), Inches(ch))
        card.fill.solid(); card.fill.fore_color.rgb = PALE_GRAY; card.line.fill.background()
        tbar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(x), Inches(y), Inches(cw), Emu(64008))
        tbar.fill.solid(); tbar.fill.fore_color.rgb = ACCENT; tbar.line.fill.background()
        # アイコン（Material Symbol → PNG）
        ipath = material_icon_png(d.get("icon", ""), (0xC5, 0x12, 0x12), px=200) if d.get("icon") else None
        ix = x + 0.26; iy = y + 0.24
        if ipath:
            slide.shapes.add_picture(ipath, Inches(ix), Inches(iy), Inches(0.6), Inches(0.6))
        # 見出し（P 政治・規制）
        hd = slide.shapes.add_textbox(Inches(x + (1.0 if ipath else 0.26)), Inches(y + 0.30), Inches(cw - 1.2), Inches(0.5))
        hp = hd.text_frame.paragraphs[0]
        r0 = hp.add_run(); r0.text = key + "  "; r0.font.size = Pt(22); r0.font.bold = True
        r0.font.color.rgb = ACCENT; _apply_font(r0, heading=True)
        r1 = hp.add_run(); r1.text = d.get("label", ""); r1.font.size = Pt(14); r1.font.bold = True
        r1.font.color.rgb = INK; _apply_font(r1, heading=True)
        # 箇条書き
        bd = slide.shapes.add_textbox(Inches(x + 0.28), Inches(y + 1.02), Inches(cw - 0.5), Inches(ch - 1.12))
        bt = bd.text_frame; bt.word_wrap = True; bt.auto_size = MSO_AUTO_SIZE.NONE
        for j, p in enumerate(d.get("points", [])):
            para = bt.paragraphs[0] if j == 0 else bt.add_paragraph()
            add_rich_runs(para, "・" + p, base_size=Pt(11.5), base_color=DARK_GRAY, bold_color=INK, line_spacing=1.3)
            para.space_after = Pt(4)
    if source:
        add_source_label(slide, source)
    add_corner_marks(slide)  # （A）deco廃止: 明背景の分析頁にエンジンが自動でコーナーマーク
    add_bottom_bar_and_footer(slide, page_num)
    return slide


def add_invention_zone_slide(prs, zone, blank, page_num=None):
    """発明アイディア（白背景版）。タイトルは他スライドと整合（add_title_shape・明朝）。
    zone: {zoneName,headline,subtitle,claimDraft,problem,inventionPoint,priorArt:{number,applicant,claim},logicSteps:[5]}"""
    s = prs.slides.add_slide(blank)
    # 普通の白背景に変更（白×赤斜線画像は表紙へ転用したため、ここでは使わない）
    s.background.fill.solid(); s.background.fill.fore_color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
    sub_y = add_title_shape(s, zone.get("headline", ""), label=str(zone.get("zoneName", "")).upper())
    if zone.get("subtitle"):
        st = s.shapes.add_textbox(Inches(MARGIN_L), Inches(sub_y + 0.02), Inches(CONTENT_W), Inches(0.32))
        set_text(st.text_frame.paragraphs[0], zone["subtitle"], Pt(12.5), MEDIUM_GRAY)
        ctop = sub_y + 0.42
    else:
        ctop = sub_y + 0.12
    ctop = max(ctop, 1.56)
    # レイアウト: 主役の請求項パネルを縦に拡大／ポイント・先行技術を下げる／ロジックは極小・半分の高さ
    CLAIM_BOT = 4.35
    cph = CLAIM_BOT - ctop
    iy, ih = 4.48, 2.06
    ly, lh = 6.60, 0.30

    # 発明のポイント（左・下方へ移動・コンパクト）
    ix, iw = 0.9, 5.35
    p1 = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(ix), Inches(iy), Inches(iw), Inches(ih))
    p1.fill.solid(); p1.fill.fore_color.rgb = PALE_GRAY; p1.line.fill.background()
    set_text(s.shapes.add_textbox(Inches(ix + 0.2), Inches(iy + 0.12), Inches(3.0), Inches(0.26)).text_frame.paragraphs[0],
             "発明のポイント", Pt(12), ACCENT, bold=True)
    set_text(s.shapes.add_textbox(Inches(ix + 0.2), Inches(iy + 0.44), Inches(1.4), Inches(0.2)).text_frame.paragraphs[0],
             "解決課題", Pt(10), MEDIUM_GRAY, bold=True)
    a = s.shapes.add_textbox(Inches(ix + 0.2), Inches(iy + 0.66), Inches(iw - 0.4), Inches(0.58)); a.text_frame.word_wrap = True; a.text_frame.auto_size = MSO_AUTO_SIZE.NONE
    add_rich_runs(a.text_frame.paragraphs[0], zone.get("problem", ""), base_size=Pt(10), base_color=DARK_GRAY, bold_color=INK, line_spacing=1.22)
    set_text(s.shapes.add_textbox(Inches(ix + 0.2), Inches(iy + 1.26), Inches(1.4), Inches(0.2)).text_frame.paragraphs[0],
             "発明の要点", Pt(10), MEDIUM_GRAY, bold=True)
    b = s.shapes.add_textbox(Inches(ix + 0.2), Inches(iy + 1.48), Inches(iw - 0.4), Inches(ih - 1.56)); b.text_frame.word_wrap = True; b.text_frame.auto_size = MSO_AUTO_SIZE.NONE
    add_rich_runs(b.text_frame.paragraphs[0], zone.get("inventionPoint", ""), base_size=Pt(10), base_color=DARK_GRAY, bold_color=INK, line_spacing=1.22)
    # 先行技術（右・クレーム1全文）
    pa = zone.get("priorArt", {})
    rx, rw = 6.42, 6.21
    p2 = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(rx), Inches(iy), Inches(rw), Inches(ih))
    p2.fill.solid(); p2.fill.fore_color.rgb = PALE_GRAY; p2.line.fill.background()
    set_text(s.shapes.add_textbox(Inches(rx + 0.2), Inches(iy + 0.12), Inches(3.0), Inches(0.26)).text_frame.paragraphs[0],
             "先行技術", Pt(12), ACCENT, bold=True)
    hd = s.shapes.add_textbox(Inches(rx + 0.2), Inches(iy + 0.44), Inches(rw - 0.4), Inches(0.26))
    hp = hd.text_frame.paragraphs[0]
    r1 = hp.add_run(); r1.text = pa.get("number", "") + "  "; r1.font.size = Pt(11); r1.font.bold = True; r1.font.color.rgb = INK; _apply_font(r1)
    r2 = hp.add_run(); r2.text = "／ " + pa.get("applicant", ""); r2.font.size = Pt(10); r2.font.color.rgb = MEDIUM_GRAY; _apply_font(r2)
    set_text(s.shapes.add_textbox(Inches(rx + 0.2), Inches(iy + 0.74), Inches(2.4), Inches(0.2)).text_frame.paragraphs[0],
             "クレーム1（全文）", Pt(9.5), MEDIUM_GRAY, bold=True)
    cl = s.shapes.add_textbox(Inches(rx + 0.2), Inches(iy + 0.98), Inches(rw - 0.4), Inches(ih - 1.06))
    clt = cl.text_frame; clt.word_wrap = True; clt.auto_size = MSO_AUTO_SIZE.NONE
    add_rich_runs(clt.paragraphs[0], pa.get("claim", ""), base_size=Pt(7.8), base_color=DARK_GRAY, bold_color=INK, line_spacing=1.16)
    # 請求項作成ロジック（最下部・極小フォント・高さ半分）
    lp = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.9), Inches(ly), Inches(11.73), Inches(lh))
    lp.fill.solid(); lp.fill.fore_color.rgb = PALE_GRAY; lp.line.fill.background()
    set_text(s.shapes.add_textbox(Inches(1.04), Inches(ly + 0.05), Inches(2.1), Inches(0.2)).text_frame.paragraphs[0],
             "請求項作成ロジック", Pt(9), ACCENT, bold=True)
    steps = zone.get("logicSteps", []); n = max(1, len(steps)); startX = 0.9 + 2.20; stepW = (11.73 - 2.4) / n
    for i, stp in enumerate(steps):
        sx = startX + i * stepW
        set_text(s.shapes.add_textbox(Inches(sx), Inches(ly + 0.05), Inches(0.22), Inches(0.18)).text_frame.paragraphs[0], str(i + 1), Pt(8), ACCENT, bold=True)
        t2 = s.shapes.add_textbox(Inches(sx + 0.20), Inches(ly + 0.03), Inches(stepW - 0.28), Inches(lh - 0.04)); t2.text_frame.word_wrap = True; t2.text_frame.auto_size = MSO_AUTO_SIZE.NONE
        set_text(t2.text_frame.paragraphs[0], stp, Pt(8), DARK_GRAY, bold=True, line_spacing=0.95)
    # 想定独立請求項案（白パネル＋赤枠・主役）★最後に描画して最前面へ
    cp = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.9), Inches(ctop), Inches(11.73), Inches(cph))
    cp.fill.solid(); cp.fill.fore_color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
    cp.line.color.rgb = ACCENT; cp.line.width = Pt(1.6)
    try: cp.shadow.inherit = False
    except Exception: pass
    set_text(s.shapes.add_textbox(Inches(1.14), Inches(ctop + 0.16), Inches(4.0), Inches(0.3)).text_frame.paragraphs[0],
             "想定独立請求項案", Pt(13), ACCENT, bold=True)
    cb = s.shapes.add_textbox(Inches(1.14), Inches(ctop + 0.52), Inches(11.25), Inches(cph - 0.64))
    cbt = cb.text_frame; cbt.word_wrap = True; cbt.auto_size = MSO_AUTO_SIZE.NONE
    add_rich_runs(cbt.paragraphs[0], zone.get("claimDraft", ""), base_size=Pt(15.5), base_color=INK, bold_color=INK, line_spacing=1.34)
    _record_section_page(page_num)
    if page_num is not None:
        set_text(s.shapes.add_textbox(Inches(12.45), Inches(7.02), Inches(0.55), Inches(0.25)).text_frame.paragraphs[0], str(page_num), Pt(10), MEDIUM_GRAY)
    return s

def add_patent_deepdive_slide(prs, title, sub_message, patent, blank, source=None, page_num=None):
    """クラスタ分析の延長：注目クラスタにある1件を『紹介』する頁。特許の分解はしない。
    左＝書誌＋解決課題、右＝このクラスタでの位置づけ・着目理由／この特許の面白さ。
    patent: {number,applicant,year,ipc,problem,position,interest}"""
    slide = prs.slides.add_slide(blank)
    sub_y = add_title_shape(slide, title, label="クラスタの注目特許")
    content_y = add_sub_message(slide, sub_message, y=sub_y) if sub_message else sub_y + 0.1
    top = content_y + 0.05; bottomY = 6.5
    # 左：書誌＋解決課題
    lx, lw = MARGIN_L, 4.7
    lbx = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(lx), Inches(top), Inches(lw), Inches(bottomY - top))
    lbx.fill.solid(); lbx.fill.fore_color.rgb = PALE_GRAY; lbx.line.fill.background()
    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(lx), Inches(top), Emu(64008), Inches(bottomY - top))
    bar.fill.solid(); bar.fill.fore_color.rgb = ACCENT; bar.line.fill.background()
    yy = top + 0.18
    for kl, kv in [("公報番号", patent.get("number", "")), ("出願人", patent.get("applicant", "")),
                   ("出願年", str(patent.get("year", ""))), ("主分類(IPC)", patent.get("ipc", ""))]:
        a = slide.shapes.add_textbox(Inches(lx + 0.22), Inches(yy), Inches(1.35), Inches(0.24))
        set_text(a.text_frame.paragraphs[0], kl, Pt(10), MEDIUM_GRAY, bold=True)
        b = slide.shapes.add_textbox(Inches(lx + 1.55), Inches(yy), Inches(lw - 1.75), Inches(0.5))
        bt = b.text_frame; bt.word_wrap = True; bt.auto_size = MSO_AUTO_SIZE.NONE
        set_text(bt.paragraphs[0], kv, Pt(11), INK, bold=True, line_spacing=1.1)
        yy += 0.5
    yy += 0.06
    pl = slide.shapes.add_textbox(Inches(lx + 0.22), Inches(yy), Inches(lw - 0.4), Inches(0.26))
    set_text(pl.text_frame.paragraphs[0], "解決課題", Pt(11), ACCENT, bold=True)
    pb = slide.shapes.add_textbox(Inches(lx + 0.22), Inches(yy + 0.32), Inches(lw - 0.42), Inches(bottomY - (yy + 0.42)))
    pbt = pb.text_frame; pbt.word_wrap = True; pbt.auto_size = MSO_AUTO_SIZE.NONE
    add_rich_runs(pbt.paragraphs[0], patent.get("problem", ""), base_size=Pt(11.5), base_color=DARK_GRAY, bold_color=INK, line_spacing=1.4)
    # 右：2枠（位置づけ・着目理由／面白さ）
    rx = lx + lw + 0.3; rw = 13.33 - rx - 0.7
    h = (bottomY - top - 0.24) / 2
    for k, (lab, key) in enumerate([("このクラスタでの位置づけ・着目理由", "position"), ("この特許の面白さ", "interest")]):
        ry = top + k * (h + 0.24)
        bx = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(rx), Inches(ry), Inches(rw), Inches(h))
        bx.fill.solid(); bx.fill.fore_color.rgb = PALE_GRAY; bx.line.fill.background()
        rb = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(rx), Inches(ry), Emu(64008), Inches(h))
        rb.fill.solid(); rb.fill.fore_color.rgb = ACCENT; rb.line.fill.background()
        lb = slide.shapes.add_textbox(Inches(rx + 0.22), Inches(ry + 0.14), Inches(rw - 0.4), Inches(0.3))
        set_text(lb.text_frame.paragraphs[0], lab, Pt(12.5), ACCENT, bold=True)
        bd = slide.shapes.add_textbox(Inches(rx + 0.22), Inches(ry + 0.52), Inches(rw - 0.42), Inches(h - 0.64))
        bdt = bd.text_frame; bdt.word_wrap = True; bdt.auto_size = MSO_AUTO_SIZE.NONE
        add_rich_runs(bdt.paragraphs[0], patent.get(key, ""), base_size=Pt(12), base_color=DARK_GRAY, bold_color=INK, line_spacing=1.45)
    if source:
        add_source_label(slide, source)
    add_corner_marks(slide)  # （A）deco廃止: 明背景の分析頁にエンジンが自動でコーナーマーク
    add_bottom_bar_and_footer(slide, page_num)
    return slide

def add_semantic_map_slide(prs, title, sub_message, points, callouts, blank,
                           label="意味マップ", source=None, page_num=None,
                           footnote=None, ring_cid=None):
    """意味マップ（UMAP散布図）スライド。SBERT意味空間でのクラスタ布置を実データで描き、
    『件数の主役』と『意味的に外れた萌芽領域』を対比して、APOLLO固有の掬い方を可視化する。

    points: [(umap_x, umap_y, cid)] 全件。
    callouts: [{"cid":"0","text":"[0] 萌芽 ＋38%","kind":"emerging",
                "anchor":(fx,fy),"box":(fx,fy)}]
       kind=emerging(クリムゾン)/mature(濃灰)/growth(中灰)。anchor=クラスタ重心の
       マップ内分率(0-1)、box=ラベル位置の分率。callout に挙げた cid だけ着色、他は淡灰。
    ring_cid: 破線リングで囲むクラスタID（萌芽の独立性を強調・任意）。
    """
    slide = prs.slides.add_slide(blank)
    sub_y = add_title_shape(slide, title, label=label)
    content_y = add_sub_message(slide, sub_message, y=sub_y) if sub_message else sub_y + 0.1

    kinds = {c["cid"]: c.get("kind", "emerging") for c in callouts}
    tmp = os.path.join(tempfile.gettempdir(), "apollo_umap_story.png")
    ar = _render_umap_scatter(points, kinds, ring_cid, tmp)   # w/h

    # マップを領域内にアスペクト比保持で配置
    bot = 6.30 if footnote else 6.55
    max_w = CONTENT_W; max_h = bot - content_y - 0.05
    mw = max_w; mh = mw / ar
    if mh > max_h:
        mh = max_h; mw = mh * ar
    mx = MARGIN_L + (max_w - mw) / 2.0; my = content_y + 0.02
    slide.shapes.add_picture(tmp, Inches(mx), Inches(my), Inches(mw), Inches(mh))

    KCOL = {"emerging": ACCENT, "mature": RGBColor(0x4D, 0x50, 0x55),
            "growth": RGBColor(0x90, 0x94, 0x99)}

    # （C）anchor/box="auto"（または未指定）対応: cidの重心から自動算出。
    # 正規化は _render_umap_scatter と一致させる（padx7%/pady10%・y反転）→ ピルが正しく重心を指す。
    def _auto_anchor(cid):
        xs = [p[0] for p in points]; ys = [p[1] for p in points]
        x0, x1 = min(xs), max(xs); y0, y1 = min(ys), max(ys)
        px = (x1 - x0) * 0.07 or 1.0; py = (y1 - y0) * 0.10 or 1.0
        x0 -= px; x1 += px; y0 -= py; y1 += py
        cc = [(p[0], p[1]) for p in points if p[2] == cid]
        if not cc:
            return (0.5, 0.5)
        ox = sum(p[0] for p in cc) / len(cc); oy = sum(p[1] for p in cc) / len(cc)
        return ((ox - x0) / (x1 - x0), 1 - (oy - y0) / (y1 - y0))

    for c in callouts:
        kc = KCOL.get(c.get("kind"), ACCENT)
        an = c.get("anchor")
        if an in (None, "auto"):
            an = _auto_anchor(c["cid"])
        bxy = c.get("box")
        if bxy in (None, "auto"):
            ox, oy = an
            bxy = (min(max(ox + (0.13 if ox < 0.5 else -0.13), 0.04), 0.92),
                   min(max(oy + (0.13 if oy < 0.5 else -0.13), 0.06), 0.92))
        ax = mx + an[0] * mw; ay = my + an[1] * mh
        bx = mx + bxy[0] * mw; by = my + bxy[1] * mh
        # 重心マーカー（小さな塗り点）
        mk = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(ax - 0.05), Inches(ay - 0.05),
                                    Inches(0.10), Inches(0.10))
        mk.fill.solid(); mk.fill.fore_color.rgb = kc; mk.line.fill.background()
        # リーダー線（マーカー→ラベル）
        try:
            ln = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT,
                                            Inches(ax), Inches(ay), Inches(bx), Inches(by))
            ln.line.color.rgb = kc; ln.line.width = Pt(1.0)
        except Exception:
            pass
        # ラベルピル（長方形・左に色バー・白半透明地）
        txt = c.get("text", "")
        pw = min(3.5, 0.20 + len(txt) * 0.135); ph = 0.40
        px = max(mx, min(bx - pw / 2, mx + mw - pw)); py = by - ph / 2
        pill = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(px), Inches(py),
                                      Inches(pw), Inches(ph))
        pill.fill.solid(); pill.fill.fore_color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
        _set_fill_alpha(pill, 90); pill.line.color.rgb = kc; pill.line.width = Pt(1.0)
        lbar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(px), Inches(py),
                                      Emu(36576), Inches(ph))
        lbar.fill.solid(); lbar.fill.fore_color.rgb = kc; lbar.line.fill.background()
        tb = slide.shapes.add_textbox(Inches(px + 0.14), Inches(py + 0.03),
                                      Inches(pw - 0.18), Inches(ph - 0.06))
        tf = tb.text_frame; tf.word_wrap = False; tf.auto_size = MSO_AUTO_SIZE.NONE
        tf.vertical_anchor = MSO_ANCHOR.MIDDLE
        set_text(tf.paragraphs[0], txt, Pt(11.5), INK, bold=True)

    if footnote:
        fbar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(MARGIN_L), Inches(bot + 0.06),
                                      Emu(36576), Inches(0.34))
        fbar.fill.solid(); fbar.fill.fore_color.rgb = ACCENT; fbar.line.fill.background()
        ft = slide.shapes.add_textbox(Inches(MARGIN_L + 0.18), Inches(bot + 0.04),
                                      Inches(CONTENT_W - 0.18), Inches(0.4))
        ftt = ft.text_frame; ftt.word_wrap = True; ftt.auto_size = MSO_AUTO_SIZE.NONE
        set_text(ftt.paragraphs[0], footnote, Pt(12.5), INK, bold=True, line_spacing=1.2)

    if source:
        add_source_label(slide, source)
    add_corner_marks(slide)  # （A）deco廃止: 明背景の分析頁にエンジンが自動でコーナーマーク
    add_bottom_bar_and_footer(slide, page_num)
    return slide


def add_method_flow_slide(prs, title, sub_message, steps, blank,
                          label="分析ロジック", source=None, page_num=None,
                          footnote=None):
    """発見の道筋（横フロー／ファネル）。分析がどの順路で結論に辿り着いたかを左→右で示す専用頁。
    各ステップ＝縦長カード、間をクリムゾン矢印で連結。

    steps: [{"no":"01", "head":"俯瞰", "what":"全○件を意味で地図化",
             "detail":"SBERT→UMAP→クラスタ化（具体1-2行）", "metric":"○件"}]
      3-5ステップ推奨。no=連番、head=短い動詞（4-6字・明朝）、what=その操作（赤の小見出し1行）、
      detail=具体（墨・1-2行）、metric=右肩の件数等（任意・赤チップ）。
    footnote: カード群の下に置く締めの1行（任意・墨太）。
    """
    n = max(1, len(steps))
    if n > 5:
        print(f"[WARN] add_method_flow_slide: steps={n} > 5。窮屈になるため5以下推奨")
    slide = prs.slides.add_slide(blank)
    sub_y = add_title_shape(slide, title, label=label)
    content_y = add_sub_message(slide, sub_message, y=sub_y) if sub_message else sub_y + 0.1

    top = content_y + 0.10
    bot = 6.42 if footnote else 6.7
    card_h = bot - top
    arrow_w = 0.42
    card_w = (CONTENT_W - arrow_w * (n - 1)) / n
    pad = 0.18
    for i, st in enumerate(steps):
        cx = MARGIN_L + i * (card_w + arrow_w)
        # カード本体（淡グレー面・長方形・影なし）
        card = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(cx), Inches(top),
                                      Inches(card_w), Inches(card_h))
        card.fill.solid(); card.fill.fore_color.rgb = PALE_GRAY; card.line.fill.background()
        # 上部クリムゾン帯（半透明）＝連番の座
        band = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(cx), Inches(top),
                                      Inches(card_w), Inches(0.60))
        band.fill.solid(); band.fill.fore_color.rgb = ACCENT; band.line.fill.background()
        _set_fill_alpha(band, 12)
        # 連番（赤・明朝・大）
        nob = slide.shapes.add_textbox(Inches(cx + pad), Inches(top + 0.08),
                                       Inches(card_w - pad * 2), Inches(0.44))
        set_text(nob.text_frame.paragraphs[0], st.get("no", ""), Pt(20), ACCENT, bold=True, heading=True)
        # 見出し（墨太・明朝）
        hb = slide.shapes.add_textbox(Inches(cx + pad), Inches(top + 0.76),
                                      Inches(card_w - pad * 2), Inches(0.5))
        set_text(hb.text_frame.paragraphs[0], st.get("head", ""), Pt(17), INK, bold=True, heading=True)
        # what（赤の小見出し）
        wb = slide.shapes.add_textbox(Inches(cx + pad), Inches(top + 1.28),
                                      Inches(card_w - pad * 2), Inches(0.72))
        wt = wb.text_frame; wt.word_wrap = True; wt.auto_size = MSO_AUTO_SIZE.NONE
        set_text(wt.paragraphs[0], st.get("what", ""), Pt(12), DEEP_RED, bold=True, line_spacing=1.3)
        # detail（墨・本文）
        db = slide.shapes.add_textbox(Inches(cx + pad), Inches(top + 2.02),
                                      Inches(card_w - pad * 2), Inches(max(0.6, card_h - 2.02 - 0.62)))
        dt = db.text_frame; dt.word_wrap = True; dt.auto_size = MSO_AUTO_SIZE.NONE
        add_rich_runs(dt.paragraphs[0], st.get("detail", ""), base_size=Pt(11.5),
                      base_color=DARK_GRAY, bold_color=INK, line_spacing=1.35)
        # metric チップ（下部・任意）
        if st.get("metric"):
            mh = 0.40
            my = top + card_h - mh - 0.16
            chip = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(cx + pad),
                                          Inches(my), Inches(card_w - pad * 2), Inches(mh))
            chip.fill.solid(); chip.fill.fore_color.rgb = ACCENT; chip.line.fill.background()
            mt = slide.shapes.add_textbox(Inches(cx + pad), Inches(my + 0.03),
                                          Inches(card_w - pad * 2), Inches(mh - 0.04))
            mp = mt.text_frame.paragraphs[0]
            set_text(mp, st.get("metric", ""), Pt(13), RGBColor(0xFF, 0xFF, 0xFF), bold=True)
            mp.alignment = PP_ALIGN.CENTER
        # 矢印（最後のカード以外）
        if i < n - 1:
            ax = cx + card_w
            arr = slide.shapes.add_textbox(Inches(ax), Inches(top + card_h / 2 - 0.35),
                                           Inches(arrow_w), Inches(0.7))
            ap = arr.text_frame.paragraphs[0]
            ap.alignment = PP_ALIGN.CENTER
            set_text(ap, "→", Pt(24), ACCENT, bold=True)

    if footnote:
        fbar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(MARGIN_L), Inches(bot + 0.14),
                                      Emu(36576), Inches(0.36))
        fbar.fill.solid(); fbar.fill.fore_color.rgb = ACCENT; fbar.line.fill.background()
        ft = slide.shapes.add_textbox(Inches(MARGIN_L + 0.18), Inches(bot + 0.12),
                                      Inches(CONTENT_W - 0.18), Inches(0.4))
        ftt = ft.text_frame; ftt.word_wrap = True; ftt.auto_size = MSO_AUTO_SIZE.NONE
        set_text(ftt.paragraphs[0], footnote, Pt(13), INK, bold=True, line_spacing=1.25)

    if source:
        add_source_label(slide, source)
    add_corner_marks(slide)  # （A）deco廃止: 明背景の分析頁にエンジンが自動でコーナーマーク
    add_bottom_bar_and_footer(slide, page_num)
    return slide


def add_patent_micro_slide(prs, title, sub_message, patents, blank,
                           label="代表特許", source=None, page_num=None):
    """ミクロ分析A: 代表特許を「番号＋タイトル＋出願人＋技術的意義1-2文」で引用。

    patents: [{"id":"特開2022-123456", "name":"発明タイトル",
               "applicant":"出願人名", "note":"技術的意義・戦略的文脈（1-2文）"}]
      1頁あたり4-6件。レポート全体で代表特許は計15件以上（複数頁に分割可）。
    """
    if len(patents) > 6:
        print(f"[WARN] add_patent_micro_slide: patents={len(patents)} > 6。1件が窮屈になるため複数頁に分割推奨")
    slide = prs.slides.add_slide(blank)
    sub_y = add_title_shape(slide, title, label=label)
    content_y = add_sub_message(slide, sub_message, y=sub_y) if sub_message else sub_y + 0.1

    n = max(1, len(patents))
    avail_h = 6.7 - content_y
    gap = 0.14
    row_h = (avail_h - gap * (n - 1)) / n
    y = content_y
    for p in patents:
        card = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(MARGIN_L), Inches(y),
                                      Inches(CONTENT_W), Inches(row_h))
        card.fill.solid(); card.fill.fore_color.rgb = PALE_GRAY; card.line.fill.background()
        # 1行目: 番号(赤) + タイトル(墨太) + 出願人(グレー)
        head = slide.shapes.add_textbox(Inches(MARGIN_L + 0.18), Inches(y + 0.08),
                                        Inches(CONTENT_W - 0.36), Inches(0.32))
        hp = head.text_frame.paragraphs[0]
        r1 = hp.add_run(); r1.text = p.get("id", "") + "  "; r1.font.size = Pt(12)
        r1.font.bold = True; r1.font.color.rgb = ACCENT; _apply_font(r1)
        r2 = hp.add_run(); r2.text = p.get("name", ""); r2.font.size = Pt(12)
        r2.font.bold = True; r2.font.color.rgb = INK; _apply_font(r2)
        r3 = hp.add_run(); r3.text = "　（" + p.get("applicant", "") + "）"
        r3.font.size = Pt(11); r3.font.color.rgb = MEDIUM_GRAY; _apply_font(r3)
        # 2行目: 技術的意義
        note = slide.shapes.add_textbox(Inches(MARGIN_L + 0.18), Inches(y + 0.40),
                                        Inches(CONTENT_W - 0.36), Inches(row_h - 0.46))
        nt = note.text_frame; nt.word_wrap = True
        set_text(nt.paragraphs[0], p.get("note", ""), Pt(11), DARK_GRAY, line_spacing=1.3)
        y += row_h + gap

    if source:
        add_source_label(slide, source)
    add_corner_marks(slide)  # （A）deco廃止: 明背景の分析頁にエンジンが自動でコーナーマーク
    add_bottom_bar_and_footer(slide, page_num)
    return slide


def add_applicant_profile_slide(prs, title, sub_message, profiles, blank,
                                label="出願人プロファイル", source=None, page_num=None):
    """ミクロ分析B: 上位5社以上を各「最低5行」の戦略分析で記述。

    profiles: [{"name":"A化学", "metrics":"182件 / クラスタ0偏在 / IPC C08L",
                "lines":["クラスタ分布の特徴…", "IPC/技術領域…", "年別推移…",
                         "競合比較・ポジション…", "戦略的含意…"]}]
      1頁あたり1-2社（5行を潰さないため）。レポート全体で5社以上。
    """
    if len(profiles) > 2:
        print(f"[WARN] add_applicant_profile_slide: profiles={len(profiles)} > 2。5行が潰れるため複数頁に分割推奨")
    slide = prs.slides.add_slide(blank)
    sub_y = add_title_shape(slide, title, label=label)
    content_y = add_sub_message(slide, sub_message, y=sub_y) if sub_message else sub_y + 0.1

    n = max(1, len(profiles))
    avail_h = 6.7 - content_y
    gap = 0.25
    card_h = (avail_h - gap * (n - 1)) / n
    y = content_y
    for pr in profiles:
        card = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(MARGIN_L), Inches(y),
                                      Inches(CONTENT_W), Inches(card_h))
        card.fill.solid(); card.fill.fore_color.rgb = LIGHT_GRAY; card.line.fill.background()
        bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(MARGIN_L), Inches(y),
                                     Emu(36576), Inches(card_h))
        bar.fill.solid(); bar.fill.fore_color.rgb = ACCENT; bar.line.fill.background()
        # ヘッダー: 社名(墨太) + メトリクス(グレー小)
        head = slide.shapes.add_textbox(Inches(MARGIN_L + 0.2), Inches(y + 0.1),
                                        Inches(CONTENT_W - 0.4), Inches(0.32))
        hp = head.text_frame.paragraphs[0]
        rn = hp.add_run(); rn.text = pr.get("name", ""); rn.font.size = Pt(15)
        rn.font.bold = True; rn.font.color.rgb = INK; _apply_font(rn)
        rm = hp.add_run(); rm.text = "　" + pr.get("metrics", "")
        rm.font.size = Pt(11); rm.font.color.rgb = MEDIUM_GRAY; _apply_font(rm)
        # 本文: 5行以上の箇条
        body = slide.shapes.add_textbox(Inches(MARGIN_L + 0.2), Inches(y + 0.46),
                                        Inches(CONTENT_W - 0.4), Inches(card_h - 0.52))
        tf = body.text_frame; tf.word_wrap = True; tf.auto_size = MSO_AUTO_SIZE.NONE
        for j, line in enumerate(pr.get("lines", [])):
            p = tf.paragraphs[0] if j == 0 else tf.add_paragraph()
            add_rich_runs(p, "・" + line, base_size=Pt(12), base_color=DARK_GRAY,
                          bold_color=INK, line_spacing=1.3)   # ** を太字処理（MD痕跡を残さない）
        y += card_h + gap

    if source:
        add_source_label(slide, source)
    add_corner_marks(slide)  # （A）deco廃止: 明背景の分析頁にエンジンが自動でコーナーマーク
    add_bottom_bar_and_footer(slide, page_num)
    return slide


from pptx.chart.data import CategoryChartData
from pptx.enum.chart import XL_CHART_TYPE, XL_LEGEND_POSITION

def _style_chart(chart, colors=None, point_colors=None, font_pt=11):
    """系列色・点別色・フォントをブランド配色で整える。"""
    try:
        chart.font.size = Pt(font_pt); chart.font.name = JA_FONT
    except Exception:
        pass
    for si, ser in enumerate(chart.series):
        if point_colors:  # 単一系列で点別に色（注目点をクリムゾン）
            for pi, pc in enumerate(point_colors):
                try:
                    pt = ser.points[pi]; pt.format.fill.solid(); pt.format.fill.fore_color.rgb = pc
                except Exception:
                    pass
        elif colors:
            ser.format.fill.solid(); ser.format.fill.fore_color.rgb = colors[si % len(colors)]

def add_bar_chart_slide(prs, title, sub_message, categories, series, blank,
                        stacked=False, colors=None, point_colors=None,
                        annotations=None, source=None, page_num=None, chart_ratio=0.62):
    """ネイティブ棒グラフ + 右に注釈（見せ方の主役・編集可能）。

    series: [("系列名", [v,...]), ...]（単一系列なら point_colors で注目点を赤に）
    colors: 系列ごとの色（既定 墨グレー段階＋クリムゾン）。stacked=True で積み上げ。
    """
    slide = prs.slides.add_slide(blank)
    sub_y = add_title_shape(slide, title)
    content_y = add_sub_message(slide, sub_message, y=sub_y)
    has_text = bool(annotations)
    gap = 0.3
    chart_w = (CONTENT_W) * (chart_ratio if has_text else 1.0) - (gap/2 if has_text else 0)
    chart_x = MARGIN_L
    chart_h = 6.5 - content_y
    cd = CategoryChartData()
    cd.categories = categories
    for name, vals in series:
        cd.add_series(name, vals)
    ctype = XL_CHART_TYPE.COLUMN_STACKED if stacked else XL_CHART_TYPE.COLUMN_CLUSTERED
    gf = slide.shapes.add_chart(ctype, Inches(chart_x), Inches(content_y),
                                Inches(chart_w), Inches(chart_h), cd)
    chart = gf.chart
    chart.has_title = False
    if len(series) > 1:
        chart.has_legend = True; chart.legend.position = XL_LEGEND_POSITION.BOTTOM
        chart.legend.include_in_layout = False
    else:
        chart.has_legend = False
    default_cols = [DEEP_RED, MEDIUM_GRAY, BAR_OLD, INK]
    _style_chart(chart, colors=colors or default_cols, point_colors=point_colors)
    try:
        chart.plots[0].gap_width = 70
    except Exception:
        pass
    # （D2）データインク最大化: 値ラベル直付け・目盛線除去・冗長な数値軸を非表示
    try:
        from pptx.enum.chart import XL_LABEL_POSITION
        plot = chart.plots[0]
        plot.has_data_labels = True
        dl = plot.data_labels
        dl.number_format = '0'; dl.number_format_is_linked = False
        dl.font.size = Pt(11); dl.font.bold = True
        try: dl.font.name = JA_FONT
        except Exception: pass
        dl.position = XL_LABEL_POSITION.INSIDE_END if stacked else XL_LABEL_POSITION.OUTSIDE_END
    except Exception:
        pass
    try:
        va = chart.value_axis
        va.has_major_gridlines = False
        if not stacked:
            va.visible = False  # 値ラベルがあるので数値軸は省く（単系列のみ）
    except Exception:
        pass
    try:
        chart.category_axis.has_major_gridlines = False
    except Exception:
        pass
    if has_text:
        tx = MARGIN_L + chart_w + gap
        add_annotation_block(slide, annotations, tx, content_y,
                             (CONTENT_W) * (1 - chart_ratio) - gap/2, chart_h - 0.2, font_size=13)
    if source:
        add_source_label(slide, source)
    add_corner_marks(slide)  # （A）deco廃止: 明背景の分析頁にエンジンが自動でコーナーマーク
    add_bottom_bar_and_footer(slide, page_num)
    return slide


def add_evidence_slide(prs, title, sub_message, items, blank,
                       label="読んだ原文（要約抜粋）", source=None, page_num=None):
    """逐語引用の証拠頁。items=[{"head":"[番号] 出願人 — 名称", "quote":"要約全文"}]。
    1頁3-4件。番号は赤・小、引用は墨の小フォント（9-10pt）。"""
    slide = prs.slides.add_slide(blank)
    sub_y = add_title_shape(slide, title, label=label)
    content_y = add_sub_message(slide, sub_message, y=sub_y) if sub_message else sub_y + 0.1
    n = max(1, len(items)); gap = 0.16
    h = (6.7 - content_y - gap * (n - 1)) / n
    y = content_y
    for it in items:
        card = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(MARGIN_L), Inches(y),
                                      Inches(CONTENT_W), Inches(h))
        card.fill.solid(); card.fill.fore_color.rgb = PALE_GRAY; card.line.fill.background()
        bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(MARGIN_L), Inches(y), Emu(36576), Inches(h))
        bar.fill.solid(); bar.fill.fore_color.rgb = ACCENT; bar.line.fill.background()
        hd = slide.shapes.add_textbox(Inches(MARGIN_L + 0.18), Inches(y + 0.06), Inches(CONTENT_W - 0.36), Inches(0.3))
        set_text(hd.text_frame.paragraphs[0], it.get("head", ""), Pt(11.5), ACCENT, bold=True)
        qt = slide.shapes.add_textbox(Inches(MARGIN_L + 0.18), Inches(y + 0.38), Inches(CONTENT_W - 0.40), Inches(h - 0.44))
        tf = qt.text_frame; tf.word_wrap = True; tf.auto_size = MSO_AUTO_SIZE.NONE
        set_text(tf.paragraphs[0], "「" + it.get("quote", "") + "」", Pt(9.5), DARK_GRAY, line_spacing=1.18)
        y += h + gap
    if source:
        add_source_label(slide, source)
    add_corner_marks(slide)  # （A）deco廃止: 明背景の分析頁にエンジンが自動でコーナーマーク
    add_bottom_bar_and_footer(slide, page_num)
    return slide


from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
import math

def add_chart_callout(slide, text: str,
                      chart_left: float, chart_top: float,
                      chart_w: float, chart_h: float,
                      x_frac: float, y_frac: float,
                      max_chars: int = 22) -> None:
    """チャート画像上にコールアウトバブルを配置する。

    Args:
        chart_left/top/w/h: チャート画像の配置座標（Inches単位）
        x_frac, y_frac   : チャート画像内の相対位置 0-1（左上が0,0）
        text             : コールアウトテキスト（短く。改行は '\\n'）
    """
    # コールアウト頂点座標
    tip_x = chart_left + chart_w * x_frac
    tip_y = chart_top  + chart_h * y_frac

    # バブルをやや右上にオフセット
    OFFSET_X, OFFSET_Y = 0.18, -0.15
    box_w, box_h = 1.35, 0.36

    bx = tip_x + OFFSET_X
    by = tip_y + OFFSET_Y

    from pptx.util import Inches, Pt, Emu
    from pptx.dml.color import RGBColor
    from pptx.enum.text import PP_ALIGN

    # --- 吹き出しライン ---
    line = slide.shapes.add_connector(
        1,  # MSO_CONNECTOR.STRAIGHT
        Inches(tip_x), Inches(tip_y),
        Inches(bx + box_w * 0.15), Inches(by + box_h)
    )
    line.line.color.rgb = RGBColor(0xC5, 0x12, 0x12)
    line.line.width = Pt(0.75)

    # --- バブル本体 ---
    txBox = slide.shapes.add_textbox(
        Inches(bx), Inches(by), Inches(box_w), Inches(box_h))
    tf = txBox.text_frame
    tf.word_wrap = True

    # 背景塗り
    from pptx.oxml.ns import qn
    from lxml import etree
    sp = txBox._element
    spPr = sp.find(qn('p:spPr'))
    solidFill = etree.SubElement(spPr, qn('a:solidFill'))
    srgbClr = etree.SubElement(solidFill, qn('a:srgbClr'))
    srgbClr.set('val', 'C51212')
    # 外枠なし
    ln = etree.SubElement(spPr, qn('a:ln'))
    noFill = etree.SubElement(ln, qn('a:noFill'))

    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.LEFT
    run = p.add_run()
    run.text = text[:max_chars]
    run.font.size = Pt(8)
    run.font.bold = True
    run.font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
    run.font.name = 'Yu Gothic'
    # padding
    tf.margin_left = Inches(0.06)
    tf.margin_top  = Inches(0.04)


# 使い方例:
# chart_l, chart_t, chart_w, chart_h = 0.9, 1.4, 7.8, 5.0
# add_chart_callout(slide, "ここが最大件数年", chart_l, chart_t, chart_w, chart_h,
#                   x_frac=0.72, y_frac=0.28)


def add_insight_tape(slide, text: str,
                     left: float, bottom_y: float,
                     width: float, height: float = 0.30) -> None:
    """チャート下端に挿入するインサイトテープ。

    Args:
        left     : テープ左端（Inches）
        bottom_y : チャート画像の下端Y座標（Inches）—— テープはここから始まる
        width    : テープ幅（Inches）
        height   : テープ高さ（Inches, default=0.30）
    """
    from pptx.util import Inches, Pt
    from pptx.dml.color import RGBColor
    from pptx.enum.text import PP_ALIGN
    from pptx.oxml.ns import qn
    from lxml import etree

    # 外枠ボックス（PALE_GRAYの薄い帯）
    box = slide.shapes.add_shape(
        1,  # MSO_SHAPE_TYPE.RECTANGLE
        Inches(left), Inches(bottom_y), Inches(width), Inches(height))
    fill = box.fill
    fill.solid()
    fill.fore_color.rgb = RGBColor(0xF6, 0xF7, 0xF8)  # PALE_GRAY
    box.line.color.rgb = RGBColor(0xD8, 0xDA, 0xDD)   # BORDER_GRAY
    box.line.width = Pt(0.5)

    # 左端クリムゾンバー（2pt幅）
    bar = slide.shapes.add_shape(
        1,
        Inches(left), Inches(bottom_y), Inches(0.04), Inches(height))
    bar.fill.solid()
    bar.fill.fore_color.rgb = RGBColor(0xC5, 0x12, 0x12)
    bar.line.fill.background()

    # テキスト
    txBox = slide.shapes.add_textbox(
        Inches(left + 0.10), Inches(bottom_y + 0.03),
        Inches(width - 0.14), Inches(height - 0.06))
    tf = txBox.text_frame
    tf.word_wrap = False
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.LEFT
    run = p.add_run()
    run.text = text
    run.font.size = Pt(9)
    run.font.color.rgb = RGBColor(0x11, 0x11, 0x11)
    run.font.name = 'Yu Gothic'


# 使い方例（チャート下端に足す）:
# chart_bottom = chart_top + chart_h  # = 1.4 + 5.0 = 6.4
# add_insight_tape(slide,
#     "上位3クラスタで全体出願の64%を占有 — 技術集中化が示す差別化余地の縮小",
#     left=0.9, bottom_y=chart_bottom, width=7.8)


SHOW_CORNER_MARKS = False   # 設計判断: コーナーマークは廃止（ナビは左端の章ガイドに一本化）

def add_corner_marks(slide,
                     slide_w_in: float = 13.33,
                     slide_h_in: float = 7.5,
                     leg: float = 0.22,
                     margin: float = 0.14) -> None:
    """スライド右上と左下にL字型コーナーマークを描く。

    Args:
        leg    : L字の一辺の長さ（Inches）
        margin : スライド端からの内側オフセット（Inches）
    """
    if not SHOW_CORNER_MARKS:   # 廃止: 既定では描かない（再有効化は SHOW_CORNER_MARKS=True）
        return
    from pptx.util import Inches, Pt, Emu
    from pptx.dml.color import RGBColor

    # （A）冪等化: 既に付与済みなら二重描画しない（deco廃止後はエンジンが自動付与するため）
    try:
        for _sh in slide.shapes:
            if getattr(_sh, "name", "") == "apollo_cmark":
                return
    except Exception:
        pass

    COLOR = RGBColor(0xC5, 0x12, 0x12)
    LW = Pt(0.75)

    def _hline(x0, y0, length):
        c = slide.shapes.add_connector(1,
            Inches(x0), Inches(y0), Inches(x0 + length), Inches(y0))
        c.line.color.rgb = COLOR
        c.line.width = LW
        try: c.name = "apollo_cmark"
        except Exception: pass

    def _vline(x0, y0, length):
        c = slide.shapes.add_connector(1,
            Inches(x0), Inches(y0), Inches(x0), Inches(y0 + length))
        c.line.color.rgb = COLOR
        c.line.width = LW

    # 右上 (top-right)
    rx = slide_w_in - margin
    ry = margin
    _hline(rx - leg, ry, leg)   # 水平
    _vline(rx, ry, leg)          # 垂直

    # 左下 (bottom-left)
    lx = margin
    ly = slide_h_in - margin
    _hline(lx, ly, leg)          # 水平
    _vline(lx, ly - leg, leg)    # 垂直


# 使い方（分析スライド末尾に呼ぶ）:
# add_corner_marks(slide)


def add_ghost_stat(slide, number_text: str, label_text: str,
                   x: float = 7.2, y: float = 1.0,
                   font_size_pt: float = 90,
                   color_hex: str = 'ECECEC') -> None:
    """背景に大きな薄いKPI数字を配置する（Layer 4）。

    Args:
        number_text : 表示する数字（例: '2,847', '38%', '#1'）
        label_text  : 数字の下に入る小ラベル（例: '総出願件数', 'YoY成長率'）
        x, y        : 配置位置（Inches）—— 既存コンテンツと重ならない位置へ
        font_size_pt: 大数字のフォントサイズ
        color_hex   : 大数字の色（默認=ECECEC、ほぼ背景に溶け込む濃さ）
    """
    from pptx.util import Inches, Pt
    from pptx.dml.color import RGBColor
    from pptx.enum.text import PP_ALIGN

    r, g, b = int(color_hex[0:2], 16), int(color_hex[2:4], 16), int(color_hex[4:6], 16)
    ghost_color = RGBColor(r, g, b)

    # 大数字
    txBig = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(4.5), Inches(1.6))
    tf = txBig.text_frame
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.LEFT
    run = p.add_run()
    run.text = number_text
    run.font.size = Pt(font_size_pt)
    run.font.bold = True
    run.font.color.rgb = ghost_color
    run.font.name = 'Century Gothic'  # 数字は欧文フォント

    # 小ラベル
    txLbl = slide.shapes.add_textbox(Inches(x + 0.05), Inches(y + 1.45), Inches(4.0), Inches(0.35))
    tf2 = txLbl.text_frame
    p2 = tf2.paragraphs[0]
    p2.alignment = PP_ALIGN.LEFT
    run2 = p2.add_run()
    run2.text = label_text
    run2.font.size = Pt(10)
    run2.font.color.rgb = RGBColor(0xBD, 0xBD, 0xBD)
    run2.font.name = 'Yu Gothic'


# 使い方例（エグゼクティブサマリー右エリアに配置）:
# add_ghost_stat(slide, '2,847', '総出願件数', x=8.2, y=2.5)
# add_ghost_stat(slide, '+38%', '直近3年CAGRvs前期比', x=8.2, y=4.5, font_size_pt=72)
